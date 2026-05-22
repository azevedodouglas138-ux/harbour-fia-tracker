import base64
import csv
import gc
import importlib
import importlib.util
import io
import json
import os
import sys
import threading
import time
from datetime import datetime, timedelta
from functools import wraps

from flask import Flask, Response, jsonify, render_template, request, send_file, session, redirect, url_for

from risk_methodology import RISK_METHODOLOGY
from liquidity_methodology import LIQUIDITY_METHODOLOGY
from ui_methodology import UI_METHODOLOGY
# Merge: o frontend usa um único window.RISK_METHODOLOGY (chaves prefixadas
# resolvem colisão: risk_* / liq_* / pt_* / cvm_* / chart_* / tab_*).
METHODOLOGY = {**RISK_METHODOLOGY, **LIQUIDITY_METHODOLOGY, **UI_METHODOLOGY}


# ---------------------------------------------------------------------------
# Lazy import de yfinance — adia o pull de pandas/numpy/lxml até a 1ª request
# que precisar de dados de mercado. Reduz baseline do worker no boot do Render
# (~80-100MB) e permite que healthcheck/login não paguem o overhead do yfinance.
# ---------------------------------------------------------------------------
def _lazy_import(name: str):
    if name in sys.modules:
        return sys.modules[name]
    spec = importlib.util.find_spec(name)
    if spec is None or spec.loader is None:
        raise ImportError(f"módulo não encontrado: {name}")
    loader = importlib.util.LazyLoader(spec.loader)
    spec.loader = loader
    module = importlib.util.module_from_spec(spec)
    sys.modules[name] = module
    loader.exec_module(module)
    return module


yf = _lazy_import("yfinance")

app = Flask(__name__)
app.secret_key = os.environ.get("SECRET_KEY", "dev-secret-change-me")

BASE_DIR = os.path.dirname(os.path.abspath(__file__))
DATA_DIR = os.path.join(BASE_DIR, "data")

LOGIN_USER  = os.environ.get("LOGIN_USER", "admin")
LOGIN_PASS  = os.environ.get("LOGIN_PASSWORD", "")
VIEWER_USER = os.environ.get("VIEWER_USER", "")
VIEWER_PASS = os.environ.get("VIEWER_PASSWORD", "")

GITHUB_TOKEN = os.environ.get("GITHUB_TOKEN", "")
GITHUB_REPO  = os.environ.get("GITHUB_REPO", "")   # "owner/repo"
GITHUB_BRANCH = os.environ.get("GITHUB_BRANCH", "main")

def _github_push(relative_path, content_str, commit_msg):
    """Push a single file to GitHub. Called in a background thread."""
    if not GITHUB_TOKEN or not GITHUB_REPO:
        return
    try:
        import requests as req
        headers = {
            "Authorization": f"token {GITHUB_TOKEN}",
            "Accept": "application/vnd.github.v3+json",
        }
        api_url = f"https://api.github.com/repos/{GITHUB_REPO}/contents/{relative_path}"
        r   = req.get(api_url, headers=headers, timeout=10)
        sha = r.json().get("sha") if r.ok else None
        payload = {
            "message": commit_msg,
            "content": base64.b64encode(content_str.encode("utf-8")).decode(),
            "branch":  GITHUB_BRANCH,
        }
        if sha:
            payload["sha"] = sha
        resp = req.put(api_url, json=payload, headers=headers, timeout=15)
        if not resp.ok:
            print(f"[github_push] {relative_path} → {resp.status_code}: {resp.text[:200]}")
    except Exception as e:
        print(f"[github_push] error: {e}")

import atexit
import queue
import signal
_github_push_queue = queue.Queue()

def _github_push_worker():
    """Worker único que processa pushes em ordem FIFO.
    Garante que o último save_X vence no GitHub (evita race condition
    quando GETs/PUTs paralelos sobrescreveriam o estado mais recente
    com conteúdo capturado em um snapshot anterior)."""
    while True:
        try:
            item = _github_push_queue.get()
            if item is None:
                _github_push_queue.task_done()
                break
            relative_path, content_str, commit_msg, done_event = item
            try:
                _github_push(relative_path, content_str, commit_msg)
            except Exception as e:
                print(f"[github_push_worker] error pushing {relative_path}: {e}")
            finally:
                if done_event is not None:
                    done_event.set()
                _github_push_queue.task_done()
        except Exception as e:
            print(f"[github_push_worker] loop error: {e}")

# Worker único, daemon — serializa todos os pushes ao GitHub
threading.Thread(target=_github_push_worker, daemon=True, name="github-push-worker").start()

def github_push_async(relative_path, content_str, commit_msg):
    """Enfileira push e retorna imediatamente. Use para writes não-críticos
    (cache, viewer_config) onde perder em shutdown não é catastrófico."""
    _github_push_queue.put((relative_path, content_str, commit_msg, None))


def _github_push_direct(relative_path, content_str, commit_msg):
    """Push síncrono ao GitHub na mesma thread. Retorna {"ok": bool, "error"?: str}.
    Não depende de worker thread externa — funciona sempre que houver creds."""
    if not GITHUB_TOKEN or not GITHUB_REPO:
        return {"ok": True, "note": "no github credentials"}
    try:
        import requests as req
        headers = {
            "Authorization": f"token {GITHUB_TOKEN}",
            "Accept": "application/vnd.github.v3+json",
        }
        api_url = f"https://api.github.com/repos/{GITHUB_REPO}/contents/{relative_path}"
        r = req.get(api_url, headers=headers, timeout=10)
        sha = r.json().get("sha") if r.ok else None
        payload = {
            "message": commit_msg,
            "content": base64.b64encode(content_str.encode("utf-8")).decode(),
            "branch":  GITHUB_BRANCH,
        }
        if sha:
            payload["sha"] = sha
        resp = req.put(api_url, json=payload, headers=headers, timeout=15)
        if resp.ok:
            return {"ok": True}
        return {"ok": False, "error": f"GitHub PUT {resp.status_code}: {resp.text[:200]}"}
    except Exception as e:
        return {"ok": False, "error": str(e)}


def github_push_sync(relative_path, content_str, commit_msg, timeout=20):
    """Push síncrono ao GitHub. BLOQUEIA na thread do request handler até
    o GitHub confirmar. Raise RuntimeError em falha — caller deve tratar
    (rollback do arquivo local + propagar erro 500 ao usuário).

    Não usa worker thread / queue (que se mostrou não-confiável em algumas
    configurações de gunicorn). O parâmetro timeout é ignorado nesta
    implementação (requests já tem timeout de 10+15s embutido)."""
    result = _github_push_direct(relative_path, content_str, commit_msg)
    if not result.get("ok"):
        raise RuntimeError(f"GitHub push falhou para {relative_path}: {result.get('error')}")

def _drain_push_queue(timeout=25):
    """Drena a fila de pushes pendentes — chamado em SIGTERM/atexit.
    Render envia SIGTERM antes de matar o container; dá ~30s para
    finalizar writes assíncronos que ainda não foram processados."""
    try:
        pending = _github_push_queue.qsize()
        if pending == 0:
            return
        print(f"[github_push] graceful shutdown: drenando {pending} push(es) pendente(s)...")
        # Espera todos os task_done() — bloqueia até a fila esvaziar
        import time as _time
        start = _time.time()
        while not _github_push_queue.empty() and (_time.time() - start) < timeout:
            _time.sleep(0.1)
        if not _github_push_queue.empty():
            print(f"[github_push] timeout drenando: {_github_push_queue.qsize()} ainda pendente(s)")
        else:
            print(f"[github_push] fila drenada em {_time.time()-start:.1f}s")
    except Exception as e:
        print(f"[github_push] drain error: {e}")

atexit.register(_drain_push_queue)
try:
    signal.signal(signal.SIGTERM, lambda *a: (_drain_push_queue(), os._exit(0)))
except (ValueError, AttributeError, OSError):
    # signal handler só funciona na thread principal; em alguns runners (gunicorn worker)
    # pode falhar. atexit acima ainda cobre o caso de exit normal.
    pass
PORTFOLIO_FILE         = os.path.join(DATA_DIR, "portfolio.json")
CACHE_FILE             = os.path.join(DATA_DIR, "cache.json")
FUND_CONFIG_FILE       = os.path.join(DATA_DIR, "fund_config.json")
QUOTA_HISTORY_FILE     = os.path.join(DATA_DIR, "quota_history.json")
VIEWER_CONFIG_FILE     = os.path.join(DATA_DIR, "viewer_config.json")
PRETRADE_HISTORY_FILE  = os.path.join(DATA_DIR, "pretrade_history.json")
PORTFOLIO_HISTORY_FILE = os.path.join(DATA_DIR, "portfolio_history.json")
INDEX_MEMBERS_FILE     = os.path.join(DATA_DIR, "index_members.json")

_price_cache = {"data": {}, "expires_at": 0}
FUNDAMENTALS_TTL = 4 * 3600
HISTORY_TTL      = 4 * 3600
STOCK_HIST_TTL   = 3600

RANGE_TO_PERIOD = {
    "1S": {"period": "5d"},
    "1M": {"period": "1mo"},
    "3M": {"period": "3mo"},
    "6M": {"period": "6mo"},
    "1A": {"period": "1y"},
}

SECTOR_PT = {
    "Energy": "Energia", "Financial Services": "Serv. Financeiros",
    "Real Estate": "Imobiliário", "Consumer Cyclical": "Consumo Cíclico",
    "Consumer Defensive": "Consumo Básico", "Healthcare": "Saúde",
    "Technology": "Tecnologia", "Industrials": "Industriais",
    "Basic Materials": "Mat. Básicos", "Communication Services": "Comunicação",
    "Utilities": "Utilidades",
}

# ---------------------------------------------------------------------------
# File helpers
# ---------------------------------------------------------------------------

def load_portfolio():
    with open(PORTFOLIO_FILE, "r", encoding="utf-8") as f:
        return json.load(f)

def _save_with_github(local_path, github_path, content, commit_msg):
    """Salva arquivo local + push síncrono ao GitHub.
    Se o push falhar, faz rollback do arquivo local e raise RuntimeError.
    Garante consistência: ou ambos sucedem ou ambos ficam no estado anterior."""
    # Backup local atual
    backup = None
    if os.path.exists(local_path):
        with open(local_path, "rb") as f:
            backup = f.read()
    # Escrever novo conteúdo localmente
    with open(local_path, "w", encoding="utf-8") as f:
        f.write(content)
    # Push GitHub; rollback local em falha
    try:
        github_push_sync(github_path, content, commit_msg)
    except Exception as e:
        if backup is not None:
            try:
                with open(local_path, "wb") as f:
                    f.write(backup)
            except Exception as rollback_err:
                print(f"[save_with_github] FALHA NO ROLLBACK de {local_path}: {rollback_err}")
        raise


def save_portfolio(data):
    content = json.dumps(data, ensure_ascii=False, indent=2)
    _save_with_github(PORTFOLIO_FILE, "data/portfolio.json", content,
                      "chore: update portfolio.json via UI")

def load_cache():
    if not os.path.exists(CACHE_FILE):
        return {}
    with open(CACHE_FILE, "r", encoding="utf-8") as f:
        return json.load(f)

def save_cache(data):
    with open(CACHE_FILE, "w", encoding="utf-8") as f:
        json.dump(data, f, ensure_ascii=False, indent=2)

def load_fund_config():
    defaults = {
        "quota_fechamento": 0, "data_fechamento": "",
        "num_cotas": None, "caixa": 0,
        "proventos_a_receber": 0, "custos_provisionados": 0,
        "performance_fee_rate": 20, "performance_fee_acumulada_rs": 0,
        "descricao_fundo": "",
        "limite_concentracao_ativo_pct": 20.0,
        "limite_concentracao_setor_pct": 40.0,
        "enable_concentracao_ativo": False,
        "enable_concentracao_setor": False,
        # Compliance de liquidez (aba 212)
        "liquidez_min_5d_pct":    80.0,
        "liquidez_max_baixa_pct": 10.0,
        "liquidez_max_zerar_dias": 30,
    }
    if not os.path.exists(FUND_CONFIG_FILE):
        return defaults
    with open(FUND_CONFIG_FILE, "r", encoding="utf-8") as f:
        data = json.load(f)
    return {**defaults, **data}

def save_fund_config(data):
    content = json.dumps(data, ensure_ascii=False, indent=2)
    _save_with_github(FUND_CONFIG_FILE, "data/fund_config.json", content,
                      "chore: update fund_config.json via UI")

def load_quota_history():
    if not os.path.exists(QUOTA_HISTORY_FILE):
        return []
    with open(QUOTA_HISTORY_FILE, "r", encoding="utf-8") as f:
        return json.load(f)

def save_quota_history(data):
    content = json.dumps(data, ensure_ascii=False, indent=2)
    _save_with_github(QUOTA_HISTORY_FILE, "data/quota_history.json", content,
                      "chore: update quota_history.json via auto-close")

def load_pretrade_history():
    if not os.path.exists(PRETRADE_HISTORY_FILE):
        return []
    with open(PRETRADE_HISTORY_FILE, "r", encoding="utf-8") as f:
        return json.load(f)

def save_pretrade_history(data):
    content = json.dumps(data, ensure_ascii=False, indent=2)
    _save_with_github(PRETRADE_HISTORY_FILE, "data/pretrade_history.json", content,
                      "chore: update pretrade_history.json via UI")

def load_portfolio_history():
    if not os.path.exists(PORTFOLIO_HISTORY_FILE):
        return []
    with open(PORTFOLIO_HISTORY_FILE, "r", encoding="utf-8") as f:
        return json.load(f)

def save_portfolio_history(data):
    content = json.dumps(data, ensure_ascii=False, indent=2)
    # SYNC: snapshots da carteira são críticos para auditoria — não podem
    # ser perdidos em restart de container (bug observado entre 13-20/05).
    _save_with_github(PORTFOLIO_HISTORY_FILE, "data/portfolio_history.json", content,
                      "chore: update portfolio_history.json")

def get_effective_fund_config():
    """Returns fund_config, overriding quota_fechamento/data_fechamento with the last history entry."""
    config = load_fund_config()
    history = load_quota_history()
    if history:
        last = history[-1]
        config["quota_fechamento"] = last["cota_fechamento"]
        config["data_fechamento"]  = last["data"]
    return config

_VIEWER_CONFIG_DEFAULTS = {
    "tab_table":         True,
    "tab_charts":        True,
    "tab_config":        True,
    "tab_history":       True,
    "tab_risk":          True,
    "tab_events":        False,
}

def load_viewer_config():
    if not os.path.exists(VIEWER_CONFIG_FILE):
        return dict(_VIEWER_CONFIG_DEFAULTS)
    with open(VIEWER_CONFIG_FILE, "r", encoding="utf-8") as f:
        data = json.load(f)
    return {**_VIEWER_CONFIG_DEFAULTS, **data}

def save_viewer_config(data):
    content = json.dumps(data, ensure_ascii=False, indent=2)
    with open(VIEWER_CONFIG_FILE, "w", encoding="utf-8") as f:
        f.write(content)
    github_push_async("data/viewer_config.json", content, "chore: update viewer_config.json via UI")

# ---------------------------------------------------------------------------
# Price fetching — always includes ^BVSP
# ---------------------------------------------------------------------------

def fetch_prices(tickers):
    all_fetch = list(set(list(tickers) + ["^BVSP"]))
    result = {}
    for ticker in all_fetch:
        try:
            info = yf.Ticker(ticker).fast_info
            price = info.last_price
            prev  = info.previous_close
            result[ticker] = {
                "price":      round(price, 2) if price else None,
                "change_pct": round((price - prev) / prev * 100, 2) if price and prev else None,
            }
        except Exception as e:
            result[ticker] = {"price": None, "change_pct": None, "error": str(e)}
    return result

def get_cached_prices(tickers):
    now = time.time()
    if now < _price_cache["expires_at"]:
        return _price_cache["data"]
    data = fetch_prices(tickers)
    _price_cache["data"] = data
    _price_cache["expires_at"] = now + 30
    return data

def invalidate_price_cache():
    _price_cache["expires_at"] = 0

# ---------------------------------------------------------------------------
# Fundamentals (weekly cache)
# ---------------------------------------------------------------------------

def _round(val, d):
    try: return round(float(val), d) if val is not None else None
    except: return None

# BDRs: mapa ticker .SA → ticker americano subjacente (para buscar forwardPE/pegRatio)
BDR_UNDERLYING = {
    "MUTC34.SA": "MU",
    "NVDC34.SA": "NVDA",
    "A1MD34.SA": "AMD",
    "MSFT34.SA": "MSFT",
    "GOGL34.SA": "GOOGL",
    "AAPL34.SA": "AAPL",
    "AMZO34.SA": "AMZN",
    "META34.SA": "META",
    "TSLA34.SA": "TSLA",
}

def fetch_fundamentals(tickers):
    # Busca fundamentals dos subjacentes americanos de BDRs de uma vez
    us_tickers = list({BDR_UNDERLYING[t] for t in tickers if t in BDR_UNDERLYING})
    us_info = {}
    for ut in us_tickers:
        try:
            us_info[ut] = yf.Ticker(ut).info
        except Exception:
            us_info[ut] = {}

    result = {}
    for ticker in tickers:
        try:
            info = yf.Ticker(ticker).info
            dy, roe, beta, ev_ebitda = info.get("dividendYield"), info.get("returnOnEquity"), info.get("beta"), info.get("enterpriseToEbitda")
            sector_en = info.get("sector")

            # Para BDRs: complementa forwardPE e pegRatio do ticker americano subjacente
            us = us_info.get(BDR_UNDERLYING.get(ticker), {})
            forward_pe = _round(info.get("forwardPE") or us.get("forwardPE"), 2)
            peg_ratio  = _round(info.get("pegRatio")  or us.get("pegRatio"),  2)

            result[ticker] = {
                "trailing_pe":         _round(info.get("trailingPE"), 2),
                "forward_pe":          forward_pe,
                "peg_ratio":           peg_ratio,
                "price_to_book":       _round(info.get("priceToBook"), 1),
                "dividend_yield":      round(dy if dy > 1 else dy * 100, 2) if dy else None,
                "market_cap":          info.get("marketCap"),
                "fifty_two_week_high": _round(info.get("fiftyTwoWeekHigh"), 2),
                "fifty_two_week_low":  _round(info.get("fiftyTwoWeekLow"), 2),
                "short_name":          info.get("shortName") or info.get("longName"),
                "beta":                _round(beta, 2),
                "enterprise_to_ebitda": _round(ev_ebitda, 1),
                "return_on_equity":    round(roe * 100, 1) if roe is not None else None,
                "sector":              SECTOR_PT.get(sector_en, sector_en) if sector_en else None,
                "average_volume":      info.get("averageVolume"),
            }
        except Exception as e:
            result[ticker] = {"error": str(e)}
    return result

def get_cached_fundamentals(tickers):
    cache = load_cache()
    now = time.time()
    # Invalida se expirado OU se peg_ratio ainda não está no cache (campo novo)
    if any(now > cache.get(t, {}).get("expires_at", 0) or "peg_ratio" not in cache.get(t, {}) or "average_volume" not in cache.get(t, {}) for t in tickers):
        fresh = fetch_fundamentals(tickers)
        for t, d in fresh.items():
            cache[t] = {**d, "expires_at": now + FUNDAMENTALS_TTL}
        save_cache(cache)
    return {t: {k: v for k, v in cache.get(t, {}).items() if k != "expires_at"} for t in tickers}

# ---------------------------------------------------------------------------
# Portfolio history vs IBOV
# ---------------------------------------------------------------------------

def compute_portfolio_history(positions, days=90):
    import pandas as pd
    tickers  = [p["yahoo_ticker"] for p in positions]
    qty_map  = {p["yahoo_ticker"]: p["quantidade"] for p in positions}
    all_tick = tickers + ["^BVSP"]
    end   = datetime.now()
    start = end - timedelta(days=days + 45)
    try:
        df = yf.download(all_tick, start=start.strftime("%Y-%m-%d"),
                         end=end.strftime("%Y-%m-%d"), progress=False, auto_adjust=True)
    except Exception as e:
        return {"series": [], "error": str(e)}
    if df.empty:
        return {"series": []}
    close = df["Close"] if isinstance(df.columns, pd.MultiIndex) else df[["Close"]].rename(columns={"Close": all_tick[0]})
    close = close.ffill()
    port  = sum(close[t] * qty_map.get(t, 0) for t in tickers if t in close.columns)
    ibov  = close["^BVSP"] if "^BVSP" in close.columns else None
    valid = port.dropna()
    valid = valid[valid > 0].tail(days)
    if len(valid) == 0:
        return {"series": []}
    base_p = float(valid.iloc[0])
    base_i = float(ibov.loc[valid.index[0]]) if ibov is not None and valid.index[0] in ibov.index else None
    series = []
    for dt in valid.index:
        pv = float(valid[dt])
        in_ = None
        if ibov is not None and base_i and dt in ibov.index:
            iv = ibov[dt]
            if iv == iv: in_ = round(float(iv) / base_i * 100, 2)
        series.append({"date": dt.strftime("%Y-%m-%d"),
                       "portfolio": round(pv / base_p * 100, 2) if base_p else None,
                       "portfolio_abs": round(pv, 2), "ibov": in_})
    return {"series": series}

def get_cached_history(positions, days=90):
    cache = load_cache()
    now   = time.time()
    key   = f"history_{days}"
    if cache.get(key) and now < cache[key].get("expires_at", 0):
        return cache[key]["data"]
    data = compute_portfolio_history(positions, days)
    cache[key] = {"data": data, "expires_at": now + HISTORY_TTL}
    save_cache(cache)
    return data

def invalidate_history_cache():
    cache = load_cache()
    for k in list(cache.keys()):
        if k.startswith("history_"): del cache[k]
    save_cache(cache)

def compute_stock_history(yahoo_ticker, range_key):
    """Fetch price history for a single ticker + ^BVSP comparison, indexed to 100."""
    import pandas as pd
    if range_key == "YTD":
        kwargs = {"start": f"{datetime.now().year}-01-01"}
    else:
        kwargs = RANGE_TO_PERIOD.get(range_key, {"period": "1mo"})
    try:
        raw = yf.download(
            [yahoo_ticker, "^BVSP"],
            progress=False,
            auto_adjust=True,
            **kwargs,
        )
    except Exception as e:
        return {"series": [], "error": str(e)}
    if raw is None or (hasattr(raw, 'empty') and raw.empty):
        return {"series": []}
    close = raw["Close"] if isinstance(raw.columns, pd.MultiIndex) else raw
    close = close.ffill()
    if yahoo_ticker not in close.columns:
        return {"series": [], "error": "ticker not found"}
    stock = close[yahoo_ticker].dropna()
    ibov  = close["^BVSP"] if "^BVSP" in close.columns else None
    if stock.empty:
        return {"series": []}
    base_s = float(stock.iloc[0])
    base_i = float(ibov.iloc[0]) if ibov is not None and not ibov.empty else None
    try:
        fi = yf.Ticker(yahoo_ticker).fast_info
        w52_high = round(float(fi.fifty_two_week_high), 2) if fi.fifty_two_week_high else None
        w52_low  = round(float(fi.fifty_two_week_low),  2) if fi.fifty_two_week_low  else None
    except Exception:
        w52_high = w52_low = None
    series = []
    for dt in stock.index:
        sv = float(stock[dt])
        iv = None
        if ibov is not None and base_i and dt in ibov.index:
            raw_iv = float(ibov[dt])
            if raw_iv == raw_iv:  # not NaN
                iv = round(raw_iv / base_i * 100, 2)
        series.append({
            "date":    dt.strftime("%Y-%m-%d"),
            "price":   round(sv, 2),
            "indexed": round(sv / base_s * 100, 2) if base_s else None,
            "ibov":    iv,
        })
    period_return = round((float(stock.iloc[-1]) / base_s - 1) * 100, 2) if base_s else None
    ibov_return   = round((series[-1]["ibov"] / 100 - 1) * 100, 2) if series and series[-1]["ibov"] else None
    vs_ibov = round(period_return - ibov_return, 2) if period_return is not None and ibov_return is not None else None
    return {
        "series":        series,
        "ticker":        yahoo_ticker,
        "period_return": period_return,
        "ibov_return":   ibov_return,
        "vs_ibov":       vs_ibov,
        "w52_high":      w52_high,
        "w52_low":       w52_low,
        "base_price":    round(base_s, 2),
    }

def get_cached_stock_history(yahoo_ticker, range_key):
    cache = load_cache()
    now   = time.time()
    key   = f"stock_hist_{yahoo_ticker}_{range_key}"
    if cache.get(key) and now < cache[key].get("expires_at", 0):
        return cache[key]["data"]
    data = compute_stock_history(yahoo_ticker, range_key)
    cache[key] = {"data": data, "expires_at": now + STOCK_HIST_TTL}
    save_cache(cache)
    return data

# ---------------------------------------------------------------------------
# Market hours check (BRT = UTC-3, sem horário de verão desde 2019)
# ---------------------------------------------------------------------------

def _brt_now():
    """Retorna datetime atual em BRT (UTC-3, sem horário de verão).
    Use sempre que gravar timestamp visível ao usuário (histórico, audit,
    arquivos). O servidor em produção (Render) roda em UTC."""
    return datetime.utcnow() - timedelta(hours=3)

def is_market_open():
    """B3: seg-sex, 10:00–17:30 BRT."""
    brt = _brt_now()
    if brt.weekday() >= 5:          # sábado=5, domingo=6
        return False
    t = brt.hour * 60 + brt.minute  # minutos desde meia-noite
    return 10 * 60 <= t < 17 * 60 + 30

# ---------------------------------------------------------------------------
# Quota & Performance Fee calculation
# ---------------------------------------------------------------------------

def calculate_quota(rows, fund_config, prices):
    quota_fech = fund_config.get("quota_fechamento") or 0
    caixa      = fund_config.get("caixa") or 0
    proventos  = fund_config.get("proventos_a_receber") or 0
    custos     = fund_config.get("custos_provisionados") or 0
    fee_rate   = (fund_config.get("performance_fee_rate") or 20) / 100

    nav_carteira = sum(r.get("valor_liquido") or 0 for r in rows)
    nav_total    = nav_carteira + caixa + proventos - custos

    brt_now    = datetime.utcnow() - timedelta(hours=3)
    today_str  = brt_now.strftime("%Y-%m-%d")
    data_fech  = fund_config.get("data_fechamento", "")
    mercado_fechado   = not is_market_open()
    today_is_official = mercado_fechado and (data_fech == today_str)

    # Fechamento oficial já registrado hoje: exibe apenas o fechamento, sem variação intraday
    if today_is_official:
        return {
            "quota_fechamento":         quota_fech,
            "data_fechamento":          data_fech,
            "cota_estimada":            quota_fech if quota_fech else None,
            "variacao_pct":             0.0,
            "variacao_rs_por_cota":     0.0,
            "retorno_fundo_pct":        0.0,
            "retorno_ibov_pct":         0.0,
            "alpha_pct":                0.0,
            "provisao_performance_pct": 0.0,
            "provisao_performance_rs":  0.0,
            "nav_total":                round(nav_total, 2),
            "mercado_fechado":          True,
            "caixa": caixa, "proventos_a_receber": proventos,
            "custos_provisionados": custos,
        }
    # Se mercado fechado mas fechamento de hoje ainda não registrado, continua calculando
    # com os preços finais do pregão (mesmo comportamento do horário de mercado aberto).

    valid = [r for r in rows if r.get("pct_total") and r.get("var_dia_pct") is not None]
    retorno_carteira = sum((r["var_dia_pct"] / 100) * (r["pct_total"] / 100) for r in valid) if valid else 0.0

    ibov_data    = prices.get("^BVSP", {})
    ibov_ret_pct = ibov_data.get("change_pct") or 0
    ibov_ret     = ibov_ret_pct / 100

    alpha        = retorno_carteira - ibov_ret
    provisao_pct = max(0.0, alpha * fee_rate)

    provisao_rs  = provisao_pct * nav_total

    cota_est = quota_fech * (1 + retorno_carteira) if quota_fech else None

    return {
        "quota_fechamento":        quota_fech,
        "data_fechamento":         fund_config.get("data_fechamento", ""),
        "cota_estimada":           round(cota_est, 8) if cota_est else None,
        "variacao_pct":            round(retorno_carteira * 100, 4),
        "variacao_rs_por_cota":    round(cota_est - quota_fech, 8) if cota_est else None,
        "retorno_fundo_pct":       round(retorno_carteira * 100, 4),
        "retorno_ibov_pct":        round(ibov_ret * 100, 4),
        "alpha_pct":               round(alpha * 100, 4),
        "provisao_performance_pct": round(provisao_pct * 100, 4),
        "provisao_performance_rs": round(provisao_rs, 2),
        "nav_total":               round(nav_total, 2),
        "mercado_fechado":         mercado_fechado,
        "caixa": caixa, "proventos_a_receber": proventos,
        "custos_provisionados": custos,
    }

# ---------------------------------------------------------------------------
# Portfolio response builder
# ---------------------------------------------------------------------------

def build_portfolio_response(portfolio, prices, fundamentals):
    import math
    _PARTICIPATION_RATE = 0.20  # 20% do ADV — padrão de mercado para fundos

    rows = []
    total_value = 0.0
    for pos in portfolio["positions"]:
        yahoo = pos["yahoo_ticker"]
        pd_   = prices.get(yahoo, {})
        fund  = fundamentals.get(yahoo, {})
        price = pd_.get("price")
        qtde  = pos["quantidade"]
        vl    = round(price * qtde, 2) if price else None
        if vl: total_value += vl
        pa    = pos.get("preco_alvo")
        upside = round((pa / price - 1) * 100, 2) if price and pa and price > 0 else None
        mc     = fund.get("market_cap")

        avg_vol = fund.get("average_volume")
        avg_daily_vol_rs = round(avg_vol * price, 2) if avg_vol and price else None
        avg_daily_vol_mm = round(avg_daily_vol_rs / 1e6, 1) if avg_daily_vol_rs else None

        manual_score = pos.get("liq_diaria_mm")
        if manual_score is None and avg_daily_vol_rs and vl and avg_daily_vol_rs > 0:
            effective_daily = avg_daily_vol_rs * _PARTICIPATION_RATE
            days_calc = vl / effective_daily
            liq_score = round(max(-30.0, min(30.0, 10.0 - 10.0 * math.log2(max(days_calc, 0.001)))), 1)
            liq_auto  = True
        else:
            liq_score = manual_score
            liq_auto  = False

        rows.append({
            "ticker": pos["ticker"], "yahoo_ticker": yahoo,
            "categoria": pos.get("categoria", "Acao"), "quantidade": qtde,
            "liq_diaria_mm": liq_score, "liq_auto": liq_auto,
            "avg_daily_vol_mm": avg_daily_vol_mm,
            "prazo_resgate_d": pos.get("prazo_resgate_d"),
            "lucro_mi_26":   pos.get("lucro_mi_26"),
            "preco_alvo": pa, "preco": price,
            "var_dia_pct": pd_.get("change_pct"),
            "valor_liquido": vl, "upside_pct": upside,
            "trailing_pe": fund.get("trailing_pe"), "forward_pe": fund.get("forward_pe"),
            "peg_ratio": fund.get("peg_ratio"),
            "price_to_book": fund.get("price_to_book"), "dividend_yield": fund.get("dividend_yield"),
            "market_cap_bi": round(mc / 1e9, 1) if mc else None,
            "week_high": fund.get("fifty_two_week_high"), "week_low": fund.get("fifty_two_week_low"),
            "short_name": fund.get("short_name"), "beta": fund.get("beta"),
            "enterprise_to_ebitda": fund.get("enterprise_to_ebitda"),
            "return_on_equity": fund.get("return_on_equity"),
            # IMPORTANTE: não cair em `categoria` (Ação/BDR) — isso corrompe
            # o HHI setorial mostrando Ação/BDR como se fossem setores.
            "sector": fund.get("sector") or "Outros",
        })
    for r in rows:
        r["pct_total"] = round(r["valor_liquido"] / total_value * 100, 2) if r["valor_liquido"] and total_value > 0 else None

    weighted_upside = round(sum(r["upside_pct"] * r["valor_liquido"] / total_value for r in rows if r["upside_pct"] and r["valor_liquido"]), 2) if total_value > 0 else None
    beta_rows = [r for r in rows if r["beta"] is not None and r["valor_liquido"]]
    weighted_beta = round(sum(r["beta"] * r["valor_liquido"] / total_value for r in beta_rows), 2) if beta_rows and total_value > 0 else None

    def _wavg(field):
        valid = [r for r in rows if r.get(field) is not None and r["valor_liquido"]]
        wt = sum(r["valor_liquido"] for r in valid)
        if not valid or wt == 0: return None
        return round(sum(r[field] * r["valor_liquido"] / wt for r in valid), 2)

    weighted_stats = {
        "w_trailing_pe":          _wavg("trailing_pe"),
        "w_forward_pe":           _wavg("forward_pe"),
        "w_peg_ratio":            _wavg("peg_ratio"),
        "w_enterprise_to_ebitda": _wavg("enterprise_to_ebitda"),
        "w_return_on_equity":     _wavg("return_on_equity"),
        "w_beta":                 weighted_beta,
        "w_price_to_book":        _wavg("price_to_book"),
        "w_dividend_yield":       _wavg("dividend_yield"),
        "w_var_dia_pct":          _wavg("var_dia_pct"),
        "w_upside_pct":           weighted_upside,
        "w_lucro_mi_26":          sum(r["lucro_mi_26"] for r in rows if r.get("lucro_mi_26")) or None,
    }

    return {
        "fund_name": portfolio["fund_name"], "total_value": round(total_value, 2),
        "weighted_upside": weighted_upside, "weighted_beta": weighted_beta,
        "weighted_stats": weighted_stats,
        "last_price_update": _brt_now().isoformat(), "rows": rows,
    }

def _build_portfolio_snapshot(data, quota, source="auto"):
    import uuid
    now = datetime.now()
    q, ws = quota or {}, data.get("weighted_stats", {})
    return {
        "id":        str(uuid.uuid4()),
        "date":      now.strftime("%Y-%m-%d"),
        "timestamp": now.isoformat(timespec="seconds"),
        "source":    source,
        "summary": {
            "total_value":            data.get("total_value"),
            "num_positions":          len(data.get("rows", [])),
            "cota_estimada":          q.get("cota_estimada"),
            "variacao_pct":           q.get("variacao_pct"),
            "w_trailing_pe":          ws.get("w_trailing_pe"),
            "w_forward_pe":           ws.get("w_forward_pe"),
            "w_peg_ratio":            ws.get("w_peg_ratio"),
            "w_enterprise_to_ebitda": ws.get("w_enterprise_to_ebitda"),
            "w_return_on_equity":     ws.get("w_return_on_equity"),
            "w_beta":                 ws.get("w_beta"),
            "w_price_to_book":        ws.get("w_price_to_book"),
            "w_dividend_yield":       ws.get("w_dividend_yield"),
            "w_var_dia_pct":          ws.get("w_var_dia_pct"),
            "w_upside_pct":           ws.get("w_upside_pct"),
            "w_lucro_mi_26":          ws.get("w_lucro_mi_26"),
        },
        "rows": data.get("rows", []),
    }

# ---------------------------------------------------------------------------
# Export
# ---------------------------------------------------------------------------

EXPORT_HEADERS = [
    "Ativo","Categoria","Setor","% Total","Valor Líquido (R$)","Preço (R$)",
    "Var. Dia %","Quantidade","Liq. Diária (mm)",
    "P/L Trailing","P/L Forward","EV/EBITDA","ROE %","Beta",
    "Lucro mi 26","P/VPA","Div. Yield %","Mkt Cap (Bi R$)","Preço Alvo (R$)","Upside %",
]

def row_to_export(r):
    return [r["ticker"],r["categoria"],r.get("sector"),r["pct_total"],r["valor_liquido"],r["preco"],
            r["var_dia_pct"],r["quantidade"],r["liq_diaria_mm"],
            r["trailing_pe"],r["forward_pe"],r.get("enterprise_to_ebitda"),r.get("return_on_equity"),r.get("beta"),
            r["lucro_mi_26"],r["price_to_book"],r["dividend_yield"],r["market_cap_bi"],r["preco_alvo"],r["upside_pct"]]

def get_export_data():
    portfolio = load_portfolio()
    tickers   = [p["yahoo_ticker"] for p in portfolio["positions"]]
    prices    = get_cached_prices(tickers)
    funds     = get_cached_fundamentals(tickers)
    return build_portfolio_response(portfolio, prices, funds)

# ---------------------------------------------------------------------------
# Routes
# ---------------------------------------------------------------------------

# ---------------------------------------------------------------------------
# Auth
# ---------------------------------------------------------------------------

_PUBLIC = {"/login", "/logout"}

def require_admin(f):
    @wraps(f)
    def decorated(*args, **kwargs):
        if session.get("role") != "admin":
            return jsonify({"error": "forbidden"}), 403
        return f(*args, **kwargs)
    return decorated

@app.before_request
def require_login():
    if request.path in _PUBLIC:
        return
    if request.path.startswith("/static/"):
        return
    if request.path == "/api/quota-history/auto-close":
        return
    if not session.get("role"):
        return redirect(url_for("login"))


# Endpoints que alocam DataFrames pandas grandes via yf.download/.history e
# precisam que o GC libere memória imediatamente (Render Starter tem só 512MB).
_HEAVY_PATHS_PREFIX = ("/api/risk/",)
_HEAVY_PATHS_EXACT = {
    "/api/attribution", "/api/history",
    "/api/performance-chart", "/api/drawdown-volatility",
    "/api/performance-indicators", "/api/monthly-returns", "/api/annual-returns",
    "/api/index-members", "/api/events",
}

@app.after_request
def _release_heavy_memory(response):
    p = request.path
    if p.startswith(_HEAVY_PATHS_PREFIX) or p in _HEAVY_PATHS_EXACT:
        gc.collect()
    return response

@app.route("/login", methods=["GET", "POST"])
def login():
    error = None
    if request.method == "POST":
        u = request.form.get("username", "").strip()
        p = request.form.get("password", "")
        if LOGIN_PASS and u == LOGIN_USER and p == LOGIN_PASS:
            session["role"] = "admin"
            return redirect("/")
        if VIEWER_PASS and u == VIEWER_USER and p == VIEWER_PASS:
            session["role"] = "viewer"
            return redirect("/")
        error = "Usuário ou senha incorretos."
    return render_template("login.html", error=error)

@app.route("/logout")
def logout():
    session.clear()
    return redirect(url_for("login"))

# ---------------------------------------------------------------------------
# Main app
# ---------------------------------------------------------------------------

def _asset_version():
    try:
        base = os.path.dirname(os.path.abspath(__file__))
        mtimes = [
            os.path.getmtime(os.path.join(base, "static", "app.js")),
            os.path.getmtime(os.path.join(base, "static", "style.css")),
        ]
        return str(int(max(mtimes)))
    except Exception:
        return "0"


@app.route("/")
def index():
    return render_template("index.html",
                           role=session.get("role", "viewer"),
                           viewer_config=load_viewer_config(),
                           risk_methodology=METHODOLOGY,
                           asset_ver=_asset_version())

@app.route("/api/portfolio")
def api_portfolio():
    portfolio = load_portfolio()
    tickers   = [p["yahoo_ticker"] for p in portfolio["positions"]]
    prices    = get_cached_prices(tickers)
    funds     = get_cached_fundamentals(tickers)
    data      = build_portfolio_response(portfolio, prices, funds)
    # Attach quota data — always uses last closing from history as base
    fund_config  = get_effective_fund_config()
    data["quota"] = calculate_quota(data["rows"], fund_config, prices)
    return jsonify(data)

@app.route("/api/prices")
def api_prices():
    portfolio = load_portfolio()
    tickers   = [p["yahoo_ticker"] for p in portfolio["positions"]]
    prices    = get_cached_prices(tickers)
    funds     = get_cached_fundamentals(tickers)
    # Expose only the fields used in real-time refresh
    fund_slim = {t: {k: funds[t].get(k) for k in ("trailing_pe","forward_pe","peg_ratio")} for t in tickers}
    return jsonify({"prices": prices, "fundamentals": fund_slim, "timestamp": _brt_now().isoformat()})

@app.route("/api/fundamentals")
def api_fundamentals():
    portfolio = load_portfolio()
    tickers   = [p["yahoo_ticker"] for p in portfolio["positions"]]
    return jsonify({"fundamentals": get_cached_fundamentals(tickers)})

@app.route("/api/history")
def api_history():
    days = max(10, min(int(request.args.get("days", 90)), 365))
    portfolio = load_portfolio()
    return jsonify(get_cached_history(portfolio["positions"], days))

def _latest_nr_cotistas():
    """Lê o último número de cotistas (nr_cotst) do informe diário CVM.
    Retorna int ou None se cvm_daily.json não existir / estiver vazio."""
    try:
        path = os.path.join(DATA_DIR, "cvm_daily.json")
        if not os.path.exists(path):
            return None
        with open(path, "r", encoding="utf-8") as f:
            data = json.load(f)
        records = data.get("records") or []
        if not records:
            return None
        records_sorted = sorted(records, key=lambda r: r.get("dt_comptc", ""))
        return records_sorted[-1].get("nr_cotst")
    except Exception:
        return None


@app.route("/api/fund-config", methods=["GET"])
def api_get_fund_config():
    # quota_fechamento + data_fechamento sempre derivados do último quota_history
    config = get_effective_fund_config()
    # num_cotistas sempre derivado do último informe CVM oficial
    config["num_cotistas"] = _latest_nr_cotistas()
    return jsonify(config)

@app.route("/api/fund-config", methods=["POST"])
@require_admin
def api_update_fund_config():
    payload = request.json
    config  = load_fund_config()
    _string_keys = {"descricao_fundo"}
    # quota_fechamento, data_fechamento e num_cotas são read-only no formulário
    # (derivados das fontes oficiais) — qualquer valor enviado é silenciosamente ignorado.
    for key in ["caixa",
                "proventos_a_receber","custos_provisionados","performance_fee_rate",
                "performance_fee_acumulada_rs","descricao_fundo",
                "limite_concentracao_ativo_pct","limite_concentracao_setor_pct",
                "enable_concentracao_ativo","enable_concentracao_setor",
                "liquidez_min_5d_pct","liquidez_max_baixa_pct","liquidez_max_zerar_dias"]:
        if key not in payload: continue
        val = payload[key]
        if key in _string_keys:
            config[key] = val
        elif key in ("enable_concentracao_ativo", "enable_concentracao_setor"):
            config[key] = bool(val)
        elif val in (None, ""):
            config[key] = None
        else:
            try: config[key] = float(val)
            except: config[key] = val
    save_fund_config(config)
    invalidate_price_cache()
    return jsonify({"ok": True})

@app.route("/api/export/csv")
@require_admin
def api_export_csv():
    data   = get_export_data()
    output = io.StringIO()
    w      = csv.writer(output)
    w.writerow(EXPORT_HEADERS)
    for r in data["rows"]: w.writerow(row_to_export(r))
    output.seek(0)
    filename = f"harbour_fia_{_brt_now().strftime('%Y%m%d')}.csv"
    return Response("\ufeff" + output.getvalue(), mimetype="text/csv; charset=utf-8",
                    headers={"Content-Disposition": f"attachment; filename={filename}"})

@app.route("/api/export/excel")
@require_admin
def api_export_excel():
    from openpyxl import Workbook
    from openpyxl.styles import Alignment, Font, PatternFill
    data = get_export_data()
    wb = Workbook(); ws = wb.active; ws.title = "Portfólio"
    hf = Font(bold=True, color="E0E3F0", name="Calibri")
    hfill = PatternFill(start_color="1A1D27", end_color="1A1D27", fill_type="solid")
    ws.append(EXPORT_HEADERS)
    for cell in ws[1]: cell.font = hf; cell.fill = hfill; cell.alignment = Alignment(horizontal="center")
    for r in data["rows"]: ws.append(row_to_export(r))
    for i, w_ in enumerate([8,9,18,8,16,10,9,12,15,11,11,9,7,7,11,11,7,10,13,14,9], 1):
        ws.column_dimensions[ws.cell(1,i).column_letter].width = w_
    ws.row_dimensions[1].height = 22
    buf = io.BytesIO(); wb.save(buf); buf.seek(0)
    filename = f"harbour_fia_{_brt_now().strftime('%Y%m%d')}.xlsx"
    return send_file(buf, as_attachment=True, download_name=filename,
                     mimetype="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet")

@app.route("/api/export/pptx", methods=["POST"])
@require_admin
def api_export_pptx():
    from pptx import Presentation
    from pptx.util import Cm, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    import base64 as _b64, io as _io

    body          = request.json or {}
    images        = body.get("images", {})
    fund_name     = body.get("fund_name", "HARBOUR IAT FIF AÇÕES RL")
    ref_date      = body.get("ref_date", "")
    descricao     = body.get("descricao", "")
    stats         = body.get("stats", {})
    indicators    = body.get("indicators", {})
    annual_years  = body.get("annual_years", [])
    risk          = body.get("risk", {})
    concentration = body.get("concentration", {})
    liquidity     = body.get("liquidity", {})
    dist          = body.get("dist", {})

    prs = Presentation()
    prs.slide_width  = Cm(33.87)
    prs.slide_height = Cm(19.05)
    blank = prs.slide_layouts[6]

    # ── Paleta Harbour Capital (brand book 2024/01) ────────────────
    # PPTX sempre usa identidade Harbour (independente do tema da tela),
    # porque é material institucional que sai do app.
    BG     = RGBColor(0x0A, 0x0F, 0x24)  # fundo principal (mais escuro que brand p/ contraste)
    ORANGE = RGBColor(0x9A, 0xB7, 0xF9)  # cor de destaque/headbar (azul claro Harbour)
    WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
    MUTED  = RGBColor(0x82, 0x9F, 0xD9)  # texto secundário (azul médio claro)
    SURF   = RGBColor(0x11, 0x19, 0x3E)  # superfície de cards/KPIs (brand oficial)
    GREEN_ = RGBColor(0x16, 0xC4, 0x7F)  # variação positiva (teal harmoniza c/ azul)
    RED_   = RGBColor(0xEF, 0x45, 0x65)  # variação negativa (coral)
    BLACK  = RGBColor(0x00, 0x00, 0x00)

    def _bg(slide):
        s = slide.shapes.add_shape(1, 0, 0, prs.slide_width, prs.slide_height)
        s.fill.solid(); s.fill.fore_color.rgb = BG; s.line.fill.background()

    def _rect(slide, l, t, w, h, color):
        s = slide.shapes.add_shape(1, l, t, w, h)
        s.fill.solid(); s.fill.fore_color.rgb = color; s.line.fill.background()

    def _txt(slide, text, l, t, w, h, size=11, bold=False, color=WHITE,
             align=PP_ALIGN.LEFT, italic=False):
        tf = slide.shapes.add_textbox(l, t, w, h).text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]; p.alignment = align
        run = p.add_run(); run.text = str(text)
        run.font.size = Pt(size); run.font.bold = bold
        run.font.italic = italic; run.font.color.rgb = color
        # Identidade Harbour: Inter no corpo; PowerPoint cai p/ Calibri se nao tiver
        run.font.name = "Inter"

    def _img(slide, b64str, l, t, w, h):
        if not b64str: return
        raw = _b64.b64decode(b64str.split(",")[-1])
        slide.shapes.add_picture(_io.BytesIO(raw), l, t, w, h)

    def _pct_color(val):
        if not isinstance(val, (int, float)): return MUTED
        return GREEN_ if val >= 0 else RED_

    def _fmt(val, is_pct=True, decimals=2):
        if not isinstance(val, (int, float)): return "—"
        if is_pct:
            sign = "+" if val > 0 else ""
            return f"{sign}{val:.{decimals}f}%"
        return f"{val:.{decimals}f}"

    # ── Slide 1: Capa ────────────────────────────────────────────────
    s1 = prs.slides.add_slide(blank)
    _bg(s1)
    _rect(s1, 0, 0, prs.slide_width, Cm(2.0), ORANGE)
    _txt(s1, fund_name, Cm(0.6), Cm(0.2), Cm(26), Cm(1.6), size=24, bold=True, color=BLACK)
    if ref_date:
        _txt(s1, f"Ref.: {ref_date}", Cm(27), Cm(0.5), Cm(6.2), Cm(1.0),
             size=10, color=BLACK, align=PP_ALIGN.RIGHT)
    stat_items = [
        ("PATRIMÔNIO",  stats.get("aum",       "—")),
        ("COTA",        stats.get("quota",      "—")),
        ("INÍCIO",      stats.get("inception",  "—")),
        ("BENCHMARK",   stats.get("benchmark",  "IBOV")),
        ("ESTRATÉGIA",  stats.get("strategy",   "Long-Only")),
    ]
    bw = Cm(5.8)
    for i, (lbl, val) in enumerate(stat_items):
        x = Cm(0.6 + i * 6.5)
        _rect(s1, x, Cm(2.4), bw, Cm(2.0), SURF)
        _txt(s1, lbl, x+Cm(0.3), Cm(2.5), bw-Cm(0.6), Cm(0.7), size=8, color=ORANGE, bold=True)
        _txt(s1, str(val), x+Cm(0.3), Cm(3.1), bw-Cm(0.6), Cm(1.0), size=11, color=WHITE, bold=True)
    if descricao:
        _txt(s1, descricao, Cm(0.6), Cm(5.0), prs.slide_width-Cm(1.2), Cm(12.0),
             size=11, color=RGBColor(0xC6, 0xD6, 0xFC))

    # ── Slide 2: Performance ─────────────────────────────────────────
    s2 = prs.slides.add_slide(blank)
    _bg(s2)
    _rect(s2, 0, 0, prs.slide_width, Cm(1.2), ORANGE)
    _txt(s2, "PERFORMANCE ACUMULADA", Cm(0.6), Cm(0.1), Cm(20), Cm(1.0), size=14, bold=True, color=BLACK)
    if images.get("perf_chart"):
        _img(s2, images["perf_chart"], Cm(0.5), Cm(1.5), Cm(24), Cm(13))
    metric_items = [
        ("RETORNO TOTAL", indicators.get("total", {}).get("ret"),   True),
        ("NO ANO",        indicators.get("no_ano", {}).get("ret"),  True),
        ("12 MESES",      indicators.get("12m",    {}).get("ret"),  True),
        ("SHARPE TOTAL",  indicators.get("total", {}).get("sharpe"), False),
        ("VOL. ANUAL",    indicators.get("total", {}).get("vol"),   True),
    ]
    for i, (lbl, val, is_pct) in enumerate(metric_items):
        y = Cm(1.8 + i * 3.0)
        _rect(s2, Cm(25.5), y, Cm(7.8), Cm(2.6), SURF)
        _txt(s2, lbl, Cm(25.8), y+Cm(0.1), Cm(7.3), Cm(0.8), size=8, color=MUTED, bold=True)
        c = _pct_color(val) if lbl not in ("SHARPE TOTAL", "VOL. ANUAL") else WHITE
        _txt(s2, _fmt(val, is_pct), Cm(25.8), y+Cm(0.9), Cm(7.3), Cm(1.4),
             size=16, bold=True, color=c)

    # ── Slide 3: Rentabilidade Anual ──────────────────────────────────
    s3 = prs.slides.add_slide(blank)
    _bg(s3)
    _rect(s3, 0, 0, prs.slide_width, Cm(1.2), ORANGE)
    _txt(s3, "RENTABILIDADE HISTÓRICA ANUAL", Cm(0.6), Cm(0.1), Cm(30), Cm(1.0), size=14, bold=True, color=BLACK)
    hdrs = ["ANO", "FUNDO", "IBOV", "CDI", "ALPHA"]
    cws  = [Cm(3.2), Cm(5.0), Cm(5.0), Cm(5.0), Cm(5.0)]
    cx   = [Cm(1.5)]
    for w in cws[:-1]: cx.append(cx[-1] + w)
    rh = Cm(0.88)
    hy = Cm(1.5)
    _rect(s3, Cm(1.0), hy, sum(cws), rh, RGBColor(0x27, 0x32, 0x73))  # header row: brand surface escuro
    for i, h in enumerate(hdrs):
        _txt(s3, h, cx[i], hy+Cm(0.05), cws[i], rh-Cm(0.1), size=9, bold=True, color=ORANGE, align=PP_ALIGN.CENTER)
    for ri, yr in enumerate(annual_years):
        y = hy + rh*(ri+1)
        _rect(s3, Cm(1.0), y, sum(cws), rh,
              RGBColor(0x1A,0x23,0x4D) if ri%2==0 else RGBColor(0x11,0x19,0x3E))  # zebra stripes em tons brand
        row_vals = [
            (yr.get("year",""),       WHITE),
            (_fmt(yr.get("fund_year")), _pct_color(yr.get("fund_year"))),
            (_fmt(yr.get("ibov_year")), _pct_color(yr.get("ibov_year"))),
            (_fmt(yr.get("cdi_year")),  _pct_color(yr.get("cdi_year"))),
            (_fmt(yr.get("alpha")),     _pct_color(yr.get("alpha"))),
        ]
        for i, (txt, c) in enumerate(row_vals):
            _txt(s3, txt, cx[i], y+Cm(0.05), cws[i], rh-Cm(0.1), size=9, color=c, align=PP_ALIGN.CENTER)
    if annual_years:
        last = annual_years[-1]
        y = hy + rh*(len(annual_years)+1)
        _rect(s3, Cm(1.0), y, sum(cws), rh, RGBColor(0x27, 0x32, 0x73))  # linha ACUMULADO destacada
        acc_vals = [
            ("ACUMULADO",                     ORANGE),
            (_fmt(last.get("fund_accum")),    _pct_color(last.get("fund_accum"))),
            (_fmt(last.get("ibov_accum")),    _pct_color(last.get("ibov_accum"))),
            (_fmt(last.get("cdi_accum")),     _pct_color(last.get("cdi_accum"))),
            ("—",                             MUTED),
        ]
        for i, (txt, c) in enumerate(acc_vals):
            _txt(s3, txt, cx[i], y+Cm(0.05), cws[i], rh-Cm(0.1), size=9, bold=True, color=c, align=PP_ALIGN.CENTER)

    # ── Slide 4: Risco ────────────────────────────────────────────────
    s4 = prs.slides.add_slide(blank)
    _bg(s4)
    _rect(s4, 0, 0, prs.slide_width, Cm(1.2), ORANGE)
    _txt(s4, "ANÁLISE DE RISCO", Cm(0.6), Cm(0.1), Cm(20), Cm(1.0), size=14, bold=True, color=BLACK)
    risk_cards = [
        ("SHARPE (TOTAL)",   risk.get("sharpe"),           False),
        ("VOLATILIDADE ANUAL", risk.get("vol"),            True),
        ("MAX. DRAWDOWN",    risk.get("max_dd"),           False),
        ("VaR 95% (1D)",     risk.get("var_95"),           False),
        ("BETA (vs IBOV)",   risk.get("beta"),             False),
        ("TRACKING ERROR",   risk.get("tracking_error"),   True),
        ("INFORMATION RATIO", risk.get("info_ratio"),      False),
        ("CAPTURE UP",       risk.get("upside_capture"),   False),
        ("CAPTURE DOWN",     risk.get("downside_capture"), True),
    ]
    cw_ = Cm(9.8); ch_ = Cm(4.2)
    for i, (lbl, val, neutral) in enumerate(risk_cards):
        col = i % 3; row = i // 3
        x = Cm(1.0 + col * 10.8); y = Cm(1.5 + row * 4.8)
        _rect(s4, x, y, cw_, ch_, SURF)
        _txt(s4, lbl, x+Cm(0.3), y+Cm(0.2), cw_-Cm(0.6), Cm(0.8), size=8, color=MUTED, bold=True)
        is_ratio = lbl in ("BETA (vs IBOV)", "INFORMATION RATIO", "SHARPE (TOTAL)", "CAPTURE UP", "CAPTURE DOWN")
        txt = _fmt(val, is_pct=not is_ratio)
        c   = WHITE if neutral or val is None else _pct_color(val)
        if lbl in ("VOLATILIDADE ANUAL", "TRACKING ERROR"): c = WHITE
        _txt(s4, txt, x+Cm(0.3), y+Cm(1.0), cw_-Cm(0.6), Cm(2.8), size=22, bold=True, color=c)

    # ── Slide 5: Concentração e Liquidez ──────────────────────────────
    s5 = prs.slides.add_slide(blank)
    _bg(s5)
    _rect(s5, 0, 0, prs.slide_width, Cm(1.2), ORANGE)
    _txt(s5, "CONCENTRAÇÃO E LIQUIDEZ", Cm(0.6), Cm(0.1), Cm(20), Cm(1.0), size=14, bold=True, color=BLACK)
    top5 = concentration.get("top5", [])
    _txt(s5, "TOP POSIÇÕES", Cm(1.0), Cm(1.5), Cm(16), Cm(0.8), size=10, color=ORANGE, bold=True)
    for i, pos in enumerate(top5[:5]):
        y = Cm(2.5 + i * 1.8)
        _rect(s5, Cm(1.0), y, Cm(16), Cm(1.5), SURF if i%2==0 else BG)
        _txt(s5, pos.get("ticker",""), Cm(1.3), y+Cm(0.2), Cm(10), Cm(1.0), size=12, bold=True, color=WHITE)
        pct = pos.get("pct_total", 0)
        _txt(s5, f"{pct:.1f}%", Cm(13), y+Cm(0.2), Cm(3.5), Cm(1.0), size=12, color=ORANGE, align=PP_ALIGN.RIGHT)
    hhi_val = concentration.get("hhi")
    hhi_lbl = concentration.get("hhi_label", "")
    top5_pct = concentration.get("top5_pct")
    liq_items = [
        ("HHI",       f"{hhi_val:.0f} ({hhi_lbl})" if isinstance(hhi_val,(int,float)) else "—"),
        ("TOP 5",     _fmt(top5_pct) if isinstance(top5_pct,(int,float)) else "—"),
        ("LIQ. 1D",   _fmt(liquidity.get("liq_1d_pct"))),
        ("LIQ. 5D",   _fmt(liquidity.get("liq_5d_pct"))),
        ("LIQ. 10D",  _fmt(liquidity.get("liq_10d_pct"))),
    ]
    _txt(s5, "INDICADORES", Cm(19.0), Cm(1.5), Cm(14), Cm(0.8), size=10, color=ORANGE, bold=True)
    for i, (lbl, val) in enumerate(liq_items):
        y = Cm(2.5 + i * 3.0)
        _rect(s5, Cm(19.0), y, Cm(14.0), Cm(2.6), SURF)
        _txt(s5, lbl, Cm(19.3), y+Cm(0.1), Cm(13.5), Cm(0.8), size=8, color=MUTED, bold=True)
        _txt(s5, val, Cm(19.3), y+Cm(0.9), Cm(13.5), Cm(1.4), size=16, bold=True, color=WHITE)

    # ── Slide 6: Distribuição de Retornos ─────────────────────────────
    s6 = prs.slides.add_slide(blank)
    _bg(s6)
    _rect(s6, 0, 0, prs.slide_width, Cm(1.2), ORANGE)
    _txt(s6, "DISTRIBUIÇÃO DE RETORNOS DIÁRIOS", Cm(0.6), Cm(0.1), Cm(30), Cm(1.0), size=14, bold=True, color=BLACK)
    if images.get("dist_chart"):
        _img(s6, images["dist_chart"], Cm(0.5), Cm(1.5), Cm(22), Cm(14))
    dist_stats = [
        ("ASSIMETRIA",     dist.get("skewness"),     False),
        ("CURTOSE",        dist.get("kurtosis"),      False),
        ("DIAS POSITIVOS", dist.get("pct_positive"),  True),
        ("MELHOR DIA",     dist.get("best_day"),      True),
        ("PIOR DIA",       dist.get("worst_day"),     False),
    ]
    for i, (lbl, val, use_pct_color) in enumerate(dist_stats):
        y = Cm(1.8 + i * 3.2)
        _rect(s6, Cm(23.5), y, Cm(9.8), Cm(2.8), SURF)
        _txt(s6, lbl, Cm(23.8), y+Cm(0.1), Cm(9.3), Cm(0.8), size=8, color=MUTED, bold=True)
        txt = _fmt(val) if isinstance(val,(int,float)) else "—"
        c   = _pct_color(val) if use_pct_color else WHITE
        _txt(s6, txt, Cm(23.8), y+Cm(0.9), Cm(9.3), Cm(1.6), size=16, bold=True, color=c)

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    filename = f"harbour_fia_{_brt_now().strftime('%Y%m%d')}.pptx"
    return send_file(buf, as_attachment=True, download_name=filename,
                     mimetype="application/vnd.openxmlformats-officedocument.presentationml.presentation")

@app.route("/api/portfolio/update", methods=["POST"])
@require_admin
def api_update_position():
    payload = request.json
    ticker  = payload.get("ticker")
    if not ticker: return jsonify({"error": "ticker required"}), 400
    with _portfolio_write_lock:
        portfolio = load_portfolio()
        updated = False
        for pos in portfolio["positions"]:
            if pos["ticker"] == ticker:
                for field in ["quantidade","liq_diaria_mm","lucro_mi_26","preco_alvo"]:
                    if field in payload:
                        val = payload[field]
                        pos[field] = float(val) if val not in (None,"") else None
                # prazo_resgate_d (int, opcional — null = usa default por categoria)
                if "prazo_resgate_d" in payload:
                    val = payload["prazo_resgate_d"]
                    pos["prazo_resgate_d"] = int(val) if val not in (None, "") else None
                updated = True; break
        if not updated: return jsonify({"error": "ticker not found"}), 404
        save_portfolio(portfolio); invalidate_price_cache(); invalidate_history_cache()
    return jsonify({"ok": True})

@app.route("/api/portfolio/add", methods=["POST"])
@require_admin
def api_add_position():
    payload = request.json
    ticker  = payload.get("ticker","").upper().strip()
    if not ticker: return jsonify({"error": "ticker required"}), 400
    yahoo_ticker = ticker + ".SA"
    try:
        price = yf.Ticker(yahoo_ticker).fast_info.last_price
        if price is None: return jsonify({"error": f"Ticker {yahoo_ticker} não encontrado"}), 400
    except Exception as e: return jsonify({"error": str(e)}), 400
    with _portfolio_write_lock:
        portfolio = load_portfolio()
        if any(p["ticker"] == ticker for p in portfolio["positions"]):
            return jsonify({"error": "Ticker já existe na carteira"}), 409
        def _f(k): v = payload.get(k); return float(v) if v not in (None,"") else None
        portfolio["positions"].append({
            "ticker": ticker, "yahoo_ticker": yahoo_ticker,
            "categoria": payload.get("categoria","Acao"),
            "quantidade": float(payload.get("quantidade",0)),
            "liq_diaria_mm": _f("liq_diaria_mm"), "lucro_mi_26": _f("lucro_mi_26"),
            "preco_alvo": _f("preco_alvo"),
        })
        save_portfolio(portfolio); invalidate_price_cache(); invalidate_history_cache()
    return jsonify({"ok": True})

@app.route("/api/portfolio/<ticker>", methods=["DELETE"])
@require_admin
def api_delete_position(ticker):
    with _portfolio_write_lock:
        portfolio = load_portfolio()
        before = len(portfolio["positions"])
        portfolio["positions"] = [p for p in portfolio["positions"] if p["ticker"] != ticker.upper()]
        if len(portfolio["positions"]) == before: return jsonify({"error": "ticker not found"}), 404
        save_portfolio(portfolio); invalidate_price_cache(); invalidate_history_cache()
    return jsonify({"ok": True})

@app.route("/api/quota-history/auto-close", methods=["POST"])
def api_auto_close():
    """Called by GitHub Actions at 17:35 BRT on weekdays to auto-register the closing NAV."""
    # Optional secret token check
    secret = os.environ.get("AUTO_CLOSE_SECRET", "")
    if secret:
        token = request.headers.get("X-Secret", "") or request.json.get("secret", "") if request.is_json else ""
        if token != secret:
            return jsonify({"error": "unauthorized"}), 401

    today = _brt_now().strftime("%Y-%m-%d")

    portfolio = load_portfolio()
    tickers   = [p["yahoo_ticker"] for p in portfolio["positions"]]
    invalidate_price_cache()
    prices    = fetch_prices(tickers)           # fresh, no cache
    funds     = get_cached_fundamentals(tickers)
    data      = build_portfolio_response(portfolio, prices, funds)
    fund_config = get_effective_fund_config()
    quota     = calculate_quota(data["rows"], fund_config, prices)

    cota_est = quota.get("cota_estimada")
    if not cota_est:
        return jsonify({"error": "cota estimada indisponível — preços não carregados"}), 500

    history = load_quota_history()
    history = [h for h in history if h["data"] != today]
    history.append({"data": today, "cota_fechamento": round(cota_est, 8)})
    history.sort(key=lambda x: x["data"])
    save_quota_history(history)

    # ── save portfolio snapshot ──────────────────────────────────────────────
    snap = _build_portfolio_snapshot(data, quota, source="auto")
    ph   = load_portfolio_history()
    ph   = [r for r in ph if not (r.get("date") == snap["date"] and r.get("source") == "auto")]
    ph.append(snap)
    ph.sort(key=lambda x: x["timestamp"])
    save_portfolio_history(ph)
    # ────────────────────────────────────────────────────────────────────────

    # ── save liquidity snapshot ─────────────────────────────────────────────
    liq_entry = _record_liquidity_snapshot(today_str=today)
    # ────────────────────────────────────────────────────────────────────────

    return jsonify({
        "ok": True, "data": today, "cota_fechamento": round(cota_est, 8),
        "liquidity": liq_entry,
    })

@app.route("/api/performance-chart")
def api_performance_chart():
    history = load_quota_history()
    if not history:
        return jsonify({"series": []})

    cache = load_cache()
    now   = time.time()
    ibov_key       = "ibov_history_full"
    benchmarks_key = "benchmarks_history_full"

    # ── IBOV ──
    ibov_map = {}
    if cache.get(ibov_key) and now < cache[ibov_key].get("expires_at", 0):
        ibov_map = cache[ibov_key]["data"]
    else:
        start  = history[0]["data"]
        end_dt = datetime.strptime(history[-1]["data"], "%Y-%m-%d") + timedelta(days=5)
        try:
            hist = yf.Ticker("^BVSP").history(
                start=start, end=end_dt.strftime("%Y-%m-%d"), timeout=10
            )
            if not hist.empty:
                ibov_map = {str(d.date()): round(float(v), 2) for d, v in hist["Close"].items()}
        except Exception as e:
            print(f"[perf-chart] IBOV error: {e}")
        ttl = HISTORY_TTL if ibov_map else 120
        cache[ibov_key] = {"data": ibov_map, "expires_at": now + ttl}
        save_cache(cache)

    # ── Additional benchmarks: SMLL, IDIV, S&P500, NASDAQ, CDI ──
    benchmark_maps = {}
    if cache.get(benchmarks_key) and now < cache[benchmarks_key].get("expires_at", 0):
        benchmark_maps = cache[benchmarks_key]["data"]
    else:
        start  = history[0]["data"]
        end_dt = datetime.strptime(history[-1]["data"], "%Y-%m-%d") + timedelta(days=5)
        # CDI first — independent of yfinance, must not be blocked by hanging requests
        try:
            cdi_daily = load_cdi_map()
            if cdi_daily:
                cumulative = 100.0
                cdi_cum = {}
                for d in sorted(cdi_daily.keys()):
                    cumulative *= (1 + cdi_daily[d] / 100)
                    cdi_cum[d] = round(cumulative, 6)
                benchmark_maps["cdi"] = cdi_cum
        except Exception as e:
            print(f"[perf-chart] CDI error: {e}")

        # SMAL11.SA = iShares Small Cap Brasil ETF (proxy SMLL index)
        # DIVO11.SA = It Now IDIV ETF (proxy IDIV index)
        extra_tickers = {
            "^SMLL":  "SMAL11.SA",
            "^IDIV":  "DIVO11.SA",
            "^GSPC":  "^GSPC",
            "^IXIC":  "^IXIC",
        }
        for out_key, yf_ticker in extra_tickers.items():
            try:
                hist = yf.Ticker(yf_ticker).history(
                    start=start, end=end_dt.strftime("%Y-%m-%d"), timeout=10
                )
                if not hist.empty:
                    benchmark_maps[out_key] = {
                        str(d.date()): round(float(v), 2)
                        for d, v in hist["Close"].items()
                    }
            except Exception as e:
                print(f"[perf-chart] {yf_ticker} error: {e}")
        ttl = HISTORY_TTL if benchmark_maps else 120
        cache[benchmarks_key] = {"data": benchmark_maps, "expires_at": now + ttl}
        save_cache(cache)

    series = [{"date": e["data"], "fund": e["cota_fechamento"], "ibov": ibov_map.get(e["data"])}
              for e in history]
    return jsonify({"series": series, "base_date": history[0]["data"], "benchmarks": benchmark_maps})

CDI_TTL = 24 * 3600

def load_cdi_map():
    """Returns {YYYY-MM-DD: daily_rate_pct} from cache or BCB API."""
    import math, requests as req
    cache    = load_cache()
    now      = time.time()
    cdi_key  = "cdi_daily"
    history  = load_quota_history()
    if not history:
        return {}

    start_date = history[0]["data"]
    end_date   = history[-1]["data"]

    cached = cache.get(cdi_key, {})
    # Serve from cache if fresh AND covers the full range
    if cached.get("expires_at", 0) > now and cached.get("end_date") == end_date:
        return cached["data"]

    start_str = datetime.strptime(start_date, "%Y-%m-%d").strftime("%d/%m/%Y")
    end_str   = datetime.strptime(end_date,   "%Y-%m-%d").strftime("%d/%m/%Y")
    url = (f"https://api.bcb.gov.br/dados/serie/bcdata.sgs.11/dados"
           f"?formato=json&dataInicial={start_str}&dataFinal={end_str}")
    try:
        resp = req.get(url, timeout=20)
        resp.raise_for_status()
        raw = resp.json()
        if not isinstance(raw, list):
            raise ValueError("resposta inesperada do BCB")
        cdi_map = {}
        for item in raw:
            d = datetime.strptime(item["data"], "%d/%m/%Y").strftime("%Y-%m-%d")
            cdi_map[d] = float(item["valor"])
        cache[cdi_key] = {"data": cdi_map, "end_date": end_date, "expires_at": now + CDI_TTL}
        save_cache(cache)
        return cdi_map
    except Exception:
        return cached.get("data", {})

@app.route("/api/drawdown-volatility")
def api_drawdown_volatility():
    import math
    history = load_quota_history()
    if not history:
        return jsonify({"series": []})

    dates = [e["data"] for e in history]
    cotas = [e["cota_fechamento"] for e in history]

    # Drawdown: (cota / peak_so_far - 1) * 100
    peak = cotas[0]
    drawdown = []
    for c in cotas:
        if c > peak:
            peak = c
        drawdown.append(round((c / peak - 1) * 100, 2))

    # Rolling annualized volatility (21-day window)
    vol_window = 21
    rolling_vol = [None] * min(vol_window, len(cotas))
    for i in range(vol_window, len(cotas)):
        sl = cotas[i - vol_window: i + 1]
        daily = [(sl[j] / sl[j-1] - 1) for j in range(1, len(sl))]
        n = len(daily)
        mean_r = sum(daily) / n
        var = sum((r - mean_r) ** 2 for r in daily) / n
        std_d = math.sqrt(var) if var > 0 else 0
        rolling_vol.append(round(std_d * math.sqrt(252) * 100, 2))

    series = [{"date": d, "drawdown": dd, "vol": v}
              for d, dd, v in zip(dates, drawdown, rolling_vol)]

    return jsonify({"series": series})

@app.route("/api/performance-indicators")
def api_performance_indicators():
    import math
    history = load_quota_history()
    if len(history) < 2:
        return jsonify({"data": {}})

    entries   = [(datetime.strptime(e["data"], "%Y-%m-%d"), e["cota_fechamento"]) for e in history]
    last_date = entries[-1][0]
    cdi_map   = load_cdi_map()

    def cdi_ann_for_window(start_dt, end_dt):
        """Annualized CDI for the date window (exclusive start, inclusive end)."""
        rates = [v for d, v in cdi_map.items()
                 if start_dt < datetime.strptime(d, "%Y-%m-%d") <= end_dt]
        if not rates:
            return None
        compound = 1.0
        for r in rates:
            compound *= (1 + r / 100)
        n = len(rates)
        return ((compound) ** (252 / n) - 1) * 100

    def compute_metrics(sl):
        if len(sl) < 2:
            return {"ret": None, "vol": None, "sharpe": None}
        dates = [d for d, _ in sl]
        cotas = [c for _, c in sl]
        ret_cum = (cotas[-1] / cotas[0] - 1) * 100
        daily = [(cotas[i] / cotas[i-1] - 1) for i in range(1, len(cotas))]
        if len(daily) < 2:
            return {"ret": round(ret_cum, 2), "vol": None, "sharpe": None}
        n      = len(daily)
        mean_r = sum(daily) / n
        var    = sum((r - mean_r) ** 2 for r in daily) / n
        std_d  = math.sqrt(var) if var > 0 else 0
        vol_ann = std_d * math.sqrt(252) * 100
        ret_ann = ((cotas[-1] / cotas[0]) ** (252 / n) - 1) * 100
        if vol_ann > 0:
            cdi = cdi_ann_for_window(dates[0], dates[-1])
            sharpe = round((ret_ann - (cdi or 0)) / vol_ann, 2)
        else:
            sharpe = None
        return {
            "ret":    round(ret_cum, 2),
            "vol":    round(vol_ann, 2) if std_d > 0 else None,
            "sharpe": sharpe,
        }

    def get_window(months_back=None, ytd=False, no_mes=False):
        if no_mes:
            cutoff = datetime(last_date.year, last_date.month, 1)
        elif ytd:
            cutoff = datetime(last_date.year, 1, 1)
        elif months_back:
            cutoff = last_date - timedelta(days=int(months_back * 365.25 / 12))
        else:
            return entries
        before = [(d, c) for d, c in entries if d < cutoff]
        after  = [(d, c) for d, c in entries if d >= cutoff]
        return ([before[-1]] if before else []) + after

    windows = [
        ("no_mes", get_window(no_mes=True)),
        ("no_ano", get_window(ytd=True)),
        ("3m",     get_window(months_back=3)),
        ("6m",     get_window(months_back=6)),
        ("12m",    get_window(months_back=12)),
        ("24m",    get_window(months_back=24)),
        ("36m",    get_window(months_back=36)),
        ("48m",    get_window(months_back=48)),
        ("60m",    get_window(months_back=60)),
        ("total",  entries),
    ]

    return jsonify({"data": {k: compute_metrics(sl) for k, sl in windows}})

@app.route("/api/monthly-returns")
def api_monthly_returns():
    history = load_quota_history()
    if not history:
        return jsonify({"years": []})

    # Load IBOV history (reuse from performance-chart cache)
    cache = load_cache()
    now   = time.time()
    ibov_key = "ibov_history_full"

    ibov_map = {}
    if cache.get(ibov_key) and now < cache[ibov_key].get("expires_at", 0):
        ibov_map = cache[ibov_key]["data"]
    else:
        start  = history[0]["data"]
        end_dt = datetime.strptime(history[-1]["data"], "%Y-%m-%d") + timedelta(days=5)
        try:
            hist = yf.Ticker("^BVSP").history(
                start=start, end=end_dt.strftime("%Y-%m-%d"), timeout=10
            )
            if not hist.empty:
                ibov_map = {str(d.date()): round(float(v), 2) for d, v in hist["Close"].items()}
        except Exception as e:
            print(f"[monthly-returns] IBOV error: {e}")
            ibov_map = {}
        ttl = HISTORY_TTL if ibov_map else 120
        cache[ibov_key] = {"data": ibov_map, "expires_at": now + ttl}
        save_cache(cache)

    # Build month-end maps: "YYYY-MM" -> last closing value of that month
    month_map_fund = {}
    for e in sorted(history, key=lambda x: x["data"]):
        month_map_fund[e["data"][:7]] = e["cota_fechamento"]

    month_map_ibov = {}
    for date_str in sorted(ibov_map):
        month_map_ibov[date_str[:7]] = ibov_map[date_str]

    inception_date = history[0]["data"]
    inception_cota = history[0]["cota_fechamento"]

    # IBOV value on or before inception date
    inception_ibov = None
    for d in sorted(ibov_map):
        if d <= inception_date:
            inception_ibov = ibov_map[d]

    all_ym = sorted(month_map_fund.keys())
    years  = sorted(set(ym[:4] for ym in all_ym))
    mnums  = ["01","02","03","04","05","06","07","08","09","10","11","12"]

    result = []
    for year in years:
        fund_months = {}
        ibov_months = {}

        for mn in mnums:
            ym      = f"{year}-{mn}"
            m       = int(mn)
            prev_ym = f"{int(year)-1}-12" if m == 1 else f"{year}-{str(m-1).zfill(2)}"

            fc = month_map_fund.get(ym)
            fp = month_map_fund.get(prev_ym)
            ic = month_map_ibov.get(ym)
            ip = month_map_ibov.get(prev_ym)

            fund_months[mn] = round((fc / fp - 1) * 100, 2) if fc and fp else None
            ibov_months[mn] = round((ic / ip - 1) * 100, 2) if ic and ip else None

        # Annual return: last December of previous year as base (inception for first year)
        prev_dec      = f"{int(year)-1}-12"
        fund_yr_start = month_map_fund.get(prev_dec) or inception_cota
        ibov_yr_start = month_map_ibov.get(prev_dec) or inception_ibov

        last_ym     = max((ym for ym in all_ym if ym.startswith(year)), default=None)
        fund_yr_end = month_map_fund.get(last_ym)
        ibov_yr_end = month_map_ibov.get(last_ym)

        fund_year  = round((fund_yr_end / fund_yr_start - 1) * 100, 2) if fund_yr_end and fund_yr_start else None
        ibov_year  = round((ibov_yr_end / ibov_yr_start - 1) * 100, 2) if ibov_yr_end and ibov_yr_start else None
        fund_accum = round((fund_yr_end / inception_cota - 1) * 100, 2) if fund_yr_end else None
        ibov_accum = round((ibov_yr_end / inception_ibov - 1) * 100, 2) if ibov_yr_end and inception_ibov else None

        result.append({
            "year": year,
            "fund_months": fund_months, "ibov_months": ibov_months,
            "fund_year": fund_year, "ibov_year": ibov_year,
            "fund_accum": fund_accum, "ibov_accum": ibov_accum,
        })

    return jsonify({"years": result, "inception_date": inception_date})

@app.route("/api/annual-returns")
def api_annual_returns():
    history = load_quota_history()
    if not history:
        return jsonify({"years": [], "inception_date": None})

    cache    = load_cache()
    now      = time.time()
    ibov_key = "ibov_history_full"

    ibov_map = {}
    if cache.get(ibov_key) and now < cache[ibov_key].get("expires_at", 0):
        ibov_map = cache[ibov_key]["data"]
    else:
        start  = history[0]["data"]
        end_dt = datetime.strptime(history[-1]["data"], "%Y-%m-%d") + timedelta(days=5)
        try:
            hist = yf.Ticker("^BVSP").history(
                start=start, end=end_dt.strftime("%Y-%m-%d"), timeout=10)
            if not hist.empty:
                ibov_map = {str(d.date()): round(float(v), 2)
                            for d, v in hist["Close"].items()}
        except Exception as e:
            print(f"[annual-returns] IBOV error: {e}")
        ttl = HISTORY_TTL if ibov_map else 120
        cache[ibov_key] = {"data": ibov_map, "expires_at": now + ttl}
        save_cache(cache)

    cdi_map = load_cdi_map()

    year_end_fund = {}
    for e in sorted(history, key=lambda x: x["data"]):
        year_end_fund[e["data"][:4]] = e["cota_fechamento"]

    year_end_ibov = {}
    for date_str in sorted(ibov_map):
        year_end_ibov[date_str[:4]] = ibov_map[date_str]

    inception_date = history[0]["data"]
    inception_cota = history[0]["cota_fechamento"]
    inception_ibov = None
    for d in sorted(ibov_map):
        if d <= inception_date:
            inception_ibov = ibov_map[d]

    years = sorted(year_end_fund.keys())
    result = []
    for year in years:
        prev_year  = str(int(year) - 1)
        fund_start = year_end_fund.get(prev_year) or inception_cota
        fund_end   = year_end_fund.get(year)
        ibov_start = year_end_ibov.get(prev_year) or inception_ibov
        ibov_end   = year_end_ibov.get(year)

        cdi_rates = [v for d, v in cdi_map.items() if d.startswith(year)]
        if cdi_rates:
            cdi_cum = 1.0
            for r in cdi_rates:
                cdi_cum *= (1 + r / 100)
            cdi_year = round((cdi_cum - 1) * 100, 2)
        else:
            cdi_year = None

        fund_year  = round((fund_end / fund_start - 1) * 100, 2) if fund_end and fund_start else None
        ibov_year  = round((ibov_end / ibov_start - 1) * 100, 2) if ibov_end and ibov_start else None
        fund_accum = round((fund_end / inception_cota - 1) * 100, 2) if fund_end else None
        ibov_accum = round((ibov_end / inception_ibov - 1) * 100, 2) if ibov_end and inception_ibov else None

        cdi_accum_rates = [v for d, v in cdi_map.items() if d >= inception_date and d[:4] <= year]
        if cdi_accum_rates:
            cum = 1.0
            for r in cdi_accum_rates:
                cum *= (1 + r / 100)
            cdi_accum = round((cum - 1) * 100, 2)
        else:
            cdi_accum = None

        alpha = round(fund_year - ibov_year, 2) if fund_year is not None and ibov_year is not None else None

        result.append({
            "year":       year,
            "fund_year":  fund_year,
            "ibov_year":  ibov_year,
            "cdi_year":   cdi_year,
            "alpha":      alpha,
            "fund_accum": fund_accum,
            "ibov_accum": ibov_accum,
            "cdi_accum":  cdi_accum,
        })

    return jsonify({"years": result, "inception_date": inception_date})

@app.route("/api/stock-history/<path:ticker>")
def api_stock_history(ticker):
    """GET /api/stock-history/PRIO3.SA?range=1M"""
    range_key = request.args.get("range", "1M").upper()
    if range_key not in ("1S", "1M", "3M", "6M", "YTD", "1A"):
        range_key = "1M"
    return jsonify(get_cached_stock_history(ticker, range_key))

@app.route("/api/quota-history", methods=["GET"])
def api_get_quota_history():
    return jsonify(load_quota_history())

@app.route("/api/quota-history", methods=["POST"])
@require_admin
def api_add_quota_history():
    payload = request.json
    data_str = payload.get("data", "").strip()
    cota     = payload.get("cota_fechamento")
    if not data_str or cota is None:
        return jsonify({"error": "data e cota_fechamento são obrigatórios"}), 400
    try:
        cota = float(cota)
    except (TypeError, ValueError):
        return jsonify({"error": "cota_fechamento inválida"}), 400
    history = load_quota_history()
    history = [h for h in history if h["data"] != data_str]
    history.append({"data": data_str, "cota_fechamento": cota})
    history.sort(key=lambda x: x["data"])
    save_quota_history(history)
    return jsonify({"ok": True})

@app.route("/api/quota-history/<date>", methods=["DELETE"])
@require_admin
def api_delete_quota_history(date):
    history = load_quota_history()
    new_history = [h for h in history if h["data"] != date]
    if len(new_history) == len(history):
        return jsonify({"error": "entrada não encontrada"}), 404
    save_quota_history(new_history)
    return jsonify({"ok": True})

@app.route("/api/viewer-config", methods=["GET"])
def api_get_viewer_config():
    return jsonify(load_viewer_config())

@app.route("/api/viewer-config", methods=["POST"])
@require_admin
def api_save_viewer_config():
    payload = request.json
    config = load_viewer_config()
    for key in _VIEWER_CONFIG_DEFAULTS:
        if key in payload:
            config[key] = bool(payload[key])
    save_viewer_config(config)
    return jsonify({"ok": True})

# ─────────────────────────────────────────────────────────────────────────────
# Attribution
# ─────────────────────────────────────────────────────────────────────────────

@app.route("/api/attribution")
def api_attribution():
    import pandas as pd
    period    = request.args.get("period", "month")
    portfolio = load_portfolio()
    positions = portfolio["positions"]
    tickers   = [p["yahoo_ticker"] for p in positions]
    prices    = get_cached_prices(tickers)

    total_value = sum(
        (prices.get(p["yahoo_ticker"], {}).get("price") or 0) * p["quantidade"]
        for p in positions
    )

    ibov_ret = None
    rows     = []

    if period == "day":
        ibov_ret = (prices.get("^BVSP", {}).get("change_pct") or 0)
        for pos in positions:
            t    = pos["yahoo_ticker"]
            p    = prices.get(t, {}).get("price") or 0
            vl   = p * pos["quantidade"]
            ret  = prices.get(t, {}).get("change_pct") or 0
            peso = vl / total_value * 100 if total_value > 0 else 0
            rows.append({
                "ticker":          pos["ticker"],
                "retorno_pct":     round(ret, 2),
                "peso_pct":        round(peso, 2),
                "contribuicao_pct": round(ret * peso / 100, 3),
                "contribuicao_bps": round(ret * peso, 1),
            })
    else:
        now_dt = datetime.now()
        if period == "week":
            start_dt = now_dt - timedelta(days=8)
        elif period == "month":
            start_dt = datetime(now_dt.year, now_dt.month, 1) - timedelta(days=1)
        elif period == "ytd":
            start_dt = datetime(now_dt.year, 1, 1) - timedelta(days=1)
        else:
            start_dt = now_dt - timedelta(days=30)
        try:
            all_tick = tickers + ["^BVSP"]
            df = yf.download(all_tick, start=start_dt.strftime("%Y-%m-%d"),
                             end=(now_dt + timedelta(days=1)).strftime("%Y-%m-%d"),
                             auto_adjust=True, progress=False)
            if df.empty:
                return jsonify({"rows": [], "error": "sem dados históricos"})
            close = df["Close"] if isinstance(df.columns, pd.MultiIndex) else df
            clean = close.dropna(how="all")
            if len(clean) < 2:
                return jsonify({"rows": [], "error": "dados insuficientes"})
            first = clean.iloc[0]
            last  = clean.iloc[-1]
            ibov_s = float(first.get("^BVSP") or 0)
            ibov_e = float(last.get("^BVSP") or 0)
            ibov_ret = round((ibov_e / ibov_s - 1) * 100, 2) if ibov_s > 0 else None
            for pos in positions:
                t    = pos["yahoo_ticker"]
                s    = float(first.get(t) or 0)
                e    = float(last.get(t) or 0)
                ret  = round((e / s - 1) * 100, 2) if s > 0 else 0
                p    = prices.get(t, {}).get("price") or 0
                vl   = p * pos["quantidade"]
                peso = vl / total_value * 100 if total_value > 0 else 0
                rows.append({
                    "ticker":          pos["ticker"],
                    "retorno_pct":     ret,
                    "peso_pct":        round(peso, 2),
                    "contribuicao_pct": round(ret * peso / 100, 3),
                    "contribuicao_bps": round(ret * peso, 1),
                })
        except Exception as e:
            return jsonify({"rows": [], "error": str(e)})

    rows.sort(key=lambda x: x["contribuicao_pct"], reverse=True)
    total_fundo = round(sum(r["contribuicao_pct"] for r in rows), 3)
    alpha = round(total_fundo - ibov_ret, 2) if ibov_ret is not None else None
    return jsonify({
        "rows":            rows,
        "total_fundo_pct": total_fundo,
        "ibov_ret_pct":    round(ibov_ret, 2) if ibov_ret is not None else None,
        "alpha_pct":       alpha,
        "period":          period,
    })


# ─────────────────────────────────────────────────────────────────────────────
# Risk Analytics
# ─────────────────────────────────────────────────────────────────────────────

RISK_TTL = 4 * 3600

STRESS_SCENARIOS = {
    "covid":      {"label": "COVID Crash",    "ibov_shock": -0.4566, "brl_shock":  0.255,  "description": "IBOV -45.7%, BRL +25.5% (fev-mar 2020)"},
    "joesley":    {"label": "Joesley Day",    "ibov_shock": -0.088,  "brl_shock":  0.085,  "description": "IBOV -8.8%, BRL +8.5% (17 mai 2017)"},
    "lula_elei":  {"label": "Eleição Lula",   "ibov_shock": -0.061,  "brl_shock":  0.037,  "description": "IBOV -6.1%, BRL +3.7% (30 out 2022)"},
    "dilma":      {"label": "Crise Dilma",    "ibov_shock": -0.128,  "brl_shock":  0.468,  "description": "IBOV -12.8%, BRL +46.8% (jan-set 2015)"},
}


def _liq_days_from_score(score):
    """Convert liq_diaria_mm score (-30 to +30) to estimated days to liquidate full position."""
    import math
    if score is None:
        return None
    return max(0.5, 2 ** ((-float(score) + 10) / 10))


def _compute_component_var_by_beta(rows, total_value, nav, portfolio_var_1d):
    if not total_value or not nav:
        return []
    rows_v = [r for r in rows if r.get("beta") is not None and r.get("valor_liquido")]
    if not rows_v:
        return []
    w_beta = sum(r["beta"] * r["valor_liquido"] / total_value for r in rows_v)
    if not w_beta:
        return []
    out = []
    for r in rows:
        w           = (r.get("valor_liquido") or 0) / total_value
        beta        = r.get("beta") or 0
        contrib_pct = (w * beta / w_beta * 100) if w_beta else 0
        var_rs      = (w * beta / w_beta * portfolio_var_1d * nav) if w_beta else 0
        out.append({
            "ticker":      r["ticker"],
            "weight_pct":  round(w * 100, 2),
            "beta":        beta,
            "contrib_pct": round(contrib_pct, 2),
            "var_1d_rs":   round(var_rs, 2),
        })
    out.sort(key=lambda x: x["contrib_pct"], reverse=True)
    return out


def _calcular_concentracao_pretrade(rows, total_value):
    """Calcula concentração por ativo e por setor a partir de rows já processados."""
    if not total_value:
        return {"por_ativo": {}, "por_setor": {}, "hhi": 0}
    por_ativo = {}
    sector_map = {}
    for r in rows:
        vl     = r.get("valor_liquido") or 0
        ytk    = r.get("yahoo_ticker") or r.get("ticker", "")
        sector = r.get("sector") or "Outros"
        por_ativo[ytk] = round(vl / total_value * 100, 4)
        sector_map[sector] = sector_map.get(sector, 0.0) + vl
    por_setor = {s: round(v / total_value * 100, 4) for s, v in sector_map.items()}
    hhi = int(round(sum((v / total_value) ** 2 for v in sector_map.values()) * 10000))
    return {"por_ativo": por_ativo, "por_setor": por_setor, "hhi": hhi}


@app.route("/api/risk/var")
def api_risk_var():
    import math
    window    = int(request.args.get("window", 252))
    cache     = load_cache()
    now       = time.time()
    cache_key = f"risk_var_{window}"
    if cache.get(cache_key) and now < cache[cache_key].get("expires_at", 0):
        return jsonify(cache[cache_key]["data"])

    history = load_quota_history()
    if len(history) < 22:
        return jsonify({"error": "Histórico insuficiente (< 22 dias)"}), 400

    fund_config = get_effective_fund_config()
    portfolio   = load_portfolio()
    tickers     = [p["yahoo_ticker"] for p in portfolio["positions"]]
    prices      = get_cached_prices(tickers)
    funds       = get_cached_fundamentals(tickers)
    pdata       = build_portfolio_response(portfolio, prices, funds)
    nav = (pdata.get("total_value") or 0) + (fund_config.get("caixa") or 0) + (fund_config.get("proventos_a_receber") or 0)

    cotas   = [e["cota_fechamento"] for e in history]
    rets    = [(cotas[i] / cotas[i - 1] - 1) for i in range(1, len(cotas))]
    rets    = rets[-window:]
    n       = len(rets)
    if n < 10:
        return jsonify({"error": "Retornos insuficientes"}), 400

    sorted_r = sorted(rets)
    idx_95   = max(1, int(math.floor(n * 0.05)))
    idx_99   = max(1, int(math.floor(n * 0.01)))

    var_95  = abs(sorted_r[idx_95 - 1])
    var_99  = abs(sorted_r[idx_99 - 1])
    cvar_95 = abs(sum(sorted_r[:idx_95]) / idx_95)
    cvar_99 = abs(sum(sorted_r[:max(1, idx_99)]) / max(1, idx_99))

    def _rs(v):  return round(v * nav, 2)
    def _pct(v): return round(v * 100, 3)

    result = {
        "window_days": window, "n_obs": n, "nav_ref": round(nav, 2),
        "var_95_1d_pct":  _pct(var_95),
        "var_99_1d_pct":  _pct(var_99),
        "cvar_95_1d_pct": _pct(cvar_95),
        "cvar_99_1d_pct": _pct(cvar_99),
        "var_95_10d_pct": _pct(var_95 * math.sqrt(10)),
        "var_99_10d_pct": _pct(var_99 * math.sqrt(10)),
        "var_95_1d_rs":   _rs(var_95),
        "var_99_1d_rs":   _rs(var_99),
        "cvar_95_1d_rs":  _rs(cvar_95),
        "cvar_99_1d_rs":  _rs(cvar_99),
        "var_95_10d_rs":  _rs(var_95 * math.sqrt(10)),
        "var_99_10d_rs":  _rs(var_99 * math.sqrt(10)),
        "component_var": _compute_component_var_by_beta(
            pdata["rows"], pdata.get("total_value", 0), nav, var_95),
        "return_distribution": {
            "mean_pct":           _pct(sum(rets) / n),
            "best_day":           _pct(max(rets)),
            "worst_day":          _pct(min(rets)),
            "positive_days_pct":  round(sum(1 for r in rets if r > 0) / n * 100, 1),
        },
    }
    cache[cache_key] = {"data": result, "expires_at": now + RISK_TTL}
    save_cache(cache)
    return jsonify(result)


@app.route("/api/risk/stress")
def api_risk_stress():
    scenario_key = request.args.get("scenario", "")
    custom_ibov  = request.args.get("ibov_shock")
    custom_brl   = request.args.get("brl_shock")

    portfolio   = load_portfolio()
    tickers     = [p["yahoo_ticker"] for p in portfolio["positions"]]
    prices      = get_cached_prices(tickers)
    funds       = get_cached_fundamentals(tickers)
    pdata       = build_portfolio_response(portfolio, prices, funds)
    fund_config = get_effective_fund_config()
    nav         = (pdata.get("total_value") or 0) + (fund_config.get("caixa") or 0) + (fund_config.get("proventos_a_receber") or 0)
    total_value = pdata.get("total_value") or 0

    def run_scenario(ibov_shock, brl_shock, label, description):
        rows_out   = []
        port_impact = 0.0
        for r in pdata["rows"]:
            w        = (r.get("valor_liquido") or 0) / total_value if total_value else 0
            beta     = r.get("beta") or 1.0
            is_bdr   = r.get("categoria", "").upper() == "BDR"
            # BRL depreciation (positive brl_shock) → BDR gains in BRL terms
            stock_imp   = beta * ibov_shock + (brl_shock if is_bdr else 0)
            pos_imp_rs  = stock_imp * (r.get("valor_liquido") or 0)
            port_impact += stock_imp * w
            rows_out.append({
                "ticker":     r["ticker"],
                "categoria":  r.get("categoria"),
                "weight_pct": round(w * 100, 2),
                "beta":       round(beta, 2),
                "impact_pct": round(stock_imp * 100, 2),
                "impact_rs":  round(pos_imp_rs, 2),
            })
        rows_out.sort(key=lambda x: x["impact_rs"])
        return {
            "label": label, "description": description,
            "ibov_shock_pct":       round(ibov_shock * 100, 2),
            "brl_shock_pct":        round(brl_shock * 100, 2),
            "portfolio_impact_pct": round(port_impact * 100, 2),
            "portfolio_impact_rs":  round(port_impact * nav, 2),
            "nav_ref": round(nav, 2),
            "positions": rows_out,
        }

    if custom_ibov is not None:
        try:
            ib  = float(custom_ibov) / 100
            brl = float(custom_brl or 0) / 100
        except ValueError:
            return jsonify({"error": "Valores inválidos"}), 400
        return jsonify(run_scenario(ib, brl, "Cenário Personalizado",
                                    f"IBOV {ib * 100:+.1f}%, BRL {brl * 100:+.1f}%"))

    if not scenario_key or scenario_key not in STRESS_SCENARIOS:
        return jsonify({"scenarios": {k: {"label": v["label"], "description": v["description"]}
                                       for k, v in STRESS_SCENARIOS.items()}})

    sc = STRESS_SCENARIOS[scenario_key]
    return jsonify(run_scenario(sc["ibov_shock"], sc["brl_shock"], sc["label"], sc["description"]))


@app.route("/api/risk/correlation")
def api_risk_correlation():
    import math
    window = int(request.args.get("window", 60))
    cache  = load_cache()
    now    = time.time()
    key    = f"risk_corr_{window}"
    if cache.get(key) and now < cache[key].get("expires_at", 0):
        return jsonify(cache[key]["data"])

    portfolio = load_portfolio()
    tickers   = [p["yahoo_ticker"] for p in portfolio["positions"]]
    all_t     = tickers + ["^BVSP"]
    end_dt    = datetime.now()
    start_dt  = end_dt - timedelta(days=int(window * 1.9) + 10)

    try:
        import pandas as pd
        df = yf.download(all_t, start=start_dt.strftime("%Y-%m-%d"),
                         end=end_dt.strftime("%Y-%m-%d"), progress=False, auto_adjust=True)
        if df.empty:
            return jsonify({"error": "Sem dados"}), 400
        close   = df["Close"] if isinstance(df.columns, pd.MultiIndex) else df
        returns = close.ffill().pct_change().dropna().tail(window)
        if len(returns) < 10:
            return jsonify({"error": "Dados insuficientes"}), 400
        corr    = returns.corr()
        cols    = [t for t in all_t if t in corr.columns]
        lmap    = {**{p["yahoo_ticker"]: p["ticker"] for p in portfolio["positions"]}, "^BVSP": "IBOV"}
        matrix  = []
        for r in cols:
            row = []
            for c in cols:
                v = corr.loc[r, c] if r in corr.index and c in corr.columns else None
                try:
                    row.append(round(float(v), 3) if v is not None and not math.isnan(float(v)) else None)
                except Exception:
                    row.append(None)
            matrix.append(row)
        result = {
            "labels": [lmap.get(t, t) for t in cols],
            "tickers": cols,
            "matrix": matrix,
            "window_days": window,
            "n_obs": len(returns),
        }
        cache[key] = {"data": result, "expires_at": now + RISK_TTL}
        save_cache(cache)
        return jsonify(result)
    except Exception as e:
        return jsonify({"error": str(e)}), 500


@app.route("/api/risk/attribution")
def api_risk_attribution():
    import math
    window = int(request.args.get("window", 60))
    cache  = load_cache()
    now    = time.time()
    key    = f"risk_attr_{window}"
    if cache.get(key) and now < cache[key].get("expires_at", 0):
        return jsonify(cache[key]["data"])

    portfolio = load_portfolio()
    tickers   = [p["yahoo_ticker"] for p in portfolio["positions"]]
    qty_map   = {p["yahoo_ticker"]: p["quantidade"] for p in portfolio["positions"]}
    end_dt    = datetime.now()
    start_dt  = end_dt - timedelta(days=int(window * 1.9) + 10)

    try:
        import pandas as pd
        df = yf.download(tickers, start=start_dt.strftime("%Y-%m-%d"),
                         end=end_dt.strftime("%Y-%m-%d"), progress=False, auto_adjust=True)
        if df.empty:
            return jsonify({"error": "Sem dados"}), 400
        close   = df["Close"] if isinstance(df.columns, pd.MultiIndex) else df
        returns = close.ffill().pct_change().dropna().tail(window)
        if len(returns) < 10:
            return jsonify({"error": "Dados insuficientes"}), 400

        prices      = get_cached_prices(tickers)
        total_value = sum((prices.get(t, {}).get("price") or 0) * qty_map.get(t, 0) for t in tickers)
        w_map       = {t: ((prices.get(t, {}).get("price") or 0) * qty_map.get(t, 0)) / total_value
                       for t in tickers} if total_value else {t: 0 for t in tickers}

        avail    = [t for t in tickers if t in returns.columns]
        port_ret = sum(returns[t] * w_map.get(t, 0) for t in avail)
        port_vol = float(port_ret.std() * math.sqrt(252) * 100)
        var_p    = float(port_ret.var())
        tlabels  = {p["yahoo_ticker"]: p["ticker"] for p in portfolio["positions"]}

        rows_out = []
        for t in avail:
            s            = returns[t]
            w            = w_map.get(t, 0)
            cov          = float(s.cov(port_ret))
            corr_p       = round(float(s.corr(port_ret)), 3)
            contrib_pct  = round(w * cov / var_p * 100, 2) if var_p > 0 else 0
            rows_out.append({
                "ticker":          tlabels.get(t, t),
                "weight_pct":      round(w * 100, 2),
                "vol_ind_pct":     round(float(s.std() * math.sqrt(252) * 100), 2),
                "corr_port":       corr_p,
                "contrib_risk_pct": contrib_pct,
                "contrib_vol_ppt": round(contrib_pct / 100 * port_vol, 2),
            })
        rows_out.sort(key=lambda x: x["contrib_risk_pct"], reverse=True)
        result = {
            "portfolio_vol_pct": round(port_vol, 2),
            "window_days": window,
            "n_obs": len(returns),
            "rows": rows_out,
        }
        cache[key] = {"data": result, "expires_at": now + RISK_TTL}
        save_cache(cache)
        return jsonify(result)
    except Exception as e:
        return jsonify({"error": str(e)}), 500


@app.route("/api/risk/rolling-beta")
def api_risk_rolling_beta():
    roll_w = int(request.args.get("roll_window", 60))
    cache  = load_cache()
    now    = time.time()
    key    = f"risk_rbeta_{roll_w}"
    if cache.get(key) and now < cache[key].get("expires_at", 0):
        return jsonify(cache[key]["data"])

    history = load_quota_history()
    if len(history) < roll_w + 5:
        return jsonify({"error": "Histórico insuficiente"}), 400

    start  = history[0]["data"]
    end_dt = datetime.strptime(history[-1]["data"], "%Y-%m-%d") + timedelta(days=5)
    try:
        ibov_h = yf.Ticker("^BVSP").history(
            start=start, end=end_dt.strftime("%Y-%m-%d"), timeout=15)
        if ibov_h.empty:
            return jsonify({"error": "Sem dados do IBOV"}), 400
        ibov_map   = {str(d.date()): float(v) for d, v in ibov_h["Close"].items()}
        ibov_dates = sorted(ibov_map)
        ibov_rets  = {ibov_dates[i]: ibov_map[ibov_dates[i]] / ibov_map[ibov_dates[i - 1]] - 1
                      for i in range(1, len(ibov_dates))}

        fund_map = {}
        prev = history[0]["cota_fechamento"]
        for e in history[1:]:
            fund_map[e["data"]] = e["cota_fechamento"] / prev - 1
            prev = e["cota_fechamento"]

        aligned = sorted(d for d in fund_map if d in ibov_rets)
        f_rets  = [fund_map[d] for d in aligned]
        i_rets  = [ibov_rets[d] for d in aligned]

        if len(aligned) < roll_w + 2:
            return jsonify({"error": "Dados insuficientes"}), 400

        series = []
        for i in range(roll_w, len(aligned)):
            sf  = f_rets[i - roll_w: i]
            si  = i_rets[i - roll_w: i]
            mf  = sum(sf) / roll_w
            mi  = sum(si) / roll_w
            cov = sum((sf[j] - mf) * (si[j] - mi) for j in range(roll_w)) / (roll_w - 1)
            var = sum((si[j] - mi) ** 2 for j in range(roll_w)) / (roll_w - 1)
            series.append({"date": aligned[i], "beta": round(cov / var, 3) if var > 0 else None})

        result = {"series": series, "roll_window": roll_w}
        cache[key] = {"data": result, "expires_at": now + RISK_TTL}
        save_cache(cache)
        return jsonify(result)
    except Exception as e:
        return jsonify({"error": str(e)}), 500


@app.route("/api/risk/liquidity")
def api_risk_liquidity():
    portfolio   = load_portfolio()
    tickers     = [p["yahoo_ticker"] for p in portfolio["positions"]]
    prices      = get_cached_prices(tickers)
    funds       = get_cached_fundamentals(tickers)
    pdata       = build_portfolio_response(portfolio, prices, funds)
    fund_config = get_effective_fund_config()
    nav         = (pdata.get("total_value") or 0) + (fund_config.get("caixa") or 0) + (fund_config.get("proventos_a_receber") or 0)
    total_value = pdata.get("total_value") or 0

    liq_1d = liq_5d = liq_10d = 0.0
    rows_out = []
    missing_score_value = 0.0
    for r in pdata["rows"]:
        vl        = r.get("valor_liquido") or 0
        score     = r.get("liq_diaria_mm")
        if score is None:
            missing_score_value += vl
        days      = _liq_days_from_score(score)
        daily_pct = (1.0 / days) if days else None
        liq1d_v   = min(vl, vl * daily_pct)      if daily_pct else 0
        liq5d_v   = min(vl, vl * daily_pct * 5)  if daily_pct else 0
        liq10d_v  = min(vl, vl * daily_pct * 10) if daily_pct else 0
        liq_1d  += liq1d_v
        liq_5d  += liq5d_v
        liq_10d += liq10d_v
        rows_out.append({
            "ticker":       r["ticker"],
            "valor_liquido": round(vl, 2),
            "weight_pct":   r.get("pct_total"),
            "liq_score":    score,
            "days_to_liq":  round(days, 1) if days else None,
            "liq_1d_pct":   round(min(100, (daily_pct or 0) * 100), 1),
            "liq_5d_pct":   round(min(100, (daily_pct or 0) * 5 * 100), 1),
            "liq_10d_pct":  round(min(100, (daily_pct or 0) * 10 * 100), 1),
        })
    rows_out.sort(key=lambda x: x.get("days_to_liq") or 9999)

    # Detecta cold cache de fundamentals: sem average_volume, o score de
    # liquidez fica None para a posição e tudo é computado como 0%.
    missing_pct = (missing_score_value / total_value * 100) if total_value else 0
    warning = None
    if missing_pct >= 30:
        warning = (f"Scores de liquidez indisponíveis para {missing_pct:.0f}% do portfólio — "
                   "dados fundamentais ainda carregando. Recarregue a página em alguns segundos.")

    return jsonify({
        "nav_ref":              round(nav, 2),
        "total_equity_rs":      round(total_value, 2),
        "portfolio_liq_1d_pct":  round(liq_1d  / total_value * 100, 1) if total_value else 0,
        "portfolio_liq_5d_pct":  round(liq_5d  / total_value * 100, 1) if total_value else 0,
        "portfolio_liq_10d_pct": round(liq_10d / total_value * 100, 1) if total_value else 0,
        "portfolio_liq_1d_rs":   round(liq_1d,  2),
        "portfolio_liq_5d_rs":   round(liq_5d,  2),
        "portfolio_liq_10d_rs":  round(liq_10d, 2),
        "rows": rows_out,
        "warning": warning,
    })


@app.route("/api/risk/tracking-error")
def api_risk_tracking_error():
    import math
    window = int(request.args.get("window", 252))
    cache  = load_cache()
    now    = time.time()
    key    = f"risk_tracking_error_{window}"
    if cache.get(key) and now < cache[key].get("expires_at", 0):
        return jsonify(cache[key]["data"])

    history = load_quota_history()
    if len(history) < 22:
        return jsonify({"error": "Histórico insuficiente"}), 400

    start  = history[0]["data"]
    end_dt = datetime.strptime(history[-1]["data"], "%Y-%m-%d") + timedelta(days=5)
    try:
        ibov_h = yf.Ticker("^BVSP").history(
            start=start, end=end_dt.strftime("%Y-%m-%d"), timeout=15)
        if ibov_h.empty:
            return jsonify({"error": "Sem dados do IBOV"}), 400
        ibov_map   = {str(d.date()): float(v) for d, v in ibov_h["Close"].items()}
        ibov_dates = sorted(ibov_map)
        ibov_rets  = {ibov_dates[i]: ibov_map[ibov_dates[i]] / ibov_map[ibov_dates[i - 1]] - 1
                      for i in range(1, len(ibov_dates))}

        fund_map = {}
        prev = history[0]["cota_fechamento"]
        for e in history[1:]:
            fund_map[e["data"]] = e["cota_fechamento"] / prev - 1
            prev = e["cota_fechamento"]

        aligned = sorted(d for d in fund_map if d in ibov_rets)
        aligned  = aligned[-window:]
        if len(aligned) < 20:
            return jsonify({"error": "Dados insuficientes"}), 400

        f_rets = [fund_map[d] for d in aligned]
        i_rets = [ibov_rets[d] for d in aligned]
        n      = len(aligned)

        excess  = [f - i for f, i in zip(f_rets, i_rets)]
        mean_ex = sum(excess) / n
        var_ex  = sum((e - mean_ex) ** 2 for e in excess) / (n - 1)
        te      = math.sqrt(var_ex) * math.sqrt(252)
        ret_ativo = mean_ex * 252
        ir      = ret_ativo / te if te > 0 else None

        result = {
            "window":               window,
            "tracking_error":       round(te * 100, 2),
            "information_ratio":    round(ir, 2) if ir is not None else None,
            "retorno_ativo_anual":  round(ret_ativo * 100, 2),
            "n_dias":               n,
        }
        cache[key] = {"data": result, "expires_at": now + RISK_TTL}
        save_cache(cache)
        return jsonify(result)
    except Exception as e:
        return jsonify({"error": str(e)}), 500


@app.route("/api/risk/sortino-calmar")
def api_risk_sortino_calmar():
    import math
    cache = load_cache()
    now   = time.time()
    key   = "risk_sortino_calmar"
    if cache.get(key) and now < cache[key].get("expires_at", 0):
        return jsonify(cache[key]["data"])

    history = load_quota_history()
    if len(history) < 22:
        return jsonify({"error": "Histórico insuficiente"}), 400

    entries   = [(datetime.strptime(e["data"], "%Y-%m-%d"), e["cota_fechamento"]) for e in history]
    last_date = entries[-1][0]
    cdi_map   = load_cdi_map()

    def cdi_ann_for_window(start_dt, end_dt):
        rates = [v for d, v in cdi_map.items()
                 if start_dt < datetime.strptime(d, "%Y-%m-%d") <= end_dt]
        if not rates:
            return None
        compound = 1.0
        for r in rates:
            compound *= (1 + r / 100)
        return (compound ** (252 / len(rates)) - 1) * 100

    def get_window(months_back=None, ytd=False, no_mes=False):
        if no_mes:
            cutoff = datetime(last_date.year, last_date.month, 1)
        elif ytd:
            cutoff = datetime(last_date.year, 1, 1)
        elif months_back:
            cutoff = last_date - timedelta(days=int(months_back * 365.25 / 12))
        else:
            return entries
        before = [(d, c) for d, c in entries if d < cutoff]
        after  = [(d, c) for d, c in entries if d >= cutoff]
        return ([before[-1]] if before else []) + after

    def compute_sortino_calmar(sl):
        if len(sl) < 5:
            return {"sortino": None, "calmar": None, "downside_vol": None, "max_dd": None}
        cotas = [c for _, c in sl]
        daily = [(cotas[i] / cotas[i - 1] - 1) for i in range(1, len(cotas))]
        if len(daily) < 2:
            return {"sortino": None, "calmar": None, "downside_vol": None, "max_dd": None}
        n       = len(daily)
        ret_ann = (cotas[-1] / cotas[0]) ** (252 / n) - 1
        # Max drawdown
        peak   = cotas[0]
        max_dd = 0.0
        for c in cotas:
            if c > peak:
                peak = c
            dd = c / peak - 1
            if dd < max_dd:
                max_dd = dd
        # Downside deviation (MAR = 0)
        neg = [r for r in daily if r < 0]
        if neg:
            downside_dev = math.sqrt(sum(r ** 2 for r in neg) / n) * math.sqrt(252)
        else:
            downside_dev = 0
        dates = [d for d, _ in sl]
        cdi   = cdi_ann_for_window(dates[0], dates[-1])
        sortino = round((ret_ann * 100 - (cdi or 0)) / (downside_dev * 100), 2) if downside_dev > 0 else None
        calmar  = round(ret_ann / abs(max_dd), 2) if max_dd < 0 else None
        return {
            "sortino":      sortino,
            "calmar":       calmar,
            "downside_vol": round(downside_dev * 100, 2),
            "max_dd":       round(max_dd * 100, 2),
        }

    windows = [
        ("no_mes", get_window(no_mes=True)),
        ("no_ano", get_window(ytd=True)),
        ("3m",     get_window(months_back=3)),
        ("6m",     get_window(months_back=6)),
        ("12m",    get_window(months_back=12)),
        ("24m",    get_window(months_back=24)),
        ("36m",    get_window(months_back=36)),
        ("total",  entries),
    ]
    result = {"windows": {k: compute_sortino_calmar(sl) for k, sl in windows}}
    cache[key] = {"data": result, "expires_at": now + RISK_TTL}
    save_cache(cache)
    return jsonify(result)


@app.route("/api/risk/capture")
def api_risk_capture():
    window = request.args.get("window", "252")
    cache  = load_cache()
    now    = time.time()
    key    = f"risk_capture_{window}"
    if cache.get(key) and now < cache[key].get("expires_at", 0):
        return jsonify(cache[key]["data"])

    history = load_quota_history()
    if len(history) < 22:
        return jsonify({"error": "Histórico insuficiente"}), 400

    start  = history[0]["data"]
    end_dt = datetime.strptime(history[-1]["data"], "%Y-%m-%d") + timedelta(days=5)
    try:
        ibov_h = yf.Ticker("^BVSP").history(
            start=start, end=end_dt.strftime("%Y-%m-%d"), timeout=15)
        if ibov_h.empty:
            return jsonify({"error": "Sem dados do IBOV"}), 400
        ibov_map   = {str(d.date()): float(v) for d, v in ibov_h["Close"].items()}
        ibov_dates = sorted(ibov_map)
        ibov_rets  = {ibov_dates[i]: ibov_map[ibov_dates[i]] / ibov_map[ibov_dates[i - 1]] - 1
                      for i in range(1, len(ibov_dates))}

        fund_map = {}
        prev = history[0]["cota_fechamento"]
        for e in history[1:]:
            fund_map[e["data"]] = e["cota_fechamento"] / prev - 1
            prev = e["cota_fechamento"]

        aligned = sorted(d for d in fund_map if d in ibov_rets)
        if window != "total":
            aligned = aligned[-int(window):]
        if len(aligned) < 10:
            return jsonify({"error": "Dados insuficientes"}), 400

        f_rets = [fund_map[d] for d in aligned]
        i_rets = [ibov_rets[d] for d in aligned]

        up_idx   = [i for i in range(len(aligned)) if i_rets[i] > 0]
        down_idx = [i for i in range(len(aligned)) if i_rets[i] < 0]

        def _capture(indices):
            if not indices:
                return None
            fund_cum = 1.0
            ibov_cum = 1.0
            for idx in indices:
                fund_cum *= (1 + f_rets[idx])
                ibov_cum *= (1 + i_rets[idx])
            ibov_tot = ibov_cum - 1
            if abs(ibov_tot) < 1e-10:
                return None
            return round((fund_cum - 1) / ibov_tot * 100, 1)

        result = {
            "window":           window,
            "upside_capture":   _capture(up_idx),
            "downside_capture": _capture(down_idx),
            "n_dias_up":        len(up_idx),
            "n_dias_down":      len(down_idx),
            "n_total":          len(aligned),
        }
        cache[key] = {"data": result, "expires_at": now + RISK_TTL}
        save_cache(cache)
        return jsonify(result)
    except Exception as e:
        return jsonify({"error": str(e)}), 500


@app.route("/api/risk/concentration")
def api_risk_concentration():
    portfolio   = load_portfolio()
    tickers     = [p["yahoo_ticker"] for p in portfolio["positions"]]
    prices      = get_cached_prices(tickers)
    funds       = get_cached_fundamentals(tickers)
    pdata       = build_portfolio_response(portfolio, prices, funds)
    total_value = pdata.get("total_value") or 0
    if not total_value:
        return jsonify({"error": "Sem dados de portfólio"}), 400

    sector_map = {}
    for r in pdata["rows"]:
        vl     = r.get("valor_liquido") or 0
        sector = r.get("sector") or "Outros"
        if sector not in sector_map:
            sector_map[sector] = {"tickers": [], "valor": 0.0}
        sector_map[sector]["tickers"].append(r["ticker"])
        sector_map[sector]["valor"] += vl

    setores = []
    hhi     = 0.0
    for setor, data in sector_map.items():
        peso = data["valor"] / total_value
        hhi += peso ** 2
        setores.append({
            "setor":    setor,
            "peso_pct": round(peso * 100, 2),
            "valor_rs": round(data["valor"], 2),
            "tickers":  data["tickers"],
        })
    setores.sort(key=lambda x: x["peso_pct"], reverse=True)

    hhi_score = int(round(hhi * 10000))
    if hhi_score < 1000:
        hhi_label = "diversificado"
    elif hhi_score < 2500:
        hhi_label = "moderado"
    else:
        hhi_label = "concentrado"

    sorted_rows = sorted(pdata["rows"], key=lambda x: x.get("pct_total") or 0, reverse=True)
    top1 = sum(r.get("pct_total") or 0 for r in sorted_rows[:1])
    top3 = sum(r.get("pct_total") or 0 for r in sorted_rows[:3])
    top5 = sum(r.get("pct_total") or 0 for r in sorted_rows[:5])

    # Sinaliza quando o cache de fundamentals está frio e a maior parte do
    # portfólio caiu em "Outros" — HHI nesse caso não representa concentração
    # setorial real e usuário deve recarregar.
    peso_outros = next((s["peso_pct"] for s in setores if s["setor"] == "Outros"), 0)
    warning = None
    if peso_outros >= 30:
        warning = (f"Setores indisponíveis para {peso_outros:.0f}% do portfólio — "
                   "dados fundamentais ainda carregando. Recarregue a página em alguns segundos.")

    return jsonify({
        "hhi":        hhi_score,
        "hhi_label":  hhi_label,
        "setores":    setores,
        "top1_pct":   round(top1, 2),
        "top3_pct":   round(top3, 2),
        "top5_pct":   round(top5, 2),
        "n_posicoes": len(pdata["rows"]),
        "n_setores":  len(setores),
        "warning":    warning,
    })


@app.route("/api/risk/fx-exposure")
def api_risk_fx_exposure():
    portfolio   = load_portfolio()
    tickers     = [p["yahoo_ticker"] for p in portfolio["positions"]]
    prices      = get_cached_prices(tickers)
    funds       = get_cached_fundamentals(tickers)
    pdata       = build_portfolio_response(portfolio, prices, funds)
    fund_config = get_effective_fund_config()
    nav         = (pdata.get("total_value") or 0) + (fund_config.get("caixa") or 0) + (fund_config.get("proventos_a_receber") or 0)
    total_value = pdata.get("total_value") or 0
    if not total_value:
        return jsonify({"error": "Sem dados de portfólio"}), 400

    bdrs     = [r for r in pdata["rows"] if r.get("categoria") == "BDR"]
    total_fx = sum(r.get("valor_liquido") or 0 for r in bdrs)
    fx_pct   = total_fx / total_value

    bdr_rows = sorted([{
        "ticker":   r["ticker"],
        "peso_pct": round((r.get("valor_liquido") or 0) / total_value * 100, 2),
        "valor_rs": round(r.get("valor_liquido") or 0, 2),
        "sector":   r.get("sector") or "—",
    } for r in bdrs], key=lambda x: x["peso_pct"], reverse=True)

    sens = {}
    for shock in [5, 10, -5, -10]:
        k = f"usd_{'plus' if shock > 0 else 'minus'}{abs(shock)}"
        sens[k] = round(fx_pct * shock, 2)

    return jsonify({
        "total_fx_exposure_pct": round(fx_pct * 100, 2),
        "total_fx_exposure_rs":  round(total_fx, 2),
        "nav_ref":               round(nav, 2),
        "bdrs":                  bdr_rows,
        "sensibilidade_pct":     sens,
    })


@app.route("/api/risk/rolling-ratios")
def api_risk_rolling_ratios():
    import math
    roll_w = int(request.args.get("roll_window", 63))
    cache  = load_cache()
    now    = time.time()
    key    = f"risk_rolling_ratios_{roll_w}"
    if cache.get(key) and now < cache[key].get("expires_at", 0):
        return jsonify(cache[key]["data"])

    history = load_quota_history()
    if len(history) < roll_w + 5:
        return jsonify({"error": "Histórico insuficiente"}), 400

    cotas       = [e["cota_fechamento"] for e in history]
    dates       = [e["data"] for e in history]
    daily       = [(cotas[i] / cotas[i - 1] - 1) for i in range(1, len(cotas))]
    daily_dates = dates[1:]

    # CDI diário médio como proxy da taxa livre de risco
    cdi_map  = load_cdi_map()
    cdi_vals = list(cdi_map.values())
    cdi_daily = (sum(cdi_vals[-252:]) / min(252, len(cdi_vals))) / 100 if cdi_vals else 0.0

    series       = []
    sharpe_vals  = []
    sortino_vals = []

    for i in range(roll_w, len(daily)):
        sl  = daily[i - roll_w: i]
        n   = len(sl)
        mn  = sum(sl) / n
        var = sum((r - mn) ** 2 for r in sl) / (n - 1) if n > 1 else 0
        std = math.sqrt(var) if var > 0 else 0
        vol_ann = std * math.sqrt(252)
        ret_ann = (1 + mn) ** 252 - 1
        rf_ann  = (1 + cdi_daily) ** 252 - 1

        sharpe  = round((ret_ann - rf_ann) / vol_ann, 2) if vol_ann > 0 else None

        neg = [r for r in sl if r < 0]
        if neg:
            dd_ann = math.sqrt(sum(r ** 2 for r in neg) / n) * math.sqrt(252)
        else:
            dd_ann = 0
        sortino = round((ret_ann - rf_ann) / dd_ann, 2) if dd_ann > 0 else None

        series.append({
            "date":    daily_dates[i],
            "sharpe":  sharpe,
            "sortino": sortino,
        })
        if sharpe is not None:
            sharpe_vals.append(sharpe)
        if sortino is not None:
            sortino_vals.append(sortino)

    current = series[-1] if series else {}
    result  = {
        "roll_window":    roll_w,
        "series":         series,
        "current_sharpe": current.get("sharpe"),
        "current_sortino":current.get("sortino"),
        "avg_sharpe":     round(sum(sharpe_vals)  / len(sharpe_vals),  2) if sharpe_vals  else None,
        "avg_sortino":    round(sum(sortino_vals) / len(sortino_vals), 2) if sortino_vals else None,
    }
    cache[key] = {"data": result, "expires_at": now + RISK_TTL}
    save_cache(cache)
    return jsonify(result)


@app.route("/api/risk/return-distribution")
def api_risk_return_distribution():
    import math
    window = int(request.args.get("window", 252))
    cache  = load_cache()
    now    = time.time()
    key    = f"risk_return_dist_{window}"
    if cache.get(key) and now < cache[key].get("expires_at", 0):
        return jsonify(cache[key]["data"])

    history = load_quota_history()
    if len(history) < 22:
        return jsonify({"error": "Histórico insuficiente"}), 400

    cotas = [e["cota_fechamento"] for e in history]
    rets  = [(cotas[i] / cotas[i - 1] - 1) for i in range(1, len(cotas))]
    rets  = rets[-window:]
    n     = len(rets)

    mean_r   = sum(rets) / n
    variance = sum((r - mean_r) ** 2 for r in rets) / (n - 1) if n > 1 else 0
    std_r    = math.sqrt(variance) if variance > 0 else 0

    skew = (sum((r - mean_r) ** 3 for r in rets) / n) / (std_r ** 3) if std_r > 0 and n > 2 else 0
    kurt = (sum((r - mean_r) ** 4 for r in rets) / n) / (std_r ** 4) - 3 if std_r > 0 and n > 3 else 0

    min_r    = min(rets)
    max_r    = max(rets)
    n_bins   = 20
    bin_size = (max_r - min_r) / n_bins if max_r > min_r else 0.001
    bins     = [min_r + i * bin_size for i in range(n_bins + 1)]
    counts   = [0] * n_bins
    for r in rets:
        idx = min(int((r - min_r) / bin_size), n_bins - 1)
        counts[idx] += 1

    sorted_r = sorted(rets)
    def pct(p): return sorted_r[max(0, int(p * n / 100) - 1)]

    # IBOV comparison
    start_h    = history[-min(window + 5, len(history))]["data"]
    end_dt     = datetime.strptime(history[-1]["data"], "%Y-%m-%d") + timedelta(days=5)
    ibov_counts = None
    ibov_mean   = None
    ibov_std    = None
    try:
        ibov_h = yf.Ticker("^BVSP").history(
            start=start_h, end=end_dt.strftime("%Y-%m-%d"), timeout=15)
        if not ibov_h.empty:
            ic    = list(ibov_h["Close"])
            ir    = [(ic[i] / ic[i - 1] - 1) for i in range(1, len(ic))]
            ir    = ir[-window:]
            ni    = len(ir)
            mi    = sum(ir) / ni
            vi    = sum((r - mi) ** 2 for r in ir) / ni if ni > 1 else 0
            ibov_mean   = round(mi * 100, 3)
            ibov_std    = round(math.sqrt(vi) * 100, 3)
            ibov_counts = [0] * n_bins
            for r in ir:
                idx = min(int((r - min_r) / bin_size), n_bins - 1)
                if 0 <= idx < n_bins:
                    ibov_counts[idx] += 1
    except Exception:
        pass

    bin_centers = [round((bins[i] + bins[i + 1]) / 2 * 100, 3) for i in range(n_bins)]

    result = {
        "window":       window,
        "n_obs":        n,
        "bin_centers":  bin_centers,
        "counts":       counts,
        "ibov_counts":  ibov_counts,
        "mean_pct":     round(mean_r * 100, 3),
        "std_pct":      round(std_r * 100, 3),
        "skewness":     round(skew, 3),
        "kurtosis":     round(kurt, 3),
        "pct_positive": round(sum(1 for r in rets if r > 0) / n * 100, 1),
        "best_day":     round(max(rets) * 100, 3),
        "worst_day":    round(min(rets) * 100, 3),
        "p5":           round(pct(5) * 100, 3),
        "p95":          round(pct(95) * 100, 3),
        "ibov_mean_pct":ibov_mean,
        "ibov_std_pct": ibov_std,
    }
    cache[key] = {"data": result, "expires_at": now + RISK_TTL}
    save_cache(cache)
    return jsonify(result)


# ─── PRÉ-TRADE ────────────────────────────────────────────────────────────────

# Lock global para mutações da carteira (portfolio.json + fund_config.caixa).
# Protege contra concorrência entre /api/pretrade/execute, /api/portfolio/update,
# /api/portfolio/add e /api/portfolio/<ticker> DELETE.
_portfolio_write_lock = threading.Lock()


def _validate_and_normalize_ops(ops_input):
    """Normaliza tickers (.SA, upper-case) e valida preço/direção.
    Retorna (ops_normalizadas, None) em caso de sucesso, ou (None, msg_erro) em erro."""
    if not ops_input:
        return None, "Informe ao menos uma operação"
    for op in ops_input:
        t = str(op.get("ticker") or "").strip()
        if not t:
            return None, "ticker é obrigatório em todas as operações"
        if "." not in t:
            t += ".SA"
        op["ticker"]        = t.upper()
        op["preco"]         = float(op.get("preco") or 0)
        op["quantidade"]    = float(op.get("quantidade") or 0)
        op["direcao"]       = str(op.get("direcao") or "compra").lower()
        op["corretagem_rs"] = float(op.get("corretagem_rs") or 0)
        if op["preco"] <= 0:
            return None, f"Preço inválido para {op['ticker']}"
        if op["direcao"] not in ("compra", "venda", "zerar"):
            return None, f"Direção inválida para {op['ticker']}"
    return ops_input, None


def _apply_operations_to_portfolio(portfolio, fund_config, ops_input, fundamentals=None):
    """Aplica operações (compra/venda/zerar) a um clone de portfolio + fund_config.

    Retorna (portfolio_novo, fund_config_novo, ops_processadas, custo_basket, prices_sim).
    - Cria posição com defaults (categoria=Acao, demais None) para tickers novos.
    - Remove posições com quantidade <= 0 ao final.
    - prices_sim é construído apenas para uso da camada de simulação (build_portfolio_response).
      Se fundamentals=None, o campo sector em ops_processadas cai para "Outros".
    """
    import copy
    fundamentals = fundamentals or {}

    portfolio_sim   = copy.deepcopy(portfolio)
    fund_config_sim = copy.deepcopy(fund_config)
    prices_sim      = {}

    ops_processadas = []
    custo_basket    = 0.0

    for op in ops_input:
        ticker     = op["ticker"]
        quantidade = op["quantidade"]
        direcao    = op["direcao"]
        preco      = op["preco"]
        corretagem = op["corretagem_rs"]

        prices_sim[ticker] = {"price": preco, "change_pct": 0.0}

        pos_existente = next((p for p in portfolio["positions"] if p["yahoo_ticker"] == ticker), None)
        pos_sim       = next((p for p in portfolio_sim["positions"] if p["yahoo_ticker"] == ticker), None)
        valor_op      = preco * quantidade

        if direcao == "compra":
            if pos_sim:
                pos_sim["quantidade"] = (pos_sim.get("quantidade") or 0) + quantidade
            else:
                portfolio_sim["positions"].append({
                    "ticker": ticker.replace(".SA", ""), "yahoo_ticker": ticker,
                    "categoria": "Acao", "quantidade": quantidade,
                    "liq_diaria_mm": None, "lucro_mi_26": None, "preco_alvo": None,
                })
            fund_config_sim["caixa"] = (fund_config_sim.get("caixa") or 0) - valor_op - corretagem
            custo_op = valor_op + corretagem
        elif direcao == "venda":
            if pos_sim:
                pos_sim["quantidade"] = max(0, (pos_sim.get("quantidade") or 0) - quantidade)
            fund_config_sim["caixa"] = (fund_config_sim.get("caixa") or 0) + valor_op - corretagem
            custo_op = -(valor_op - corretagem)
        else:  # zerar
            qtd_atual  = (pos_sim or {}).get("quantidade") or 0
            portfolio_sim["positions"] = [p for p in portfolio_sim["positions"] if p["yahoo_ticker"] != ticker]
            fund_config_sim["caixa"] = (fund_config_sim.get("caixa") or 0) + preco * qtd_atual - corretagem
            quantidade = qtd_atual
            valor_op   = preco * quantidade
            custo_op   = -(valor_op - corretagem)

        custo_basket += custo_op
        sector_ativo  = (fundamentals.get(ticker) or {}).get("sector") or "Outros"
        ops_processadas.append({
            "ticker":         ticker.replace(".SA", ""),
            "yahoo_ticker":   ticker,
            "is_novo":        pos_existente is None and direcao == "compra",
            "sector":         sector_ativo,
            "direcao":        direcao,
            "quantidade":     quantidade,
            "preco":          preco,
            "valor_total_rs": round(valor_op, 2),
            "corretagem_rs":  corretagem,
            "custo_op_rs":    round(custo_op, 2),
        })

    # Remover posições zeradas
    portfolio_sim["positions"] = [p for p in portfolio_sim["positions"] if (p.get("quantidade") or 0) > 0]

    return portfolio_sim, fund_config_sim, ops_processadas, custo_basket, prices_sim


@app.route("/api/pretrade/simulate", methods=["POST"])
@require_admin
def api_pretrade_simulate():
    payload = request.json or {}

    # Aceita basket {operacoes: [...]} ou formato legado {ticker, quantidade, ...}
    if "operacoes" in payload:
        ops_input = payload["operacoes"]
    else:
        ops_input = [{
            "ticker":        payload.get("ticker", "").strip(),
            "quantidade":    payload.get("quantidade", 0),
            "direcao":       payload.get("direcao", "compra"),
            "preco":         payload.get("preco", 0),
            "corretagem_rs": payload.get("corretagem_rs", 0),
        }]

    ops_input, err = _validate_and_normalize_ops(ops_input)
    if err:
        return jsonify({"error": err}), 400

    portfolio    = load_portfolio()
    fund_config  = get_effective_fund_config()
    tickers_cart = [p["yahoo_ticker"] for p in portfolio["positions"]]

    tickers_all = list(tickers_cart)
    for op in ops_input:
        if op["ticker"] not in tickers_all:
            tickers_all.append(op["ticker"])

    prices       = get_cached_prices(tickers_all)
    fundamentals = get_cached_fundamentals(tickers_all)

    # ── ESTADO ANTES ──
    pdata_antes = build_portfolio_response(portfolio, prices, fundamentals)
    total_antes = pdata_antes.get("total_value") or 0
    quota_antes = calculate_quota(pdata_antes["rows"], fund_config, prices)
    conc_antes  = _calcular_concentracao_pretrade(pdata_antes["rows"], total_antes)

    # Grupo I antes (ações + BDRs — Res. CVM 175)
    _GRUPO1_CATS = {"Acao", "BDR", "Acao BDR"}
    nav_antes = quota_antes.get("nav_total") or total_antes
    valor_g1_antes = sum(
        (r.get("valor_liquido") or 0)
        for r in pdata_antes["rows"]
        if (r.get("categoria") or "Acao") in _GRUPO1_CATS
    )
    pct_g1_antes = (valor_g1_antes / nav_antes * 100) if nav_antes else 0

    # ── APLICAR (em clone) ──
    portfolio_sim, fund_config_sim, ops_processadas, custo_basket, prices_overrides = (
        _apply_operations_to_portfolio(portfolio, fund_config, ops_input, fundamentals)
    )
    # Mesclar preços de mercado com os preços simulados das operações
    import copy
    prices_sim = copy.deepcopy(prices)
    prices_sim.update(prices_overrides)

    # ── ESTADO DEPOIS ──
    pdata_depois = build_portfolio_response(portfolio_sim, prices_sim, fundamentals)
    total_depois = pdata_depois.get("total_value") or 0
    quota_depois = calculate_quota(pdata_depois["rows"], fund_config_sim, prices_sim)
    conc_depois  = _calcular_concentracao_pretrade(pdata_depois["rows"], total_depois) if total_depois else {"por_ativo": {}, "por_setor": {}, "hhi": 0}

    nav_depois = quota_depois.get("nav_total") or total_depois
    valor_g1_depois = sum(
        (r.get("valor_liquido") or 0)
        for r in pdata_depois["rows"]
        if (r.get("categoria") or "Acao") in _GRUPO1_CATS
    )
    pct_g1_depois = (valor_g1_depois / nav_depois * 100) if nav_depois else 0

    cota_antes_v  = quota_antes.get("cota_estimada") or quota_antes.get("quota_fechamento") or 0
    cota_depois_v = quota_depois.get("cota_estimada") or quota_depois.get("quota_fechamento") or 0
    imp_por_cota  = round(cota_depois_v - cota_antes_v, 8)
    caixa_depois  = fund_config_sim.get("caixa") or 0

    # ── COMPLIANCE — Resolução CVM 175 + limites internos (condicionais) ──
    lim_ativo  = float(fund_config.get("limite_concentracao_ativo_pct") or 20.0)
    lim_setor  = float(fund_config.get("limite_concentracao_setor_pct") or 40.0)
    ativo_on   = bool(fund_config.get("enable_concentracao_ativo", False))
    setor_on   = bool(fund_config.get("enable_concentracao_setor", False))

    def _status_max(valor, limite):
        if valor > limite:          return "violacao"
        if valor > limite * 0.85:   return "alerta"
        return "ok"

    def _status_min(valor, minimo):
        if valor < minimo * 0.85:   return "violacao"
        if valor < minimo:          return "alerta"
        return "ok"

    compliance = [
        # Regra regulatória obrigatória: mínimo 67% Grupo I (Res. CVM 175)
        {
            "regra": "Mín. 67% em Ações/BDRs — Grupo I (Res. CVM 175)",
            "limite_pct": 67.0,
            "valor_antes_pct": round(pct_g1_antes, 2),
            "valor_depois_pct": round(pct_g1_depois, 2),
            "status": _status_min(pct_g1_depois, 67.0),
            "tipo": "minimo",
        },
    ]

    # Limites internos por ativo (apenas se habilitado)
    if ativo_on:
        tickers_tocados = {op["yahoo_ticker"] for op in ops_processadas}
        for ytk in sorted(tickers_tocados):
            pct_d = conc_depois["por_ativo"].get(ytk, 0.0)
            pct_a = conc_antes["por_ativo"].get(ytk, 0.0)
            compliance.append({
                "regra": f"Conc. por Ativo — {ytk.replace('.SA','')} (interno)",
                "limite_pct": lim_ativo,
                "valor_antes_pct": round(pct_a, 2),
                "valor_depois_pct": round(pct_d, 2),
                "status": _status_max(pct_d, lim_ativo),
                "tipo": "maximo",
            })

    # Limites internos por setor (apenas se habilitado)
    if setor_on:
        setores_tocados = {op["sector"] for op in ops_processadas}
        for setor in sorted(setores_tocados):
            pct_d = conc_depois["por_setor"].get(setor, 0.0)
            pct_a = conc_antes["por_setor"].get(setor, 0.0)
            compliance.append({
                "regra": f"Conc. por Setor — {setor} (interno)",
                "limite_pct": lim_setor,
                "valor_antes_pct": round(pct_a, 2),
                "valor_depois_pct": round(pct_d, 2),
                "status": _status_max(pct_d, lim_setor),
                "tipo": "maximo",
            })

    # Caixa disponível
    if caixa_depois < 0:
        compliance.append({
            "regra": "Caixa Disponível",
            "limite_pct": 0,
            "valor_antes_pct": round(fund_config.get("caixa") or 0, 2),
            "valor_depois_pct": round(caixa_depois, 2),
            "status": "alerta",
            "tipo": "caixa",
        })

    return jsonify({
        "operacoes": ops_processadas,
        "basket": {
            "custo_total_rs": round(custo_basket, 2),
            "num_operacoes":  len(ops_processadas),
        },
        "antes": {
            "nav_total":       round(nav_antes, 2),
            "cota_estimada":   round(cota_antes_v, 8),
            "weighted_beta":   round(pdata_antes.get("weighted_beta") or 0, 4),
            "weighted_upside": round(pdata_antes.get("weighted_upside") or 0, 2),
            "hhi":   conc_antes["hhi"],
            "caixa": round(fund_config.get("caixa") or 0, 2),
            "pct_grupo1": round(pct_g1_antes, 2),
        },
        "depois": {
            "nav_total":       round(nav_depois, 2),
            "cota_estimada":   round(cota_depois_v, 8),
            "weighted_beta":   round(pdata_depois.get("weighted_beta") or 0, 4),
            "weighted_upside": round(pdata_depois.get("weighted_upside") or 0, 2),
            "hhi":   conc_depois["hhi"],
            "caixa": round(caixa_depois, 2),
            "pct_grupo1": round(pct_g1_depois, 2),
        },
        "impactos": {
            "variacao_cota_pct":  round((cota_depois_v / cota_antes_v - 1) * 100, 4) if cota_antes_v else 0,
            "variacao_nav_rs":    round(nav_depois - nav_antes, 2),
            "variacao_beta":      round((pdata_depois.get("weighted_beta") or 0) - (pdata_antes.get("weighted_beta") or 0), 4),
            "variacao_upside_pp": round((pdata_depois.get("weighted_upside") or 0) - (pdata_antes.get("weighted_upside") or 0), 2),
            "variacao_hhi":       conc_depois["hhi"] - conc_antes["hhi"],
            "impacto_por_cota_rs": imp_por_cota,
        },
        "compliance": compliance,
        "parametros_compliance": {
            "grupo1_minimo_pct":             67.0,
            "enable_concentracao_ativo":     ativo_on,
            "limite_concentracao_ativo_pct": lim_ativo,
            "enable_concentracao_setor":     setor_on,
            "limite_concentracao_setor_pct": lim_setor,
        },
        "rows_depois": [
            {
                "ticker":        r["ticker"],
                "yahoo_ticker":  r.get("yahoo_ticker", ""),
                "pct_total":     round(r.get("pct_total") or 0, 2),
                "valor_liquido": round(r.get("valor_liquido") or 0, 2),
            }
            for r in sorted(pdata_depois["rows"], key=lambda x: x.get("pct_total") or 0, reverse=True)
        ],
        "rows_antes": [
            {
                "ticker":        r["ticker"],
                "yahoo_ticker":  r.get("yahoo_ticker", ""),
                "pct_total":     round(r.get("pct_total") or 0, 2),
                "valor_liquido": round(r.get("valor_liquido") or 0, 2),
            }
            for r in sorted(pdata_antes["rows"], key=lambda x: x.get("pct_total") or 0, reverse=True)
        ],
    })


# ─── EVENTOS CORPORATIVOS ──────────────────────────────────────────────────────

EVENTS_TTL = 24 * 3600

def _evento_descricao(tipo, confirmado):
    mapa = {
        "RESULTADO": "Earnings Date" if confirmado else "Earnings Date (estimado)",
        "DIVIDENDO": "Proventos declarados" if confirmado else "Proventos pagos",
        "EX-DIV":   "Data ex-dividendo",
        "SPLIT":    "Desdobramento de ações",
    }
    return mapa.get(tipo, tipo)

def fetch_events(tickers):
    import pandas as pd
    from concurrent.futures import ThreadPoolExecutor, as_completed
    hoje = _brt_now().date()

    def _fetch_one(yticker):
        eventos = []
        try:
            t = yf.Ticker(yticker)

            # Earnings date via calendar
            try:
                cal = t.calendar
                if cal:
                    earnings = cal.get("Earnings Date") or cal.get("earnings_date")
                    if earnings is not None:
                        datas = earnings if isinstance(earnings, (list, tuple)) else [earnings]
                        for d in datas[:1]:  # só o primeiro (mais provável)
                            try:
                                dstr = str(d.date()) if hasattr(d, "date") else str(d)[:10]
                                eventos.append({"tipo": "RESULTADO", "data": dstr, "valor": None, "confirmado": False})
                            except Exception:
                                pass
            except Exception:
                pass

            # Dividendos históricos (últimos 180 dias)
            try:
                divs = t.dividends
                if divs is not None and len(divs) > 0:
                    corte = pd.Timestamp(hoje) - pd.Timedelta(days=180)
                    divs_rec = divs[divs.index >= corte]
                    for ts, valor in divs_rec.items():
                        try:
                            dstr = str(ts.date()) if hasattr(ts, "date") else str(ts)[:10]
                            eventos.append({"tipo": "DIVIDENDO", "data": dstr, "valor": _round(float(valor), 4), "confirmado": True})
                        except Exception:
                            pass
            except Exception:
                pass

            # Ex-dividend date via info
            try:
                info = t.info or {}
                ex_div_ts = info.get("exDividendDate")
                if ex_div_ts:
                    import datetime as _dt
                    if isinstance(ex_div_ts, (int, float)):
                        ex_date = _dt.date.fromtimestamp(ex_div_ts)
                    elif hasattr(ex_div_ts, "date"):
                        ex_date = ex_div_ts.date()
                    else:
                        ex_date = None
                    if ex_date:
                        last_div = info.get("lastDividendValue") or info.get("dividendRate")
                        eventos.append({"tipo": "EX-DIV", "data": str(ex_date), "valor": _round(float(last_div), 4) if last_div else None, "confirmado": True})
            except Exception:
                pass

        except Exception:
            pass

        return yticker, eventos

    result = {}
    with ThreadPoolExecutor(max_workers=5) as ex:
        futures = {ex.submit(_fetch_one, t): t for t in tickers}
        for f in as_completed(futures):
            try:
                yticker, evs = f.result()
                result[yticker] = evs
            except Exception:
                pass
    return result

def get_cached_events(tickers):
    cache = load_cache()
    now   = time.time()
    key   = "events_v1"
    if key in cache and now < cache[key].get("expires_at", 0):
        return cache[key]["data"]
    fresh = fetch_events(tickers)
    cache[key] = {"data": fresh, "expires_at": now + EVENTS_TTL}
    save_cache(cache)
    return fresh

@app.route("/api/events")
def api_events():
    portfolio   = load_portfolio()
    positions   = portfolio["positions"]
    tickers     = [p["yahoo_ticker"] for p in positions]
    ticker_map  = {p["yahoo_ticker"]: p["ticker"] for p in positions}

    dias_futuro      = int(request.args.get("dias_futuro", 90))
    incl_historico   = request.args.get("incluir_historico", "1") != "0"
    force            = request.args.get("force", "0") == "1"

    if force:
        cache = load_cache()
        cache.pop("events_v1", None)
        save_cache(cache)

    eventos_por_ticker = get_cached_events(tickers)
    hoje = _brt_now().date()

    todos_eventos = []
    for yahoo_ticker, eventos in eventos_por_ticker.items():
        for ev in eventos:
            try:
                from datetime import date as _date
                data_ev = datetime.strptime(ev["data"], "%Y-%m-%d").date()
            except Exception:
                continue
            dias = (data_ev - hoje).days
            eh_futuro = 0 <= dias <= dias_futuro
            eh_hist   = incl_historico and -180 <= dias < 0
            if not (eh_futuro or eh_hist):
                continue
            todos_eventos.append({
                "ticker":         ticker_map.get(yahoo_ticker, yahoo_ticker.replace(".SA", "")),
                "yahoo_ticker":   yahoo_ticker,
                "tipo":           ev["tipo"],
                "data":           ev["data"],
                "dias_ate_evento": dias,
                "valor":          ev.get("valor"),
                "confirmado":     ev.get("confirmado", False),
                "descricao":      _evento_descricao(ev["tipo"], ev.get("confirmado", False)),
            })

    todos_eventos.sort(key=lambda x: x["data"])

    por_ativo = {}
    for ev in todos_eventos:
        por_ativo.setdefault(ev["yahoo_ticker"], []).append(ev)

    return jsonify({
        "eventos":    todos_eventos,
        "por_ativo":  por_ativo,
        "gerado_em":  _brt_now().isoformat(),
    })


# ---------------------------------------------------------------------------
# PRÉ-TRADE HISTORY — persistência e PDF
# ---------------------------------------------------------------------------

def _generate_pretrade_pdf(record):
    """Gera PDF de auditoria de uma simulação de pré-trade. Retorna bytes."""
    from io import BytesIO
    from reportlab.lib import colors
    from reportlab.lib.pagesizes import A4
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
    from reportlab.lib.units import cm
    from reportlab.platypus import (
        SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle, HRFlowable
    )

    buf = BytesIO()
    doc = SimpleDocTemplate(
        buf, pagesize=A4,
        leftMargin=1.8*cm, rightMargin=1.8*cm,
        topMargin=1.5*cm, bottomMargin=1.5*cm,
    )

    # ── Cores ──
    C_BG      = colors.HexColor("#0d0d1a")
    C_HDR     = colors.HexColor("#1a1a2e")
    C_HDR_TXT = colors.white
    C_OK      = colors.HexColor("#00cc88")
    C_ALERTA  = colors.HexColor("#f5a623")
    C_VIOL    = colors.HexColor("#cc3333")
    C_COMPRA  = colors.HexColor("#f5a623")
    C_VENDA   = colors.HexColor("#00cc88")
    C_ZERAR   = colors.HexColor("#cc3333")
    C_BODY    = colors.HexColor("#cccccc")
    C_MUTED   = colors.HexColor("#888888")
    C_LINE    = colors.HexColor("#333333")

    def status_color(s):
        return C_OK if s == "ok" else C_ALERTA if s == "alerta" else C_VIOL

    def dir_color(d):
        return C_COMPRA if d == "compra" else C_VENDA if d == "venda" else C_ZERAR

    styles = getSampleStyleSheet()
    mono   = "Courier"

    st_title = ParagraphStyle("title", fontName="Helvetica-Bold", fontSize=14,
                               textColor=C_HDR_TXT, spaceAfter=4)
    st_sub   = ParagraphStyle("sub",   fontName="Helvetica",      fontSize=9,
                               textColor=C_MUTED,   spaceAfter=2)
    st_sec   = ParagraphStyle("sec",   fontName="Helvetica-Bold", fontSize=10,
                               textColor=C_HDR_TXT, spaceBefore=10, spaceAfter=4)
    st_foot  = ParagraphStyle("foot",  fontName="Helvetica",      fontSize=8,
                               textColor=C_MUTED,   alignment=1, spaceBefore=6)

    def tbl_style(header_rows=1, extra=None):
        base = [
            ("BACKGROUND",  (0, 0), (-1, header_rows - 1), C_HDR),
            ("TEXTCOLOR",   (0, 0), (-1, header_rows - 1), C_HDR_TXT),
            ("FONTNAME",    (0, 0), (-1, header_rows - 1), "Helvetica-Bold"),
            ("FONTSIZE",    (0, 0), (-1, -1), 8),
            ("FONTNAME",    (0, header_rows), (-1, -1), mono),
            ("TEXTCOLOR",   (0, header_rows), (-1, -1), C_BODY),
            ("BACKGROUND",  (0, header_rows), (-1, -1), C_BG),
            ("ROWBACKGROUNDS", (0, header_rows), (-1, -1), [C_BG, colors.HexColor("#0f0f22")]),
            ("GRID",        (0, 0), (-1, -1), 0.3, C_LINE),
            ("TOPPADDING",  (0, 0), (-1, -1), 4),
            ("BOTTOMPADDING", (0, 0), (-1, -1), 4),
            ("LEFTPADDING", (0, 0), (-1, -1), 6),
            ("RIGHTPADDING", (0, 0), (-1, -1), 6),
        ]
        if extra:
            base.extend(extra)
        return TableStyle(base)

    def fmt(v, dec=2):
        try:
            return f"{float(v):,.{dec}f}".replace(",", "X").replace(".", ",").replace("X", ".")
        except Exception:
            return str(v) if v is not None else "—"

    def fmtR(v, dec=2):
        return f"R$ {fmt(v, dec)}" if v is not None else "—"

    # ── Timestamp e label ──
    ts_raw = record.get("timestamp", "")
    try:
        dt = datetime.fromisoformat(ts_raw)
        data_str = dt.strftime("%d/%m/%Y")
        hora_str = dt.strftime("%H:%M:%S")
    except Exception:
        data_str = hora_str = ts_raw

    rec_id    = record.get("id", "")[:8]
    label_str = record.get("label", "") or ""
    antes     = record.get("antes", {})
    depois    = record.get("depois", {})
    impactos  = record.get("impactos", {})
    operacoes = record.get("operacoes", [])
    basket    = record.get("basket", {})
    compliance = record.get("compliance", [])
    rows_antes  = {r["ticker"]: r for r in record.get("rows_antes", [])}
    rows_depois = record.get("rows_depois", [])

    story = []

    # ── Cabeçalho ──
    story.append(Paragraph("HARBOUR IAT FIF AÇÕES RL — RELATÓRIO PRÉ-TRADE", st_title))
    story.append(Paragraph(
        f"Data: {data_str}  |  Hora: {hora_str}  |  ID: {rec_id}" +
        (f"  |  Ref: {label_str}" if label_str else ""),
        st_sub
    ))
    story.append(HRFlowable(width="100%", thickness=0.5, color=C_LINE, spaceAfter=6))

    # ── Basket de Operações ──
    # ── Seção 0: Parâmetros de Compliance Utilizados ──
    params = record.get("parametros_compliance", {})
    if params:
        story.append(Paragraph("0. PARÂMETROS DE COMPLIANCE UTILIZADOS", st_sec))
        param_rows = [["REGRA", "STATUS", "VALOR"]]
        param_rows.append([
            "Grupo I Mín. (CVM 175)",
            "SEMPRE ATIVO",
            f"mín {fmt(params.get('grupo1_minimo_pct', 67.0))}%",
        ])
        param_rows.append([
            "Conc. por Ativo (interno)",
            "ATIVO" if params.get("enable_concentracao_ativo") else "INATIVO",
            f"máx {fmt(params.get('limite_concentracao_ativo_pct', 20.0))}%",
        ])
        param_rows.append([
            "Conc. por Setor (interno)",
            "ATIVO" if params.get("enable_concentracao_setor") else "INATIVO",
            f"máx {fmt(params.get('limite_concentracao_setor_pct', 40.0))}%",
        ])
        param_extra = [
            ("TEXTCOLOR", (1, 1), (1, 1), C_OK),    # Grupo I sempre ativo → verde
        ]
        for i, row in enumerate(param_rows[1:], start=1):
            status = row[1]
            if status == "ATIVO":
                param_extra.append(("TEXTCOLOR", (1, i), (1, i), C_OK))
            elif status == "INATIVO":
                param_extra.append(("TEXTCOLOR", (1, i), (1, i), C_MUTED))
        param_tbl = Table(param_rows, colWidths=[6*cm, 3.5*cm, 3.5*cm])
        param_tbl.setStyle(tbl_style(extra=param_extra))
        story.append(param_tbl)
        story.append(Spacer(1, 6))

    story.append(Paragraph("1. BASKET DE OPERAÇÕES", st_sec))
    dir_label = {"compra": "COMPRA", "venda": "VENDA", "zerar": "ZERAR"}
    op_data = [["ATIVO", "DIREÇÃO", "QTDE", "PREÇO", "VALOR TOTAL", "CORRETAGEM"]]
    for op in operacoes:
        op_data.append([
            op.get("ticker", ""),
            dir_label.get(op.get("direcao", ""), op.get("direcao", "")),
            fmt(op.get("quantidade", 0), 0),
            fmtR(op.get("preco")),
            fmtR(op.get("valor_total_rs")),
            fmtR(op.get("corretagem_rs", 0)),
        ])
    # Linha total
    op_data.append([
        f"CUSTO LÍQUIDO TOTAL: {fmtR(basket.get('custo_total_rs'))}",
        "", "", "", "", "",
    ])

    # Cores por direção nas linhas de operações
    op_extra = []
    for i, op in enumerate(operacoes, start=1):
        c = dir_color(op.get("direcao", ""))
        op_extra.append(("TEXTCOLOR", (1, i), (1, i), c))
    op_extra.append(("SPAN",       (0, len(operacoes) + 1), (-1, len(operacoes) + 1)))
    op_extra.append(("TEXTCOLOR",  (0, len(operacoes) + 1), (-1, len(operacoes) + 1), C_ALERTA))
    op_extra.append(("FONTNAME",   (0, len(operacoes) + 1), (-1, len(operacoes) + 1), "Helvetica-Bold"))
    op_extra.append(("BACKGROUND", (0, len(operacoes) + 1), (-1, len(operacoes) + 1), C_HDR))

    op_tbl = Table(op_data, colWidths=[3.5*cm, 1.8*cm, 1.8*cm, 2.2*cm, 2.8*cm, 2.8*cm])
    op_tbl.setStyle(tbl_style(extra=op_extra))
    story.append(op_tbl)
    story.append(Spacer(1, 6))

    # ── Impacto no Fundo ──
    story.append(Paragraph("2. IMPACTO NO FUNDO", st_sec))
    metrics = [
        ("Cota Estimada",     fmt(antes.get("cota_estimada"), 8), fmt(depois.get("cota_estimada"), 8), f"{fmt(impactos.get('variacao_cota_pct'), 4)}%"),
        ("NAV Total",         fmtR(antes.get("nav_total")),        fmtR(depois.get("nav_total")),        fmtR(impactos.get("variacao_nav_rs"))),
        ("Caixa Resultante",  fmtR(antes.get("caixa")),            fmtR(depois.get("caixa")),            fmtR((depois.get("caixa") or 0) - (antes.get("caixa") or 0))),
        ("Grupo I (Ações/BDRs)", f"{fmt(antes.get('pct_grupo1'))}%", f"{fmt(depois.get('pct_grupo1'))}%", f"{fmt((depois.get('pct_grupo1') or 0) - (antes.get('pct_grupo1') or 0))} pp"),
        ("Beta Pond.",        fmt(antes.get("weighted_beta"), 4),  fmt(depois.get("weighted_beta"), 4),  fmt(impactos.get("variacao_beta"), 4)),
        ("Upside Pond.",      f"{fmt(antes.get('weighted_upside'))}%", f"{fmt(depois.get('weighted_upside'))}%", f"{fmt(impactos.get('variacao_upside_pp'))} pp"),
        ("HHI Concentração",  str(antes.get("hhi", "—")),          str(depois.get("hhi", "—")),          str(impactos.get("variacao_hhi", "—"))),
    ]
    imp_data = [["MÉTRICA", "ANTES", "DEPOIS", "Δ"]] + [list(m) for m in metrics]

    # Cor vermelha para caixa negativo
    imp_extra = []
    for i, m in enumerate(metrics, start=1):
        nome = m[0]
        if nome == "Caixa Resultante":
            try:
                val = float(str(depois.get("caixa") or 0))
                if val < 0:
                    imp_extra.append(("TEXTCOLOR", (2, i), (2, i), C_VIOL))
            except Exception:
                pass
        if nome == "Grupo I (Ações/BDRs)":
            try:
                val = float(str(depois.get("pct_grupo1") or 0))
                if val < 67:
                    imp_extra.append(("TEXTCOLOR", (2, i), (2, i), C_VIOL))
            except Exception:
                pass

    imp_tbl = Table(imp_data, colWidths=[4.2*cm, 3.8*cm, 3.8*cm, 2.9*cm])
    imp_tbl.setStyle(tbl_style(extra=imp_extra))
    story.append(imp_tbl)
    story.append(Spacer(1, 6))

    # ── Checklist de Compliance ──
    story.append(Paragraph("3. CHECKLIST DE COMPLIANCE — RESOLUÇÃO CVM 175", st_sec))
    comp_data = [["REGRA", "LIMITE", "ANTES", "DEPOIS", "STATUS"]]
    comp_extra = []
    for i, c in enumerate(compliance, start=1):
        s = c.get("status", "ok")
        sc = status_color(s)
        sl = "OK" if s == "ok" else "ALERTA" if s == "alerta" else "VIOLAÇÃO"
        tipo = c.get("tipo", "")
        lim_str = f"mín {fmt(c.get('limite_pct'))}%" if tipo == "minimo" else \
                  f"máx {fmt(c.get('limite_pct'))}%" if tipo != "caixa" else "—"
        if tipo == "caixa":
            antes_str  = fmtR(c.get("valor_antes_pct"))
            depois_str = fmtR(c.get("valor_depois_pct"))
        else:
            antes_str  = f"{fmt(c.get('valor_antes_pct'))}%"
            depois_str = f"{fmt(c.get('valor_depois_pct'))}%"
        comp_data.append([c.get("regra", ""), lim_str, antes_str, depois_str, sl])
        comp_extra.append(("TEXTCOLOR", (4, i), (4, i), sc))
        comp_extra.append(("FONTNAME",  (4, i), (4, i), "Helvetica-Bold"))
        if s != "ok":
            comp_extra.append(("BACKGROUND", (0, i), (-1, i), colors.HexColor("#1a0a0a")))

    comp_tbl = Table(comp_data, colWidths=[5.5*cm, 2.2*cm, 2.5*cm, 2.5*cm, 2*cm])
    comp_tbl.setStyle(tbl_style(extra=comp_extra))
    story.append(comp_tbl)
    story.append(Spacer(1, 6))

    # ── Carteira Antes vs Depois ──
    story.append(Paragraph("4. CARTEIRA — ANTES vs DEPOIS", st_sec))
    cart_data = [["ATIVO", "ANTES %", "DEPOIS %", "Δ pp"]]
    tickers_ops = {op.get("ticker") for op in operacoes}
    for r in rows_depois:
        tk  = r.get("ticker", "")
        ra  = rows_antes.get(tk, {"pct_total": 0})
        d_pp = (r.get("pct_total") or 0) - (ra.get("pct_total") or 0)
        cart_data.append([tk, f"{fmt(ra.get('pct_total'))}%", f"{fmt(r.get('pct_total'))}%", f"{'+' if d_pp > 0 else ''}{fmt(d_pp)}"])
    # Ativos zerados
    for tk, ra in rows_antes.items():
        if not any(r.get("ticker") == tk for r in rows_depois):
            cart_data.append([f"{tk} [ZERADO]", f"{fmt(ra.get('pct_total'))}%", "0,00%", f"-{fmt(ra.get('pct_total'))}"])

    cart_extra = []
    for i, row in enumerate(cart_data[1:], start=1):
        tk = row[0].split(" ")[0]
        if tk in tickers_ops:
            cart_extra.append(("TEXTCOLOR", (0, i), (0, i), C_ALERTA))
        if "[ZERADO]" in row[0]:
            cart_extra.append(("TEXTCOLOR", (0, i), (-1, i), C_VIOL))

    cart_tbl = Table(cart_data, colWidths=[4*cm, 3*cm, 3*cm, 3*cm])
    cart_tbl.setStyle(tbl_style(extra=cart_extra))
    story.append(cart_tbl)

    # ── Rodapé de auditoria ──
    story.append(Spacer(1, 10))
    story.append(HRFlowable(width="100%", thickness=0.5, color=C_LINE, spaceAfter=4))
    story.append(Paragraph(
        f"Este relatório foi gerado automaticamente para fins de auditoria e compliance conforme "
        f"Resolução CVM 175. ID: {record.get('id', '')} | Gerado em: {data_str} {hora_str}",
        st_foot
    ))

    doc.build(story)
    return buf.getvalue()


@app.route("/api/pretrade/history", methods=["GET"])
@require_admin
def api_pretrade_history_list():
    history = load_pretrade_history()
    history.sort(key=lambda x: x.get("timestamp", ""), reverse=True)
    return jsonify(history)


@app.route("/api/pretrade/history/save", methods=["POST"])
@require_admin
def api_pretrade_history_save():
    import uuid
    payload = request.json or {}
    required = ["antes", "depois", "operacoes", "compliance"]
    for k in required:
        if k not in payload:
            return jsonify({"error": f"Campo obrigatório ausente: {k}"}), 400

    record = {
        "id":        str(uuid.uuid4()),
        "timestamp": _brt_now().isoformat(timespec="seconds"),
        "label":     (payload.get("label") or "").strip()[:120],
        "operacoes": payload.get("operacoes", []),
        "basket":    payload.get("basket", {}),
        "antes":     payload.get("antes", {}),
        "depois":    payload.get("depois", {}),
        "impactos":  payload.get("impactos", {}),
        "compliance": payload.get("compliance", []),
        "rows_antes": payload.get("rows_antes", []),
        "rows_depois": payload.get("rows_depois", []),
        "parametros_compliance": payload.get("parametros_compliance", {}),
    }

    history = load_pretrade_history()
    history.append(record)
    save_pretrade_history(history)

    return jsonify({"id": record["id"], "timestamp": record["timestamp"]}), 201


@app.route("/api/pretrade/execute", methods=["POST"])
@require_admin
def api_pretrade_execute():
    """Aplica as operações da simulação na carteira real (portfolio.json + fund_config.caixa).

    Payload: {
      operacoes: [...]          # mesma estrutura do /simulate
      pretrade_history_id?: str # opcional — marca registro como executado
      compliance_override?: bool# opcional — apenas para auditoria
    }
    """
    payload = request.json or {}
    ops_input = payload.get("operacoes") or []
    pretrade_history_id = payload.get("pretrade_history_id")
    compliance_override = bool(payload.get("compliance_override", False))

    ops_input, err = _validate_and_normalize_ops(ops_input)
    if err:
        return jsonify({"error": err}), 400

    executed_at = _brt_now().isoformat(timespec="seconds")

    with _portfolio_write_lock:
        portfolio    = load_portfolio()
        fund_config  = load_fund_config()
        tickers_all  = list({p["yahoo_ticker"] for p in portfolio["positions"]} | {op["ticker"] for op in ops_input})
        fundamentals = get_cached_fundamentals(tickers_all)

        # 1. Aplicar in-memory
        portfolio_novo, fund_config_novo, ops_processadas, custo_basket, _ = (
            _apply_operations_to_portfolio(portfolio, fund_config, ops_input, fundamentals)
        )

        # 2. Marcar registro no histórico (se id veio) — ANTES de salvar a carteira,
        # para que em caso de falha do save_portfolio possamos reverter o executed_at.
        history = None
        record  = None
        if pretrade_history_id:
            history = load_pretrade_history()
            record  = next((r for r in history if r.get("id") == pretrade_history_id), None)
            if record is None:
                return jsonify({"error": "Registro de histórico não encontrado"}), 404
            if record.get("executed_at"):
                return jsonify({
                    "error":       "Operação já executada",
                    "executed_at": record["executed_at"],
                }), 409
            record["executed_at"]          = executed_at
            record["compliance_override"]  = compliance_override
            save_pretrade_history(history)

        # 3. Persistir carteira e fund_config. Se falhar, reverter executed_at.
        try:
            save_portfolio(portfolio_novo)
            save_fund_config(fund_config_novo)
        except Exception as e:
            if record is not None and history is not None:
                record.pop("executed_at", None)
                record.pop("compliance_override", None)
                save_pretrade_history(history)
            return jsonify({"error": f"Falha ao salvar carteira: {e}"}), 500

        invalidate_price_cache()
        invalidate_history_cache()

    return jsonify({
        "ok":                  True,
        "executed_at":         executed_at,
        "operacoes_aplicadas": len(ops_processadas),
        "custo_basket_rs":     round(custo_basket, 2),
        "caixa_apos":          round(fund_config_novo.get("caixa") or 0, 2),
    }), 200


@app.route("/api/pretrade/history/<record_id>", methods=["DELETE"])
@require_admin
def api_pretrade_history_delete(record_id):
    history = load_pretrade_history()
    new_history = [r for r in history if r.get("id") != record_id]
    if len(new_history) == len(history):
        return jsonify({"error": "Registro não encontrado"}), 404
    save_pretrade_history(new_history)
    return jsonify({"ok": True})


@app.route("/api/pretrade/history/<record_id>/pdf", methods=["GET"])
@require_admin
def api_pretrade_history_pdf(record_id):
    history = load_pretrade_history()
    record  = next((r for r in history if r.get("id") == record_id), None)
    if not record:
        return jsonify({"error": "Registro não encontrado"}), 404

    try:
        pdf_bytes = _generate_pretrade_pdf(record)
    except Exception as e:
        return jsonify({"error": f"Erro ao gerar PDF: {str(e)}"}), 500

    ts = record.get("timestamp", "")[:10].replace("-", "")
    filename = f"pretrade_{ts}_{record_id[:8]}.pdf"

    from flask import make_response
    resp = make_response(pdf_bytes)
    resp.headers["Content-Type"]        = "application/pdf"
    resp.headers["Content-Disposition"] = f'attachment; filename="{filename}"'
    return resp


# ---------------------------------------------------------------------------
# Portfolio History
# ---------------------------------------------------------------------------

@app.route("/api/portfolio-history", methods=["GET"])
@require_admin
def api_portfolio_history_list():
    history = load_portfolio_history()
    result  = [{k: v for k, v in r.items() if k != "rows"} for r in history]
    result.sort(key=lambda x: x.get("timestamp", ""), reverse=True)
    return jsonify(result)

@app.route("/api/portfolio-history/save", methods=["POST"])
@require_admin
def api_portfolio_history_save():
    portfolio   = load_portfolio()
    tickers     = [p["yahoo_ticker"] for p in portfolio["positions"]]
    prices      = fetch_prices(tickers)
    funds       = get_cached_fundamentals(tickers)
    data        = build_portfolio_response(portfolio, prices, funds)
    fund_config = get_effective_fund_config()
    quota       = calculate_quota(data["rows"], fund_config, prices)
    if not quota.get("cota_estimada"):
        return jsonify({"error": "Cota estimada indisponível — preços não carregados"}), 500
    snap = _build_portfolio_snapshot(data, quota, source="manual")
    ph   = load_portfolio_history()
    ph.append(snap)
    ph.sort(key=lambda x: x["timestamp"])
    save_portfolio_history(ph)
    return jsonify({"id": snap["id"], "timestamp": snap["timestamp"]}), 201

@app.route("/api/portfolio-history/<record_id>", methods=["GET"])
@require_admin
def api_portfolio_history_detail(record_id):
    history = load_portfolio_history()
    record  = next((r for r in history if r.get("id") == record_id), None)
    if not record:
        return jsonify({"error": "Registro não encontrado"}), 404
    return jsonify(record)

@app.route("/api/portfolio-history/<record_id>", methods=["DELETE"])
@require_admin
def api_portfolio_history_delete(record_id):
    history     = load_portfolio_history()
    new_history = [r for r in history if r.get("id") != record_id]
    if len(new_history) == len(history):
        return jsonify({"error": "Registro não encontrado"}), 404
    save_portfolio_history(new_history)
    return jsonify({"ok": True})

@app.route("/api/portfolio-history/<record_id>/pdf", methods=["GET"])
@require_admin
def api_portfolio_history_pdf(record_id):
    history = load_portfolio_history()
    record  = next((r for r in history if r.get("id") == record_id), None)
    if not record:
        return jsonify({"error": "Registro não encontrado"}), 404
    try:
        pdf_bytes = _generate_portfolio_snapshot_pdf(record)
    except Exception as e:
        return jsonify({"error": f"Erro ao gerar PDF: {str(e)}"}), 500
    ts = record.get("date", record.get("timestamp", ""))[:10].replace("-", "")
    filename = f"carteira_{ts}_{record_id[:8]}.pdf"
    from flask import make_response
    resp = make_response(pdf_bytes)
    resp.headers["Content-Type"]        = "application/pdf"
    resp.headers["Content-Disposition"] = f'attachment; filename="{filename}"'
    return resp


def _generate_portfolio_snapshot_pdf(record):
    """Gera PDF de auditoria de um snapshot da carteira. Retorna bytes."""
    from io import BytesIO
    from reportlab.lib import colors
    from reportlab.lib.pagesizes import A4, landscape
    from reportlab.lib.styles import ParagraphStyle
    from reportlab.lib.units import cm
    from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle, HRFlowable

    buf = BytesIO()
    PAGE = landscape(A4)
    doc = SimpleDocTemplate(
        buf, pagesize=PAGE,
        leftMargin=1.4*cm, rightMargin=1.4*cm,
        topMargin=1.2*cm, bottomMargin=1.2*cm,
    )

    C_BG      = colors.HexColor("#0d0d1a")
    C_ALT     = colors.HexColor("#0f0f22")
    C_HDR     = colors.HexColor("#1a1a2e")
    C_HDR_TXT = colors.white
    C_OK      = colors.HexColor("#00cc88")
    C_NEG     = colors.HexColor("#cc3333")
    C_WARN    = colors.HexColor("#f5a623")
    C_BODY    = colors.HexColor("#cccccc")
    C_MUTED   = colors.HexColor("#888888")
    C_LINE    = colors.HexColor("#333333")
    C_MANUAL  = colors.HexColor("#00aacc")

    mono = "Courier"

    st_title = ParagraphStyle("title", fontName="Helvetica-Bold", fontSize=13,
                               textColor=C_HDR_TXT, spaceAfter=3)
    st_sub   = ParagraphStyle("sub",   fontName="Helvetica",      fontSize=8,
                               textColor=C_MUTED,   spaceAfter=2)
    st_sec   = ParagraphStyle("sec",   fontName="Helvetica-Bold", fontSize=9,
                               textColor=C_HDR_TXT, spaceBefore=8, spaceAfter=3)
    st_foot  = ParagraphStyle("foot",  fontName="Helvetica",      fontSize=7,
                               textColor=C_MUTED,   alignment=1, spaceBefore=4)

    def tbl_style(extra=None):
        base = [
            ("BACKGROUND",    (0, 0), (-1, 0), C_HDR),
            ("TEXTCOLOR",     (0, 0), (-1, 0), C_HDR_TXT),
            ("FONTNAME",      (0, 0), (-1, 0), "Helvetica-Bold"),
            ("FONTSIZE",      (0, 0), (-1, -1), 7),
            ("FONTNAME",      (0, 1), (-1, -1), mono),
            ("TEXTCOLOR",     (0, 1), (-1, -1), C_BODY),
            ("ROWBACKGROUNDS",(0, 1), (-1, -1), [C_BG, C_ALT]),
            ("GRID",          (0, 0), (-1, -1), 0.25, C_LINE),
            ("TOPPADDING",    (0, 0), (-1, -1), 3),
            ("BOTTOMPADDING", (0, 0), (-1, -1), 3),
            ("LEFTPADDING",   (0, 0), (-1, -1), 4),
            ("RIGHTPADDING",  (0, 0), (-1, -1), 4),
            ("ALIGN",         (0, 0), (-1, -1), "RIGHT"),
            ("ALIGN",         (0, 0), (0, -1),  "LEFT"),
        ]
        if extra:
            base.extend(extra)
        return TableStyle(base)

    def fmt(v, dec=2):
        try:
            return f"{float(v):,.{dec}f}".replace(",", "X").replace(".", ",").replace("X", ".")
        except Exception:
            return str(v) if v is not None else "—"

    def fmtpct(v):
        return fmt(v) + "%" if v is not None else "—"

    def fmtmm(v):
        return "R$" + fmt(v/1e6, 1) + "MM" if v is not None else "—"

    # ── Metadados ──
    ts_raw = record.get("timestamp", "")
    try:
        dt       = datetime.fromisoformat(ts_raw)
        data_str = dt.strftime("%d/%m/%Y")
        hora_str = dt.strftime("%H:%M:%S")
    except Exception:
        data_str = hora_str = ts_raw

    rec_id  = record.get("id", "")[:8]
    source  = record.get("source", "auto")
    s       = record.get("summary", {})
    rows    = record.get("rows", [])

    story = []

    # ── Cabeçalho ──
    src_label = "MANUAL" if source == "manual" else "AUTO (FECHAMENTO)"
    story.append(Paragraph("HARBOUR IAT FIF AÇÕES RL — SNAPSHOT DA CARTEIRA", st_title))
    story.append(Paragraph(
        f"Data: {data_str}  |  Hora: {hora_str}  |  ID: {rec_id}  |  Origem: {src_label}",
        st_sub
    ))
    story.append(HRFlowable(width="100%", thickness=0.5, color=C_LINE, spaceAfter=5))

    # ── Bloco de resumo (uma linha de métricas) ──
    nav_val  = s.get("total_value")
    cota_val = s.get("cota_estimada")
    var_val  = s.get("variacao_pct")
    n_pos    = s.get("num_positions", len(rows))

    sum_data = [[
        "NAV TOTAL", "COTA ESTIMADA", "VAR. DIA", "POSIÇÕES",
        "UPSIDE POND.", "BETA POND.", "P/L FWD POND.", "ROE POND.",
    ], [
        fmtmm(nav_val),
        fmt(cota_val, 8) if cota_val else "—",
        fmtpct(var_val),
        str(n_pos),
        fmtpct(s.get("w_upside_pct")),
        fmt(s.get("w_beta"), 2) if s.get("w_beta") is not None else "—",
        fmt(s.get("w_forward_pe"), 1) if s.get("w_forward_pe") is not None else "—",
        fmtpct(s.get("w_return_on_equity")),
    ]]

    sum_extra = [("ALIGN", (0, 1), (-1, 1), "CENTER")]
    if var_val is not None:
        c = C_OK if float(var_val) >= 0 else C_NEG
        sum_extra.append(("TEXTCOLOR", (2, 1), (2, 1), c))

    W = PAGE[0] - 2.8*cm
    sum_tbl = Table(sum_data, colWidths=[W/8]*8)
    sum_tbl.setStyle(tbl_style(extra=sum_extra))
    story.append(sum_tbl)
    story.append(Spacer(1, 6))

    # ── Tabela 1: Composição & Preços ──
    story.append(Paragraph("1. COMPOSIÇÃO E PREÇOS", st_sec))
    t1_hdr = ["ATIVO", "CATEG.", "SETOR", "% TOTAL", "VALOR LÍQ.", "PREÇO", "VAR. DIA", "QTDE", "P. ALVO", "UPSIDE"]
    t1_rows = [t1_hdr]
    for r in sorted(rows, key=lambda x: (x.get("pct_total") or 0), reverse=True):
        vd = r.get("var_dia_pct")
        up = r.get("upside_pct")
        t1_rows.append([
            r.get("ticker", ""),
            r.get("categoria", ""),
            (r.get("sector") or "")[:18],
            fmtpct(r.get("pct_total")),
            fmtmm(r.get("valor_liquido")),
            fmt(r.get("preco"), 2) if r.get("preco") is not None else "—",
            fmtpct(vd),
            fmt(r.get("quantidade"), 0) if r.get("quantidade") is not None else "—",
            fmt(r.get("preco_alvo"), 2) if r.get("preco_alvo") is not None else "—",
            fmtpct(up),
        ])

    # Linha de totais
    total_nav = s.get("total_value")
    t1_rows.append([
        "TOTAL", "", "", "100,00%",
        fmtmm(total_nav), "", "", "", "", fmtpct(s.get("w_upside_pct")),
    ])

    t1_extra = []
    for i, r in enumerate(rows, start=1):
        vd = r.get("var_dia_pct")
        up = r.get("upside_pct")
        if vd is not None:
            t1_extra.append(("TEXTCOLOR", (6, i), (6, i), C_OK if float(vd) >= 0 else C_NEG))
        if up is not None:
            t1_extra.append(("TEXTCOLOR", (9, i), (9, i), C_OK if float(up) >= 0 else C_NEG))
    # Total row
    ti = len(rows) + 1
    t1_extra += [
        ("BACKGROUND", (0, ti), (-1, ti), C_HDR),
        ("TEXTCOLOR",  (0, ti), (-1, ti), C_WARN),
        ("FONTNAME",   (0, ti), (-1, ti), "Helvetica-Bold"),
    ]

    cw1 = [2.0*cm, 1.6*cm, 3.2*cm, 1.5*cm, 2.2*cm, 1.6*cm, 1.6*cm, 2.0*cm, 1.6*cm, 1.6*cm]
    t1_tbl = Table(t1_rows, colWidths=cw1)
    t1_tbl.setStyle(tbl_style(extra=t1_extra))
    story.append(t1_tbl)
    story.append(Spacer(1, 6))

    # ── Tabela 2: Valuation & Fundamentos ──
    story.append(Paragraph("2. VALUATION E FUNDAMENTOS", st_sec))
    t2_hdr = ["ATIVO", "P/L TRAIL.", "P/L FWD.", "PEG", "EV/EBITDA", "ROE %", "BETA", "P/VPA", "DIV. YIELD", "MKT CAP", "LUCRO MI 26"]
    t2_rows = [t2_hdr]
    for r in sorted(rows, key=lambda x: (x.get("pct_total") or 0), reverse=True):
        t2_rows.append([
            r.get("ticker", ""),
            fmt(r.get("trailing_pe"), 1)    if r.get("trailing_pe")    is not None else "—",
            fmt(r.get("forward_pe"), 1)     if r.get("forward_pe")     is not None else "—",
            fmt(r.get("peg_ratio"), 2)      if r.get("peg_ratio")      is not None else "—",
            fmt(r.get("enterprise_to_ebitda"), 1) if r.get("enterprise_to_ebitda") is not None else "—",
            fmtpct(r.get("return_on_equity")),
            fmt(r.get("beta"), 2)           if r.get("beta")           is not None else "—",
            fmt(r.get("price_to_book"), 2)  if r.get("price_to_book")  is not None else "—",
            fmtpct(r.get("dividend_yield")),
            fmt(r.get("market_cap_bi"), 1) + "bi" if r.get("market_cap_bi") is not None else "—",
            fmt(r.get("lucro_mi_26"), 0)    if r.get("lucro_mi_26")    is not None else "—",
        ])

    # Linha de médias ponderadas
    t2_rows.append([
        "POND.",
        fmt(s.get("w_trailing_pe"), 1)          if s.get("w_trailing_pe")          is not None else "—",
        fmt(s.get("w_forward_pe"), 1)            if s.get("w_forward_pe")           is not None else "—",
        fmt(s.get("w_peg_ratio"), 2)             if s.get("w_peg_ratio")            is not None else "—",
        fmt(s.get("w_enterprise_to_ebitda"), 1)  if s.get("w_enterprise_to_ebitda") is not None else "—",
        fmtpct(s.get("w_return_on_equity")),
        fmt(s.get("w_beta"), 2)                  if s.get("w_beta")                 is not None else "—",
        fmt(s.get("w_price_to_book"), 2)         if s.get("w_price_to_book")        is not None else "—",
        fmtpct(s.get("w_dividend_yield")),
        "—",
        fmt(s.get("w_lucro_mi_26"), 0)           if s.get("w_lucro_mi_26")          is not None else "—",
    ])

    ti2 = len(rows) + 1
    t2_extra = [
        ("BACKGROUND", (0, ti2), (-1, ti2), C_HDR),
        ("TEXTCOLOR",  (0, ti2), (-1, ti2), C_WARN),
        ("FONTNAME",   (0, ti2), (-1, ti2), "Helvetica-Bold"),
    ]

    cw2 = [2.0*cm, 1.8*cm, 1.8*cm, 1.5*cm, 2.2*cm, 1.8*cm, 1.5*cm, 1.5*cm, 2.0*cm, 2.0*cm, 2.2*cm]
    t2_tbl = Table(t2_rows, colWidths=cw2)
    t2_tbl.setStyle(tbl_style(extra=t2_extra))
    story.append(t2_tbl)

    # ── Rodapé ──
    story.append(Spacer(1, 8))
    story.append(HRFlowable(width="100%", thickness=0.5, color=C_LINE, spaceAfter=3))
    story.append(Paragraph(
        f"Relatório gerado automaticamente para fins de auditoria interna. "
        f"ID: {record.get('id', '')}  |  Gerado em: {data_str} {hora_str}  |  Origem: {src_label}",
        st_foot
    ))

    doc.build(story)
    return buf.getvalue()


# ---------------------------------------------------------------------------
# Portfolio History — Análise avançada (aba 213: HISTÓRICO DA CARTEIRA)
# ---------------------------------------------------------------------------

def _ph_extract_metrics(snap):
    """Extrai métricas-chave do snapshot para time series e timeline."""
    summary = snap.get("summary") or {}
    rows    = snap.get("rows") or []
    nav     = summary.get("total_value") or 0

    # HHI por ativo (concentração)
    if nav > 0:
        hhi_ativo = round(sum(((r.get("valor_liquido") or 0) / nav) ** 2 for r in rows) * 10000)
    else:
        hhi_ativo = 0

    # Pesos por setor
    setores = {}
    for r in rows:
        sec = r.get("sector") or "Outros"
        setores[sec] = setores.get(sec, 0.0) + (r.get("valor_liquido") or 0)
    setor_pcts = {s: round(v / nav * 100, 2) for s, v in setores.items()} if nav > 0 else {}

    # HHI por setor
    if nav > 0:
        hhi_setor = round(sum((v / nav) ** 2 for v in setores.values()) * 10000)
    else:
        hhi_setor = 0

    # Grupo I (Ações + BDRs) — Res. CVM 175
    _GRUPO1 = {"Acao", "BDR", "Acao BDR"}
    valor_g1 = sum((r.get("valor_liquido") or 0) for r in rows
                   if (r.get("categoria") or "Acao") in _GRUPO1)
    pct_g1 = round(valor_g1 / nav * 100, 2) if nav > 0 else 0

    return {
        "nav":            nav,
        "num_positions":  summary.get("num_positions") or len(rows),
        "cota_estimada":  summary.get("cota_estimada"),
        "weighted_beta":  summary.get("w_beta"),
        "weighted_upside": summary.get("w_upside_pct"),
        "hhi_ativo":      hhi_ativo,
        "hhi_setor":      hhi_setor,
        "pct_grupo1":     pct_g1,
        "setor_pcts":     setor_pcts,
    }


@app.route("/api/portfolio-history/timeline", methods=["GET"])
@require_admin
def api_portfolio_history_timeline():
    """Timeline com summary expandido por snapshot, ordenado por data desc."""
    history = load_portfolio_history()
    items = []
    for snap in history:
        m = _ph_extract_metrics(snap)
        items.append({
            "id":            snap.get("id"),
            "date":          snap.get("date"),
            "timestamp":     snap.get("timestamp"),
            "source":        snap.get("source") or "auto",
            "nav":           m["nav"],
            "num_positions": m["num_positions"],
            "cota_estimada": m["cota_estimada"],
            "variacao_pct":  (snap.get("summary") or {}).get("variacao_pct"),
            "pct_grupo1":    m["pct_grupo1"],
            "hhi_ativo":     m["hhi_ativo"],
        })
    items.sort(key=lambda x: x.get("timestamp") or "", reverse=True)
    return jsonify({
        "snapshots":  items,
        "first_date": items[-1]["date"] if items else None,
        "last_date":  items[0]["date"]  if items else None,
        "total":      len(items),
    })


def _ph_find_snap_by_date(history, date_str):
    """Localiza o snapshot mais recente em ou antes da data dada (formato YYYY-MM-DD).
    Se não houver nenhum, retorna None."""
    candidates = [s for s in history if (s.get("date") or "") <= date_str]
    if not candidates:
        return None
    candidates.sort(key=lambda s: s.get("timestamp") or "")
    return candidates[-1]


def _diff_snapshots(snap_from, snap_to):
    """Compara duas snapshots e retorna estrutura de diff."""
    rows_from = {r["ticker"]: r for r in (snap_from.get("rows") or [])}
    rows_to   = {r["ticker"]: r for r in (snap_to.get("rows") or [])}
    all_tickers = sorted(set(rows_from) | set(rows_to))

    nav_from = (snap_from.get("summary") or {}).get("total_value") or 0
    nav_to   = (snap_to.get("summary") or {}).get("total_value")   or 0

    posicoes = []
    for tk in all_tickers:
        rf = rows_from.get(tk) or {}
        rt = rows_to.get(tk)   or {}
        qf = rf.get("quantidade") or 0
        qt = rt.get("quantidade") or 0
        vf = rf.get("valor_liquido") or 0
        vt = rt.get("valor_liquido") or 0
        pf = rf.get("pct_total") or 0
        pt = rt.get("pct_total") or 0
        status = "manteve"
        if qf == 0 and qt > 0:   status = "novo"
        elif qf > 0 and qt == 0: status = "removido"
        elif qt > qf:            status = "aumentou"
        elif qt < qf:            status = "reduziu"
        posicoes.append({
            "ticker":          tk,
            "categoria":       (rt.get("categoria") or rf.get("categoria") or "Acao"),
            "sector":          (rt.get("sector") or rf.get("sector") or "Outros"),
            "qtde_from":       qf, "qtde_to": qt, "delta_qtde": qt - qf,
            "valor_from":      round(vf, 2), "valor_to": round(vt, 2),
            "delta_valor":     round(vt - vf, 2),
            "pct_from":        round(pf, 2), "pct_to": round(pt, 2),
            "delta_pct_pp":    round(pt - pf, 2),
            "preco_from":      rf.get("preco"),
            "preco_to":        rt.get("preco"),
            "status":          status,
        })
    posicoes.sort(key=lambda x: abs(x["delta_valor"] or 0), reverse=True)

    m_from = _ph_extract_metrics(snap_from)
    m_to   = _ph_extract_metrics(snap_to)

    # Diff por setor (deltas em pp)
    setores_diff = []
    setores_all = sorted(set(m_from["setor_pcts"]) | set(m_to["setor_pcts"]))
    for sec in setores_all:
        pf = m_from["setor_pcts"].get(sec, 0)
        pt = m_to["setor_pcts"].get(sec, 0)
        setores_diff.append({
            "sector":     sec,
            "pct_from":   pf,
            "pct_to":     pt,
            "delta_pp":   round(pt - pf, 2),
        })
    setores_diff.sort(key=lambda x: abs(x["delta_pp"]), reverse=True)

    return {
        "snap_from": {
            "id":        snap_from.get("id"),
            "date":      snap_from.get("date"),
            "timestamp": snap_from.get("timestamp"),
            "source":    snap_from.get("source"),
        },
        "snap_to": {
            "id":        snap_to.get("id"),
            "date":      snap_to.get("date"),
            "timestamp": snap_to.get("timestamp"),
            "source":    snap_to.get("source"),
        },
        "summary_diff": {
            "nav_from":        round(nav_from, 2),
            "nav_to":          round(nav_to, 2),
            "delta_nav":       round(nav_to - nav_from, 2),
            "delta_nav_pct":   round((nav_to / nav_from - 1) * 100, 4) if nav_from > 0 else 0,
            "n_pos_from":      m_from["num_positions"],
            "n_pos_to":        m_to["num_positions"],
            "delta_n_pos":     m_to["num_positions"] - m_from["num_positions"],
            "pct_g1_from":     m_from["pct_grupo1"],
            "pct_g1_to":       m_to["pct_grupo1"],
            "hhi_ativo_from":  m_from["hhi_ativo"],
            "hhi_ativo_to":    m_to["hhi_ativo"],
            "delta_hhi_ativo": m_to["hhi_ativo"] - m_from["hhi_ativo"],
            "hhi_setor_from":  m_from["hhi_setor"],
            "hhi_setor_to":    m_to["hhi_setor"],
            "delta_hhi_setor": m_to["hhi_setor"] - m_from["hhi_setor"],
        },
        "posicoes":     posicoes,
        "setores_diff": setores_diff,
        "n_novos":      sum(1 for p in posicoes if p["status"] == "novo"),
        "n_removidos":  sum(1 for p in posicoes if p["status"] == "removido"),
        "n_alterados":  sum(1 for p in posicoes if p["status"] in ("aumentou", "reduziu")),
    }


@app.route("/api/portfolio-history/diff", methods=["GET"])
@require_admin
def api_portfolio_history_diff():
    """Diff entre dois snapshots por data. Params: from=YYYY-MM-DD, to=YYYY-MM-DD."""
    d_from = (request.args.get("from") or "").strip()
    d_to   = (request.args.get("to")   or "").strip()
    if not d_from or not d_to:
        return jsonify({"error": "params from e to obrigatórios (YYYY-MM-DD)"}), 400
    history = load_portfolio_history()
    if not history:
        return jsonify({"error": "sem snapshots no histórico"}), 404
    snap_from = _ph_find_snap_by_date(history, d_from)
    snap_to   = _ph_find_snap_by_date(history, d_to)
    if not snap_from or not snap_to:
        return jsonify({"error": "snapshots não encontrados para uma das datas"}), 404
    return jsonify(_diff_snapshots(snap_from, snap_to))


@app.route("/api/portfolio-history/timeseries", methods=["GET"])
@require_admin
def api_portfolio_history_timeseries():
    """Time series de métricas do portfolio_history.
    Sempre retorna todas — frontend escolhe quais plotar."""
    history = load_portfolio_history()
    history_sorted = sorted(history, key=lambda s: s.get("timestamp") or "")

    # Coletar todos os setores únicos
    setores_all = set()
    for snap in history_sorted:
        m = _ph_extract_metrics(snap)
        setores_all.update(m["setor_pcts"].keys())
    setores_all = sorted(setores_all)

    points = []
    for snap in history_sorted:
        m = _ph_extract_metrics(snap)
        points.append({
            "date":          snap.get("date"),
            "timestamp":     snap.get("timestamp"),
            "nav":           m["nav"],
            "num_positions": m["num_positions"],
            "hhi_ativo":     m["hhi_ativo"],
            "hhi_setor":     m["hhi_setor"],
            "weighted_beta": m["weighted_beta"],
            "weighted_upside": m["weighted_upside"],
            "pct_grupo1":    m["pct_grupo1"],
            "setor_pcts":    {s: m["setor_pcts"].get(s, 0) for s in setores_all},
        })
    return jsonify({
        "points":      points,
        "all_sectors": setores_all,
        "total":       len(points),
    })


def _infer_operations(snap_from, snap_to):
    """Detecta variações de quantidade entre dois snapshots e infere operações."""
    rows_from = {r["ticker"]: r for r in (snap_from.get("rows") or [])}
    rows_to   = {r["ticker"]: r for r in (snap_to.get("rows") or [])}
    all_tickers = sorted(set(rows_from) | set(rows_to))

    ops = []
    for tk in all_tickers:
        rf = rows_from.get(tk) or {}
        rt = rows_to.get(tk)   or {}
        qf = rf.get("quantidade") or 0
        qt = rt.get("quantidade") or 0
        delta = qt - qf
        if delta == 0:
            continue
        # preço médio: média simples dos dois preços (aproximação)
        p_from = rf.get("preco")
        p_to   = rt.get("preco")
        if p_from and p_to:
            preco_estimado = (p_from + p_to) / 2.0
        else:
            preco_estimado = p_to or p_from
        valor_estimado = abs(delta) * (preco_estimado or 0)
        if delta > 0:
            direcao = "compra"
        elif qt == 0:
            direcao = "zerou"
        else:
            direcao = "venda"
        ops.append({
            "ticker":          tk,
            "direcao":         direcao,
            "qtde_from":       qf,
            "qtde_to":         qt,
            "delta_qtde":      delta,
            "preco_estimado":  round(preco_estimado, 2) if preco_estimado else None,
            "valor_estimado":  round(valor_estimado, 2),
        })
    return ops


def _crossmatch_pretrade(operations, pretrade_history, date_from, date_to):
    """Para cada operação inferida, tenta encontrar match em pretrade_history.executed_at
    dentro da janela [date_from, date_to]. Match por ticker + direção + delta qtde compatível."""
    # Pré-filtrar pretrade_history executados dentro da janela
    candidates = []
    for rec in pretrade_history:
        exec_at = rec.get("executed_at")
        if not exec_at:
            continue
        exec_date = exec_at[:10]
        if not (date_from <= exec_date <= date_to):
            continue
        for op in (rec.get("operacoes") or []):
            candidates.append({
                "pretrade_id": rec.get("id"),
                "executed_at": exec_at,
                "ticker":      (op.get("ticker") or "").replace(".SA", "").upper(),
                "direcao":     op.get("direcao"),
                "quantidade":  op.get("quantidade") or 0,
            })

    for inferred in operations:
        tk = inferred["ticker"].upper()
        dir_inferred = inferred["direcao"]
        # zerou matcheia com venda no pretrade
        dir_search = "venda" if dir_inferred in ("venda", "zerou") else "compra"
        matches = [c for c in candidates
                   if c["ticker"] == tk and c["direcao"] == dir_search]
        if matches:
            # Pega o de quantidade mais próxima
            matches.sort(key=lambda c: abs(c["quantidade"] - abs(inferred["delta_qtde"])))
            best = matches[0]
            inferred["pretrade_id"]     = best["pretrade_id"]
            inferred["pretrade_exec_at"] = best["executed_at"]
            inferred["rastreado"]       = True
        else:
            inferred["pretrade_id"]      = None
            inferred["pretrade_exec_at"] = None
            inferred["rastreado"]        = False
    return operations


@app.route("/api/portfolio-history/operations", methods=["GET"])
@require_admin
def api_portfolio_history_operations():
    """Operações inferidas entre snapshots consecutivos no histórico, com
    cross-match contra pretrade_history (mostra quais foram via pré-trade
    e quais foram alterações manuais sem registro)."""
    d_from = (request.args.get("from") or "").strip()
    d_to   = (request.args.get("to")   or "").strip()
    history = load_portfolio_history()
    if not history:
        return jsonify({"operations": [], "total": 0})

    history_sorted = sorted(history, key=lambda s: s.get("timestamp") or "")
    # Filtrar por data se passada
    if d_from or d_to:
        history_sorted = [s for s in history_sorted
                          if (not d_from or (s.get("date") or "") >= d_from)
                          and (not d_to or (s.get("date") or "") <= d_to)]

    pretrade_h = load_pretrade_history()
    range_from = d_from or (history_sorted[0]["date"] if history_sorted else "")
    range_to   = d_to   or (history_sorted[-1]["date"] if history_sorted else "")

    all_ops = []
    for i in range(1, len(history_sorted)):
        prev = history_sorted[i-1]
        curr = history_sorted[i]
        ops_pair = _infer_operations(prev, curr)
        for op in ops_pair:
            op["date_from"] = prev.get("date")
            op["date_to"]   = curr.get("date")
            op["snap_id_from"] = prev.get("id")
            op["snap_id_to"]   = curr.get("id")
        all_ops.extend(ops_pair)

    # Cross-match com pretrade
    _crossmatch_pretrade(all_ops, pretrade_h, range_from, range_to)

    # Ordenar por data (descendente)
    all_ops.sort(key=lambda o: (o.get("date_to") or "", o.get("ticker") or ""), reverse=True)
    n_rastreadas = sum(1 for o in all_ops if o["rastreado"])
    return jsonify({
        "operations":         all_ops,
        "total":              len(all_ops),
        "n_rastreadas":       n_rastreadas,
        "n_nao_rastreadas":   len(all_ops) - n_rastreadas,
        "range_from":         range_from,
        "range_to":           range_to,
    })


@app.route("/api/portfolio-history/diff/pdf", methods=["GET"])
@require_admin
def api_portfolio_history_diff_pdf():
    """PDF de auditoria de diff entre 2 datas."""
    d_from = (request.args.get("from") or "").strip()
    d_to   = (request.args.get("to")   or "").strip()
    if not d_from or not d_to:
        return jsonify({"error": "params from e to obrigatórios"}), 400
    history = load_portfolio_history()
    snap_from = _ph_find_snap_by_date(history, d_from)
    snap_to   = _ph_find_snap_by_date(history, d_to)
    if not snap_from or not snap_to:
        return jsonify({"error": "snapshots não encontrados"}), 404
    diff = _diff_snapshots(snap_from, snap_to)
    try:
        pdf_bytes = _generate_portfolio_diff_pdf(diff)
    except Exception as e:
        return jsonify({"error": f"Erro ao gerar PDF: {str(e)}"}), 500
    filename = f"carteira_diff_{snap_from.get('date')}_a_{snap_to.get('date')}.pdf"
    from flask import make_response
    resp = make_response(pdf_bytes)
    resp.headers["Content-Type"]        = "application/pdf"
    resp.headers["Content-Disposition"] = f'attachment; filename="{filename}"'
    return resp


def _generate_portfolio_diff_pdf(diff):
    """Gera PDF de auditoria do diff entre 2 snapshots da carteira."""
    from io import BytesIO
    from reportlab.lib import colors
    from reportlab.lib.pagesizes import A4, landscape
    from reportlab.lib.styles import ParagraphStyle
    from reportlab.lib.units import cm
    from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle, HRFlowable

    buf = BytesIO()
    doc = SimpleDocTemplate(buf, pagesize=landscape(A4),
                            leftMargin=1.2*cm, rightMargin=1.2*cm,
                            topMargin=1.2*cm, bottomMargin=1.2*cm)

    C_BG  = colors.HexColor("#0d0d1a"); C_HDR = colors.HexColor("#1a1a2e")
    C_TXT = colors.white; C_OK = colors.HexColor("#00cc88")
    C_NEG = colors.HexColor("#cc3333"); C_WARN = colors.HexColor("#f5a623")
    C_BODY = colors.HexColor("#cccccc"); C_MUTED = colors.HexColor("#888888")
    C_LINE = colors.HexColor("#333333")

    st_title = ParagraphStyle("title", fontName="Helvetica-Bold", fontSize=14, textColor=C_TXT, spaceAfter=4)
    st_sub   = ParagraphStyle("sub",   fontName="Helvetica",      fontSize=9,  textColor=C_MUTED, spaceAfter=2)
    st_sec   = ParagraphStyle("sec",   fontName="Helvetica-Bold", fontSize=10, textColor=C_TXT, spaceBefore=10, spaceAfter=4)
    st_foot  = ParagraphStyle("foot",  fontName="Helvetica",      fontSize=7,  textColor=C_MUTED, alignment=1, spaceBefore=6)

    def tbl_style(header_rows=1, extra=None):
        base = [
            ("BACKGROUND", (0,0), (-1,header_rows-1), C_HDR),
            ("TEXTCOLOR",  (0,0), (-1,header_rows-1), C_TXT),
            ("FONTNAME",   (0,0), (-1,header_rows-1), "Helvetica-Bold"),
            ("FONTSIZE",   (0,0), (-1,-1), 7),
            ("FONTNAME",   (0,header_rows), (-1,-1), "Courier"),
            ("TEXTCOLOR",  (0,header_rows), (-1,-1), C_BODY),
            ("BACKGROUND", (0,header_rows), (-1,-1), C_BG),
            ("ROWBACKGROUNDS", (0,header_rows), (-1,-1), [C_BG, colors.HexColor("#0f0f22")]),
            ("GRID", (0,0), (-1,-1), 0.3, C_LINE),
            ("TOPPADDING",  (0,0), (-1,-1), 3),
            ("BOTTOMPADDING", (0,0), (-1,-1), 3),
            ("LEFTPADDING", (0,0), (-1,-1), 4),
            ("RIGHTPADDING", (0,0), (-1,-1), 4),
        ]
        if extra: base.extend(extra)
        return TableStyle(base)

    def fmt(v, dec=2):
        try: return f"{float(v):,.{dec}f}".replace(",", "X").replace(".", ",").replace("X", ".")
        except: return str(v) if v is not None else "—"

    def fmtR(v, dec=2):
        return f"R$ {fmt(v, dec)}" if v is not None else "—"

    sf, st = diff["snap_from"], diff["snap_to"]
    sd = diff["summary_diff"]

    story = []
    story.append(Paragraph("HARBOUR IAT FIF AÇÕES RL — DIFF DE CARTEIRA", st_title))
    story.append(Paragraph(
        f"De: {sf.get('date')} ({sf.get('source','—')})  →  Para: {st.get('date')} ({st.get('source','—')})",
        st_sub
    ))
    story.append(HRFlowable(width="100%", thickness=0.5, color=C_LINE, spaceAfter=6))

    # ── Sumário ──
    story.append(Paragraph("1. SUMÁRIO DAS MUDANÇAS", st_sec))
    sum_rows = [
        ["MÉTRICA", "DE", "PARA", "Δ"],
        ["NAV Total",         fmtR(sd["nav_from"]),     fmtR(sd["nav_to"]),     f"{fmtR(sd['delta_nav'])} ({fmt(sd['delta_nav_pct'],2)}%)"],
        ["Nº Posições",       str(sd["n_pos_from"]),    str(sd["n_pos_to"]),    f"{sd['delta_n_pos']:+d}"],
        ["Grupo I (Ações/BDR) %", f"{fmt(sd['pct_g1_from'])}%", f"{fmt(sd['pct_g1_to'])}%", f"{fmt(sd['pct_g1_to'] - sd['pct_g1_from'])} pp"],
        ["HHI Ativo",         str(sd["hhi_ativo_from"]), str(sd["hhi_ativo_to"]), f"{sd['delta_hhi_ativo']:+d}"],
        ["HHI Setor",         str(sd["hhi_setor_from"]), str(sd["hhi_setor_to"]), f"{sd['delta_hhi_setor']:+d}"],
        ["Novos / Removidos / Alterados", "—", "—",
         f"+{diff['n_novos']}  / -{diff['n_removidos']}  / Δ{diff['n_alterados']}"],
    ]
    sum_tbl = Table(sum_rows, colWidths=[7*cm, 4.5*cm, 4.5*cm, 5*cm])
    sum_tbl.setStyle(tbl_style())
    story.append(sum_tbl)
    story.append(Spacer(1, 6))

    # ── Posições alteradas ──
    story.append(Paragraph("2. POSIÇÕES (ordenadas por |Δ valor| desc)", st_sec))
    pos_rows = [["ATIVO", "STATUS", "QTDE DE", "QTDE PARA", "Δ QTDE", "VALOR DE", "VALOR PARA", "Δ VALOR", "% DE", "% PARA", "Δ pp"]]
    pos_extra = []
    for i, p in enumerate(diff["posicoes"], start=1):
        st_color = C_OK if p["status"] == "novo" else C_NEG if p["status"] == "removido" else C_WARN if p["status"] in ("aumentou","reduziu") else C_MUTED
        pos_rows.append([
            p["ticker"],
            p["status"].upper(),
            fmt(p["qtde_from"], 0),
            fmt(p["qtde_to"], 0),
            f"{int(p['delta_qtde']):+d}",
            fmtR(p["valor_from"], 0),
            fmtR(p["valor_to"], 0),
            fmtR(p["delta_valor"], 0),
            f"{fmt(p['pct_from'])}%",
            f"{fmt(p['pct_to'])}%",
            f"{p['delta_pct_pp']:+.2f}",
        ])
        pos_extra.append(("TEXTCOLOR", (1, i), (1, i), st_color))
    pos_tbl = Table(pos_rows, colWidths=[2*cm, 2*cm, 2*cm, 2*cm, 1.8*cm, 2.5*cm, 2.5*cm, 2.5*cm, 1.5*cm, 1.5*cm, 1.5*cm], repeatRows=1)
    pos_tbl.setStyle(tbl_style(extra=pos_extra))
    story.append(pos_tbl)
    story.append(Spacer(1, 6))

    # ── Setores ──
    story.append(Paragraph("3. ROTAÇÃO SETORIAL", st_sec))
    sec_rows = [["SETOR", "% DE", "% PARA", "Δ pp"]]
    for s in diff["setores_diff"]:
        sec_rows.append([s["sector"], f"{fmt(s['pct_from'])}%", f"{fmt(s['pct_to'])}%", f"{s['delta_pp']:+.2f}"])
    sec_tbl = Table(sec_rows, colWidths=[8*cm, 3*cm, 3*cm, 3*cm])
    sec_tbl.setStyle(tbl_style())
    story.append(sec_tbl)

    story.append(Spacer(1, 8))
    story.append(Paragraph(
        "Diff baseado em snapshots de portfolio_history.json. NOVOS = entraram entre as datas; "
        "REMOVIDOS = sairam; AUMENTOU/REDUZIU = mudança de quantidade.",
        st_foot
    ))

    doc.build(story)
    return buf.getvalue()


# ---------------------------------------------------------------------------
# Index Members (Bloomberg-style Leaders / Laggards)
# ---------------------------------------------------------------------------

INDEX_TICKER_MAP = {
    "IBOV":    "^BVSP",
    "IBRX100": "BRAX11.SA",
    "SMLL":    "SMAL11.SA",
}
INDEX_MEMBERS_TTL = 3600  # 1 hour cache

@app.route("/api/index-members")
def api_index_members():
    """GET /api/index-members?index=IBOV&start=2026-01-01&end=2026-04-02"""
    import pandas as pd

    index_name = request.args.get("index", "IBOV").upper()
    start_str  = request.args.get("start", "")
    end_str    = request.args.get("end", "")

    if index_name not in INDEX_TICKER_MAP:
        return jsonify({"error": f"Índice inválido: {index_name}"}), 400

    if not os.path.exists(INDEX_MEMBERS_FILE):
        return jsonify({"error": "index_members.json não encontrado"}), 404

    with open(INDEX_MEMBERS_FILE, "r", encoding="utf-8") as f:
        all_members = json.load(f)

    members = all_members.get(index_name, [])
    if not members:
        return jsonify({"error": f"Índice {index_name} sem constituintes"}), 404

    # Parse dates
    try:
        dt_start = datetime.strptime(start_str, "%Y-%m-%d") if start_str else datetime(datetime.now().year, 1, 1)
        dt_end   = datetime.strptime(end_str,   "%Y-%m-%d") if end_str   else datetime.now()
    except ValueError:
        return jsonify({"error": "Formato de data inválido (use AAAA-MM-DD)"}), 400

    # Cache lookup
    cache_key = f"idx_members_{index_name}_{dt_start.date()}_{dt_end.date()}"
    cache = load_cache()
    now   = time.time()
    if cache.get(cache_key) and now < cache[cache_key].get("expires_at", 0):
        return jsonify(cache[cache_key]["data"])

    # Normalise weights to sum to 100
    total_w = sum(m.get("weight_pct", 0) for m in members)
    if total_w > 0:
        for m in members:
            m["weight_pct"] = round(m.get("weight_pct", 0) / total_w * 100, 4)

    # Build tickers list (members + index level)
    index_ticker = INDEX_TICKER_MAP[index_name]
    member_tickers = [m["ticker"] for m in members]
    all_tickers = member_tickers + [index_ticker]

    # Portfolio tickers for cross-reference
    try:
        portfolio = load_portfolio()
        portfolio_tickers = set(
            p.get("yahoo_ticker") or p.get("ticker", "")
            for p in portfolio.get("positions", [])
        )
    except Exception:
        portfolio_tickers = set()

    # Fetch price history in bulk
    try:
        df = yf.download(
            all_tickers,
            start=dt_start,
            end=dt_end + timedelta(days=1),
            auto_adjust=True,
            progress=False,
            threads=True,
        )
    except Exception as e:
        return jsonify({"error": f"Erro ao buscar dados: {str(e)}"}), 500

    # Extract Close prices
    if isinstance(df.columns, pd.MultiIndex):
        close = df["Close"] if "Close" in df.columns.get_level_values(0) else df.iloc[:, :]
    else:
        close = df

    def get_series(ticker):
        if ticker in close.columns:
            return close[ticker].dropna()
        return pd.Series(dtype=float)

    # Index level start/end
    idx_series = get_series(index_ticker)
    idx_start_val = float(idx_series.iloc[0])  if len(idx_series) >= 2 else None
    idx_end_val   = float(idx_series.iloc[-1]) if len(idx_series) >= 2 else None
    idx_change    = (idx_end_val - idx_start_val) if idx_start_val and idx_end_val else None
    idx_return_pct = ((idx_end_val / idx_start_val) - 1) * 100 if idx_start_val and idx_end_val else None

    results = []
    for m in members:
        ticker = m["ticker"]
        name   = m["name"]
        weight = m["weight_pct"]

        series = get_series(ticker)
        if len(series) < 2:
            continue

        start_price = float(series.iloc[0])
        end_price   = float(series.iloc[-1])
        if start_price == 0:
            continue
        change_pct = (end_price / start_price - 1) * 100

        # Points contributed to index (in index units)
        if idx_start_val and idx_start_val > 0:
            points = (change_pct / 100) * (weight / 100) * idx_start_val
        else:
            points = 0.0

        # % of total index move from this stock
        if idx_return_pct and idx_return_pct != 0:
            idx_mv_pct = ((change_pct / 100) * (weight / 100)) / (idx_return_pct / 100) * 100
        else:
            idx_mv_pct = 0.0

        results.append({
            "ticker":      ticker,
            "name":        name,
            "weight_pct":  round(weight, 2),
            "end_price":   round(end_price, 2),
            "change_pct":  round(change_pct, 2),
            "points":      round(points, 3),
            "idx_mv_pct":  round(idx_mv_pct, 2),
            "in_portfolio": ticker in portfolio_tickers,
        })

    results.sort(key=lambda x: x["change_pct"], reverse=True)
    max_rows = 35
    leaders  = results[:max_rows]
    laggards = sorted(results, key=lambda x: x["change_pct"])[:max_rows]

    data = {
        "index":        index_name,
        "start":        dt_start.strftime("%Y-%m-%d"),
        "end":          dt_end.strftime("%Y-%m-%d"),
        "idx_start":    round(idx_start_val, 2)  if idx_start_val  else None,
        "idx_end":      round(idx_end_val, 2)    if idx_end_val    else None,
        "idx_change":   round(idx_change, 2)     if idx_change     else None,
        "idx_return_pct": round(idx_return_pct, 2) if idx_return_pct else None,
        "n_total":      len(results),
        "n_up":         sum(1 for r in results if r["change_pct"] > 0),
        "n_down":       sum(1 for r in results if r["change_pct"] < 0),
        "n_unch":       sum(1 for r in results if r["change_pct"] == 0),
        "leaders":      leaders,
        "laggards":     laggards,
    }

    cache[cache_key] = {"data": data, "expires_at": now + INDEX_MEMBERS_TTL}
    save_cache(cache)
    return jsonify(data)


# =============================================================================
# CVM OFICIAL — Informe Diário do fundo (HARBOUR IAT FIF)
# =============================================================================

import cvm_daily_fetcher as _cvm_daily

# Refresh do informe diário CVM roda no GitHub Actions (.github/workflows/cvm-daily.yml),
# que executa fetcher no runner e dá commit do data/cvm_daily.json. O Flask só serve o
# JSON pronto — sem thread daemon dentro do worker, sem download de ZIP em memória.


def _cvm_daily_summary(records, quota_history):
    """Computes aggregates (captação líq 30d/YTD/12m, var cotistas, diff cota) + cota history lookup."""
    if not records:
        return {
            "total_rows": 0,
            "cota_cvm_atual": None,
            "data_cota": None,
            "pl_atual": None,
            "nr_cotst_atual": None,
            "captc_liq_30d": 0.0,
            "captc_liq_ytd": 0.0,
            "captc_liq_12m": 0.0,
            "var_cotst_30d": 0,
            "var_cotst_ytd": 0,
            "diff_cota_pct": None,
        }

    records_sorted = sorted(records, key=lambda r: r["dt_comptc"])
    last = records_sorted[-1]
    first = records_sorted[0]

    last_date = datetime.strptime(last["dt_comptc"], "%Y-%m-%d").date()
    d30 = last_date - timedelta(days=30)
    d365 = last_date - timedelta(days=365)
    year_start = last_date.replace(month=1, day=1)

    def _sum_net(records_slice):
        return sum((r.get("captc_dia") or 0) - (r.get("resg_dia") or 0) for r in records_slice)

    slice_30 = [r for r in records_sorted if r["dt_comptc"] >= d30.isoformat()]
    slice_ytd = [r for r in records_sorted if r["dt_comptc"] >= year_start.isoformat()]
    slice_12m = [r for r in records_sorted if r["dt_comptc"] >= d365.isoformat()]

    # Var cotistas
    def _first_of(slice_):
        return slice_[0] if slice_ else None

    r30 = _first_of(slice_30)
    rytd = _first_of(slice_ytd)
    var_cot_30 = (last.get("nr_cotst") or 0) - (r30.get("nr_cotst") or 0) if r30 else 0
    var_cot_ytd = (last.get("nr_cotst") or 0) - (rytd.get("nr_cotst") or 0) if rytd else 0

    # Diff cota (oficial x calculada) — procura o dia mais recente em que ambos existem
    diff_cota_pct = None
    diff_cota_data = None
    if quota_history:
        try:
            calc_map = {}
            for q in quota_history:
                d = q.get("data")
                v = q.get("cota_fechamento") or q.get("quota")
                if d and v:
                    calc_map[d] = float(v)
            for rec in reversed(records_sorted):
                d = rec.get("dt_comptc")
                cota_cvm = rec.get("vl_quota")
                if d in calc_map and cota_cvm:
                    cota_calc = calc_map[d]
                    if cota_calc:
                        diff_cota_pct = (float(cota_cvm) - cota_calc) / cota_calc * 100.0
                        diff_cota_data = d
                        break
        except Exception:
            diff_cota_pct = None
            diff_cota_data = None

    return {
        "total_rows": len(records_sorted),
        "cota_cvm_atual": last.get("vl_quota"),
        "data_cota": last.get("dt_comptc"),
        "pl_atual": last.get("vl_patrim_liq"),
        "nr_cotst_atual": last.get("nr_cotst"),
        "captc_liq_30d": _sum_net(slice_30),
        "captc_liq_ytd": _sum_net(slice_ytd),
        "captc_liq_12m": _sum_net(slice_12m),
        "var_cotst_30d": var_cot_30,
        "var_cotst_ytd": var_cot_ytd,
        "diff_cota_pct": diff_cota_pct,
        "diff_cota_data": diff_cota_data,
        "first_record": first.get("dt_comptc"),
    }


@app.route("/api/cvm/fund-daily", methods=["GET"])
def api_cvm_fund_daily():
    """Returns the full set of CVM daily records + registration metadata."""
    storage = _cvm_daily.load_storage()
    cadastro = _cvm_daily.load_cadastro()
    status = _cvm_daily.get_status()
    return jsonify({
        "cnpj":        storage.get("cnpj"),
        "cota_inicio": storage.get("cota_inicio"),
        "last_refresh": storage.get("last_refresh"),
        "records":     storage.get("records", []),
        "cadastro":    cadastro,
        "scheduler":   status,
    })


@app.route("/api/cvm/fund-daily/summary", methods=["GET"])
def api_cvm_fund_daily_summary():
    """Quick summary + last 30 records for cards/charts."""
    storage = _cvm_daily.load_storage()
    records = storage.get("records") or []
    # carrega quota_history calculada para comparação
    qh_path = os.path.join(BASE_DIR, "data", "quota_history.json")
    quota_history = []
    if os.path.exists(qh_path):
        try:
            with open(qh_path, "r", encoding="utf-8") as f:
                quota_history = json.load(f)
        except Exception:
            quota_history = []

    summary = _cvm_daily_summary(records, quota_history)
    summary["last_refresh"] = storage.get("last_refresh")
    summary["cnpj"] = storage.get("cnpj")
    return jsonify(summary)


@app.route("/api/cvm/fund-daily/refresh", methods=["POST"])
@require_admin
def api_cvm_fund_daily_refresh():
    """Refresh manual do mês atual + M-1. Síncrono — para uso eventual via UI admin.
    O refresh diário automático roda no GitHub Actions (cvm-daily.yml) para não
    consumir memória do worker do Render."""
    result = _cvm_daily.refresh_current()
    return jsonify({"ok": True, "mode": "refresh", "result": result})


@app.route("/api/cvm/fund-daily/backfill", methods=["POST"])
@require_admin
def api_cvm_fund_daily_backfill():
    """Backfill completo desde a cota_inicio. Síncrono — só dispare em ambiente
    com folga de memória (preferir rodar local ou via GitHub Actions, nunca no
    Render Starter durante operação)."""
    result = _cvm_daily.backfill_since(_cvm_daily.COTA_INICIO[:7])
    return jsonify({"ok": True, "mode": "backfill", "result": result})


@app.route("/api/cvm/fund-daily/status", methods=["GET"])
def api_cvm_fund_daily_status():
    return jsonify(_cvm_daily.get_status())


# ─── LIQUIDEZ ─────────────────────────────────────────────────────────────────

LIQUIDITY_HISTORY_FILE = os.path.join(DATA_DIR, "liquidity_history.json")

# Buckets de dias úteis (mesmos do sistema de referência do usuário)
LIQUIDEZ_BUCKETS = [1, 2, 3, 4, 5, 10, 21, 30, 42, 63, 84, 105, 126, 180, 252, 360, 540]

# Prazo de resgate default por categoria (D+x)
PRAZO_RESGATE_POR_CATEGORIA = {
    "Acao": 2, "BDR": 2, "Acao BDR": 2,
    "FundoRF": 0, "FundoMM": 1, "Caixa": 0,
}

# Cenários: (volume_disponivel_pct, percentile_resgate)
LIQUIDEZ_CENARIOS = {
    "neutro": {"vol_mult": 1.00, "percentile": 50, "label": "Neutro"},
    "stress": {"vol_mult": 0.50, "percentile": 75, "label": "Stress"},
    "crise":  {"vol_mult": 0.30, "percentile": 95, "label": "Crise"},
}

# Faixas de classificação
def _classify_liquidity(days):
    if days is None: return "sem_dados"
    if days < 3:  return "alta"
    if days <= 7: return "media"
    if days <= 30: return "baixa"
    return "muito_baixa"


def _prazo_resgate_position(pos):
    """Retorna o prazo de resgate em dias úteis para a posição.
    Usa override pos['prazo_resgate_d'] se setado, senão default por categoria."""
    override = pos.get("prazo_resgate_d")
    if override is not None:
        try:
            return int(override)
        except (TypeError, ValueError):
            pass
    cat = pos.get("categoria") or "Acao"
    return PRAZO_RESGATE_POR_CATEGORIA.get(cat, 2)


def _calc_days_market_only(valor, avg_vol_rs, volume_mult=0.20):
    """Dias necessários para EXECUTAR as vendas (sem contar settlement).
    Cada dia útil de venda absorve volume_mult × volume_médio do ativo.
    Default: 20% do ADV — teto seguro de participação para não impactar preço."""
    if not valor or valor <= 0:
        return 0.0
    if not avg_vol_rs or avg_vol_rs <= 0:
        return 0.0  # sem dado → assume liquidez instantânea no mercado
    return valor / (avg_vol_rs * volume_mult)


def _calc_days_to_liquidate(valor, avg_vol_rs, prazo_resgate, volume_mult=0.20):
    """Total de dias úteis até 100% da posição estar liquidada EM CAIXA.
    = prazo_settlement + (dias de execução - 1, se > 1)
    Para ativos onde 1 dia de venda já basta: total = prazo_settlement.
    Para ativos pouco líquidos: total = settlement + extra dias necessários."""
    days_market = _calc_days_market_only(valor, avg_vol_rs, volume_mult)
    return prazo_resgate + max(days_market - 1, 0.0)


def _liquidatable_fraction(B, prazo_resgate, days_market):
    """Fração da posição (0.0 a 1.0) liquidada e em caixa até o bucket B (dias úteis).
    Modelo settlement-aware: durante prazo_settlement, caixa ainda não chegou (0%).
    A partir daí, cada novo dia de bucket adiciona 1 fatia das vendas que já settled."""
    if B < prazo_resgate:
        return 0.0
    sell_days_settled = B - prazo_resgate + 1
    # max(days_market, 1.0) garante que ativo super-líquido (days_market < 1) seja
    # contado como 1 dia de venda — o caixa total chega em D+settlement.
    return min(sell_days_settled / max(days_market, 1.0), 1.0)


def _percentile(sorted_vals, p):
    """Percentil simples (linear interpolation). p em 0-100."""
    if not sorted_vals:
        return 0.0
    if len(sorted_vals) == 1:
        return sorted_vals[0]
    k = (len(sorted_vals) - 1) * (p / 100.0)
    f = int(k)
    c = min(f + 1, len(sorted_vals) - 1)
    if f == c:
        return sorted_vals[f]
    return sorted_vals[f] + (sorted_vals[c] - sorted_vals[f]) * (k - f)


def _calc_redemption_curve(cvm_records, buckets, percentile=50):
    """Para cada bucket B (dias), retorna o percentil das somas rolantes
    de resg_pct sobre janelas de B dias úteis no histórico CVM.

    resg_pct_d = resg_dia / vl_patrim_liq do dia.
    Janela rolante de B dias úteis → soma → percentile_p sobre todas as janelas.
    """
    if not cvm_records:
        return [0.0] * len(buckets), 0

    # Ordenar por data e calcular resg_pct diário
    records_sorted = sorted(cvm_records, key=lambda r: r.get("dt_comptc", ""))
    resg_pct_series = []
    for r in records_sorted:
        pl = r.get("vl_patrim_liq") or 0
        resg = r.get("resg_dia") or 0
        if pl > 0:
            resg_pct_series.append(resg / pl * 100.0)  # em %
        else:
            resg_pct_series.append(0.0)

    n = len(resg_pct_series)
    out = []
    for B in buckets:
        B = min(B, n)  # cap na qtde de records disponíveis
        if B <= 0:
            out.append(0.0)
            continue
        # Janelas rolantes de B dias úteis — soma de cada janela
        windows = []
        for i in range(n - B + 1):
            windows.append(sum(resg_pct_series[i:i+B]))
        if not windows:
            out.append(0.0)
            continue
        windows.sort()
        out.append(round(_percentile(windows, percentile), 4))

    return out, n


def _build_liquidity_snapshot(scenario="neutro"):
    """Calcula o snapshot completo da liquidez do fundo no cenário pedido.
    Retorna dict pronto para serializar como JSON."""
    cen = LIQUIDEZ_CENARIOS.get(scenario, LIQUIDEZ_CENARIOS["neutro"])
    volume_mult = 0.20 * cen["vol_mult"]

    portfolio    = load_portfolio()
    tickers      = [p["yahoo_ticker"] for p in portfolio["positions"]]
    prices       = get_cached_prices(tickers)
    fundamentals = get_cached_fundamentals(tickers)

    # ── Por ativo ──
    por_ativo = []
    nav = 0.0
    for pos in portfolio["positions"]:
        yh    = pos["yahoo_ticker"]
        price = (prices.get(yh) or {}).get("price")
        qtde  = pos.get("quantidade") or 0
        valor = round(price * qtde, 2) if price else 0.0
        nav  += valor

        avg_vol = (fundamentals.get(yh) or {}).get("average_volume")
        avg_vol_rs = (avg_vol * price) if (avg_vol and price) else None

        prazo_resgate = _prazo_resgate_position(pos)
        days_market   = _calc_days_market_only(valor, avg_vol_rs, volume_mult)
        days          = _calc_days_to_liquidate(valor, avg_vol_rs, prazo_resgate, volume_mult)

        # Liquidez por bucket — RAMPA settlement-aware:
        # B < prazo_settlement: 0% (caixa ainda não chegou)
        # B >= settlement: rampa linear conforme vendas vão settling
        # Mais realista que step-function (que pulava de 0 a 100% num bucket só).
        liq_por_bucket = [
            _liquidatable_fraction(B, prazo_resgate, days_market) * 100.0
            for B in LIQUIDEZ_BUCKETS
        ]

        por_ativo.append({
            "ticker":          pos["ticker"],
            "yahoo_ticker":    yh,
            "categoria":       pos.get("categoria") or "Acao",
            "prazo_resgate_d": prazo_resgate,
            "valor_bruto":     valor,
            "avg_vol_rs":      round(avg_vol_rs, 2) if avg_vol_rs else None,
            "dias_zerar":      round(days, 2),
            "liq_por_bucket":  liq_por_bucket,
            "classificacao":   _classify_liquidity(days),
        })

    # Proporção e cumulativo por bucket
    for a in por_ativo:
        a["proporcao_pct"] = round(a["valor_bruto"] / nav * 100, 2) if nav > 0 else 0.0
        # liquidez ponderada (proporção da carteira) por bucket
        a["liq_ponderada_por_bucket"] = [
            round(a["proporcao_pct"] * (frac / 100.0), 4)
            for frac in a["liq_por_bucket"]
        ]

    # Liquidez ativos cumulativa (% do PL) por bucket
    liquidez_ativos = [0.0] * len(LIQUIDEZ_BUCKETS)
    for a in por_ativo:
        for i, lp in enumerate(a["liq_ponderada_por_bucket"]):
            liquidez_ativos[i] += lp
    liquidez_ativos = [round(v, 2) for v in liquidez_ativos]

    # ── Resgate projetado (do histórico CVM) ──
    storage = _cvm_daily.load_storage()
    cvm_records = storage.get("records") or []
    resgate_projetado, n_records = _calc_redemption_curve(
        cvm_records, LIQUIDEZ_BUCKETS, percentile=cen["percentile"]
    )

    # ── Índice liquidez = liquidez_ativos / max(resgate_projetado, 0.01) ──
    indice_liquidez = []
    for la, rp in zip(liquidez_ativos, resgate_projetado):
        denom = max(rp, 0.01)
        indice_liquidez.append(round(la / denom, 4))

    return {
        "scenario":          scenario,
        "scenario_label":    cen["label"],
        "scenario_vol_mult": cen["vol_mult"],
        "scenario_percentile": cen["percentile"],
        "buckets":           LIQUIDEZ_BUCKETS,
        "liquidez_ativos":   liquidez_ativos,
        "resgate_projetado": resgate_projetado,
        "indice_liquidez":   indice_liquidez,
        "por_ativo":         sorted(por_ativo, key=lambda a: a["valor_bruto"], reverse=True),
        "nav":               round(nav, 2),
        "n_records_resgate": n_records,
    }


@app.route("/api/liquidity/snapshot", methods=["GET"])
@require_admin
def api_liquidity_snapshot():
    scenario = (request.args.get("scenario") or "neutro").lower()
    if scenario not in LIQUIDEZ_CENARIOS:
        return jsonify({"error": f"scenario inválido (use: {list(LIQUIDEZ_CENARIOS.keys())})"}), 400
    snap = _build_liquidity_snapshot(scenario=scenario)
    return jsonify(snap)


def _build_liquidity_market(window_days=60):
    """KPIs + faixas + matriz por ativo para a sub-aba MERCADO."""
    snap = _build_liquidity_snapshot(scenario="neutro")
    por_ativo = snap["por_ativo"]
    nav = snap["nav"]

    # Volume médio ponderado: sum(valor_i × avg_vol_i) / NAV
    vol_total = 0.0
    vol_ponderado_num = 0.0
    for a in por_ativo:
        vol = a.get("avg_vol_rs") or 0
        if vol > 0:
            vol_total += vol  # média simples também útil
        vol_ponderado_num += (a["valor_bruto"] or 0) * vol
    vol_medio_ponderado = (vol_ponderado_num / nav) if nav > 0 else 0.0

    # Faixas de liquidez (% do PL por faixa)
    faixas = {"alta": 0.0, "media": 0.0, "baixa": 0.0, "muito_baixa": 0.0, "sem_dados": 0.0}
    for a in por_ativo:
        faixas[a["classificacao"]] = faixas.get(a["classificacao"], 0.0) + (a["proporcao_pct"] or 0.0)
    faixas = {k: round(v, 2) for k, v in faixas.items()}

    pct_alta = faixas.get("alta", 0.0)

    # Prazo médio ponderado = sum(valor × dias) / NAV
    prazo_medio = sum((a["valor_bruto"] or 0) * (a["dias_zerar"] or 0) for a in por_ativo) / nav if nav > 0 else 0.0

    # Matriz por ativo (sub-aba MERCADO)
    matriz = []
    for a in por_ativo:
        vol = a.get("avg_vol_rs") or 0
        pct_vol_diario = (a["valor_bruto"] / vol * 100) if (vol and vol > 0) else None
        matriz.append({
            "ticker":         a["ticker"],
            "valor_carteira": a["valor_bruto"],
            "vol_medio_rs":   vol if vol > 0 else None,
            "pct_vol_diario": round(pct_vol_diario, 2) if pct_vol_diario is not None else None,
            "dias_zerar":     a["dias_zerar"],
            "classificacao":  a["classificacao"],
        })

    # Histórico do prazo médio ponderado (de liquidity_history.json)
    prazo_historico = []
    try:
        if os.path.exists(LIQUIDITY_HISTORY_FILE):
            with open(LIQUIDITY_HISTORY_FILE, "r", encoding="utf-8") as f:
                hist = json.load(f)
            prazo_historico = [
                {"data": h.get("data"), "prazo": h.get("prazo_medio_zerar")}
                for h in hist if h.get("data") and h.get("prazo_medio_zerar") is not None
            ]
            prazo_historico.sort(key=lambda x: x["data"])
    except Exception as e:
        print(f"[liquidity/market] erro lendo history: {e}")

    return {
        "window_days":          window_days,
        "kpis": {
            "valor_carteira":      round(nav, 2),
            "vol_medio_ponderado": round(vol_medio_ponderado, 2),
            "pct_alta_liquidez":   round(pct_alta, 2),
            "prazo_medio_zerar":   round(prazo_medio, 2),
        },
        "faixas":               faixas,
        "prazo_medio_historico": prazo_historico,
        "por_ativo":            matriz,
    }


@app.route("/api/liquidity/market", methods=["GET"])
@require_admin
def api_liquidity_market():
    try:
        window_days = int(request.args.get("window_days") or 60)
    except (TypeError, ValueError):
        window_days = 60
    return jsonify(_build_liquidity_market(window_days=window_days))


def _build_liquidity_compliance():
    """Avalia regras de compliance de liquidez baseadas no snapshot atual."""
    fc = load_fund_config()
    min_5d_pct       = float(fc.get("liquidez_min_5d_pct")       or 80.0)
    max_baixa_pct    = float(fc.get("liquidez_max_baixa_pct")    or 10.0)
    max_zerar_dias   = float(fc.get("liquidez_max_zerar_dias")   or 30.0)

    snap_neutro = _build_liquidity_snapshot(scenario="neutro")
    snap_stress = _build_liquidity_snapshot(scenario="stress")
    market      = _build_liquidity_market()

    # Posição de bucket 5 dias na lista
    idx_5d = LIQUIDEZ_BUCKETS.index(5) if 5 in LIQUIDEZ_BUCKETS else 4
    pct_5d_neutro = snap_neutro["liquidez_ativos"][idx_5d]
    indice_5d_stress = snap_stress["indice_liquidez"][idx_5d]
    pct_baixa_muito = (market["faixas"].get("baixa", 0) + market["faixas"].get("muito_baixa", 0))
    prazo_medio = market["kpis"]["prazo_medio_zerar"]

    def status_min(valor, lim):
        if valor < lim * 0.85: return "violacao"
        if valor < lim:        return "alerta"
        return "ok"

    def status_max(valor, lim):
        if valor > lim:        return "violacao"
        if valor > lim * 0.85: return "alerta"
        return "ok"

    regras = [
        {
            "nome":         f"Mín {min_5d_pct:.0f}% liquidatável em 5 dias úteis",
            "descricao":    "Capacidade de honrar resgates em D+5 (boa prática ANBIMA).",
            "tipo":         "minimo",
            "limite":       min_5d_pct,
            "valor_atual":  pct_5d_neutro,
            "unidade":      "%",
            "status":       status_min(pct_5d_neutro, min_5d_pct),
        },
        {
            "nome":         f"Máx {max_baixa_pct:.0f}% em ativos com prazo > 7 dias",
            "descricao":    "Concentração em ativos pouco líquidos (Baixa + Muito baixa).",
            "tipo":         "maximo",
            "limite":       max_baixa_pct,
            "valor_atual":  round(pct_baixa_muito, 2),
            "unidade":      "%",
            "status":       status_max(pct_baixa_muito, max_baixa_pct),
        },
        {
            "nome":         f"Prazo médio ponderado < {max_zerar_dias:.0f} dias",
            "descricao":    "Tempo médio ponderado para zerar a carteira inteira.",
            "tipo":         "maximo",
            "limite":       max_zerar_dias,
            "valor_atual":  prazo_medio,
            "unidade":      "dias",
            "status":       status_max(prazo_medio, max_zerar_dias),
        },
        {
            "nome":         "Índice Liquidez D+5 em cenário Stress ≥ 1.0",
            "descricao":    "Capacidade de honrar resgate sob stress moderado (50% volume, P75 resgate).",
            "tipo":         "minimo",
            "limite":       1.0,
            "valor_atual":  indice_5d_stress,
            "unidade":      "ratio",
            "status":       status_min(indice_5d_stress, 1.0),
        },
    ]

    pior_status = "ok"
    for r in regras:
        if r["status"] == "violacao":
            pior_status = "violacao"
        elif r["status"] == "alerta" and pior_status == "ok":
            pior_status = "alerta"

    return {
        "regras":      regras,
        "pior_status": pior_status,
        "thresholds":  {
            "liquidez_min_5d_pct":    min_5d_pct,
            "liquidez_max_baixa_pct": max_baixa_pct,
            "liquidez_max_zerar_dias": max_zerar_dias,
        },
    }


@app.route("/api/liquidity/compliance", methods=["GET"])
@require_admin
def api_liquidity_compliance():
    return jsonify(_build_liquidity_compliance())


@app.route("/api/liquidity/history", methods=["GET"])
@require_admin
def api_liquidity_history():
    if not os.path.exists(LIQUIDITY_HISTORY_FILE):
        return jsonify([])
    try:
        with open(LIQUIDITY_HISTORY_FILE, "r", encoding="utf-8") as f:
            return jsonify(json.load(f))
    except Exception as e:
        return jsonify({"error": str(e)}), 500


def load_liquidity_history():
    if not os.path.exists(LIQUIDITY_HISTORY_FILE):
        return []
    with open(LIQUIDITY_HISTORY_FILE, "r", encoding="utf-8") as f:
        return json.load(f)


def save_liquidity_history(data):
    content = json.dumps(data, ensure_ascii=False, indent=2)
    _save_with_github(LIQUIDITY_HISTORY_FILE, "data/liquidity_history.json", content,
                      "chore: update liquidity_history.json via auto-close")


def _generate_liquidity_pdf(snapshot, market, compliance):
    """Gera PDF de auditoria de liquidez. Retorna bytes."""
    from io import BytesIO
    from reportlab.lib import colors
    from reportlab.lib.pagesizes import A4, landscape
    from reportlab.lib.styles import ParagraphStyle
    from reportlab.lib.units import cm
    from reportlab.platypus import (
        SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle, HRFlowable, PageBreak,
    )

    buf = BytesIO()
    doc = SimpleDocTemplate(
        buf, pagesize=landscape(A4),
        leftMargin=1.2*cm, rightMargin=1.2*cm,
        topMargin=1.2*cm, bottomMargin=1.2*cm,
    )

    C_BG      = colors.HexColor("#0d0d1a")
    C_HDR     = colors.HexColor("#1a1a2e")
    C_HDR_TXT = colors.white
    C_OK      = colors.HexColor("#00cc88")
    C_ALERTA  = colors.HexColor("#f5a623")
    C_VIOL    = colors.HexColor("#cc3333")
    C_BODY    = colors.HexColor("#cccccc")
    C_MUTED   = colors.HexColor("#888888")
    C_LINE    = colors.HexColor("#333333")

    def status_color(s):
        return C_OK if s == "ok" else C_ALERTA if s == "alerta" else C_VIOL

    mono = "Courier"
    st_title = ParagraphStyle("title", fontName="Helvetica-Bold", fontSize=14, textColor=C_HDR_TXT, spaceAfter=4)
    st_sub   = ParagraphStyle("sub",   fontName="Helvetica",      fontSize=9,  textColor=C_MUTED, spaceAfter=2)
    st_sec   = ParagraphStyle("sec",   fontName="Helvetica-Bold", fontSize=10, textColor=C_HDR_TXT, spaceBefore=10, spaceAfter=4)
    st_foot  = ParagraphStyle("foot",  fontName="Helvetica",      fontSize=7,  textColor=C_MUTED, alignment=1, spaceBefore=6)

    def tbl_style(header_rows=1, extra=None):
        base = [
            ("BACKGROUND",  (0, 0), (-1, header_rows - 1), C_HDR),
            ("TEXTCOLOR",   (0, 0), (-1, header_rows - 1), C_HDR_TXT),
            ("FONTNAME",    (0, 0), (-1, header_rows - 1), "Helvetica-Bold"),
            ("FONTSIZE",    (0, 0), (-1, -1), 7),
            ("FONTNAME",    (0, header_rows), (-1, -1), mono),
            ("TEXTCOLOR",   (0, header_rows), (-1, -1), C_BODY),
            ("BACKGROUND",  (0, header_rows), (-1, -1), C_BG),
            ("ROWBACKGROUNDS", (0, header_rows), (-1, -1), [C_BG, colors.HexColor("#0f0f22")]),
            ("GRID",        (0, 0), (-1, -1), 0.3, C_LINE),
            ("TOPPADDING",  (0, 0), (-1, -1), 3),
            ("BOTTOMPADDING", (0, 0), (-1, -1), 3),
            ("LEFTPADDING", (0, 0), (-1, -1), 4),
            ("RIGHTPADDING", (0, 0), (-1, -1), 4),
        ]
        if extra:
            base.extend(extra)
        return TableStyle(base)

    def fmt(v, dec=2):
        try:
            return f"{float(v):,.{dec}f}".replace(",", "X").replace(".", ",").replace("X", ".")
        except Exception:
            return str(v) if v is not None else "—"

    def fmtR(v, dec=2):
        return f"R$ {fmt(v, dec)}" if v is not None else "—"

    story = []

    # ── Cabeçalho ──
    now_str  = _brt_now().strftime("%d/%m/%Y %H:%M")
    scenario_label = snapshot.get("scenario_label", snapshot.get("scenario", "Neutro"))
    story.append(Paragraph("HARBOUR IAT FIF AÇÕES RL — RELATÓRIO DE LIQUIDEZ", st_title))
    story.append(Paragraph(
        f"Emitido em: {now_str}  |  Cenário: {scenario_label}  |  NAV: {fmtR(snapshot.get('nav'))}  |  Histórico CVM: {snapshot.get('n_records_resgate', 0)} registros",
        st_sub
    ))
    story.append(HRFlowable(width="100%", thickness=0.5, color=C_LINE, spaceAfter=6))

    # ── Seção 1: KPIs ──
    story.append(Paragraph("1. KPIs DE LIQUIDEZ", st_sec))
    k = market["kpis"]
    kpi_rows = [
        ["MÉTRICA", "VALOR"],
        ["Valor em carteira (NAV ativos)", fmtR(k["valor_carteira"])],
        ["Volume médio ponderado", fmtR(k["vol_medio_ponderado"]) + " / dia"],
        ["% em alta liquidez (<3d)", f"{fmt(k['pct_alta_liquidez'])}%"],
        ["Prazo médio ponderado p/ zerar", f"{fmt(k['prazo_medio_zerar'], 2)} dias"],
    ]
    kpi_tbl = Table(kpi_rows, colWidths=[9*cm, 6*cm])
    kpi_tbl.setStyle(tbl_style())
    story.append(kpi_tbl)
    story.append(Spacer(1, 6))

    # ── Seção 2: Compliance ──
    story.append(Paragraph("2. COMPLIANCE CVM / ANBIMA", st_sec))
    comp_rows = [["REGRA", "ATUAL", "LIMITE", "STATUS"]]
    comp_extra = []
    for i, r in enumerate(compliance["regras"], start=1):
        unidade = r.get("unidade", "")
        atual = (f"{fmt(r['valor_atual'])}%" if unidade == "%"
                 else f"{fmt(r['valor_atual'], 2)} d" if unidade == "dias"
                 else f"{fmt(r['valor_atual'], 2)}")
        limite = (f"{fmt(r['limite'])}%" if unidade == "%"
                  else f"{fmt(r['limite'], 0)} d" if unidade == "dias"
                  else f"{fmt(r['limite'], 2)}")
        lbl = "OK" if r["status"] == "ok" else "ALERTA" if r["status"] == "alerta" else "VIOLAÇÃO"
        comp_rows.append([r["nome"], atual, limite, lbl])
        comp_extra.append(("TEXTCOLOR", (3, i), (3, i), status_color(r["status"])))
    comp_tbl = Table(comp_rows, colWidths=[14*cm, 4*cm, 4*cm, 3.5*cm])
    comp_tbl.setStyle(tbl_style(extra=comp_extra))
    story.append(comp_tbl)
    story.append(Spacer(1, 6))

    # ── Seção 3: Faixas de liquidez ──
    story.append(Paragraph("3. FAIXAS DE LIQUIDEZ (% DO PL)", st_sec))
    f = market["faixas"]
    faixas_rows = [
        ["FAIXA", "PRAZO", "% DO PL"],
        ["Alta liquidez",      "< 3 dias",   f"{fmt(f.get('alta', 0))}%"],
        ["Média liquidez",     "4-7 dias",   f"{fmt(f.get('media', 0))}%"],
        ["Baixa liquidez",     "8-30 dias",  f"{fmt(f.get('baixa', 0))}%"],
        ["Muito baixa",        "> 30 dias",  f"{fmt(f.get('muito_baixa', 0))}%"],
    ]
    faixas_extra = [
        ("TEXTCOLOR", (0, 1), (0, 1), C_OK),
        ("TEXTCOLOR", (0, 2), (0, 2), colors.HexColor("#3399ff")),
        ("TEXTCOLOR", (0, 3), (0, 3), C_ALERTA),
        ("TEXTCOLOR", (0, 4), (0, 4), C_VIOL),
    ]
    faixas_tbl = Table(faixas_rows, colWidths=[6*cm, 4*cm, 4*cm])
    faixas_tbl.setStyle(tbl_style(extra=faixas_extra))
    story.append(faixas_tbl)
    story.append(PageBreak())

    # ── Seção 4: Tabela Liquidez Ativos (heatmap) ──
    story.append(Paragraph("4. LIQUIDEZ ATIVOS — CUMULATIVO POR BUCKET (% DO PL)", st_sec))
    buckets = snapshot["buckets"]
    header = ["ATIVO", "PRAZO", "VALOR", "%"] + [str(b) for b in buckets]
    ativos_rows = [header]
    ativos_extra = []
    for i, a in enumerate(snapshot["por_ativo"], start=1):
        row = [
            a["ticker"],
            f"{a['prazo_resgate_d']}d",
            fmtR(a["valor_bruto"], 0),
            f"{fmt(a['proporcao_pct'])}%",
        ]
        for v in a["liq_ponderada_por_bucket"]:
            row.append(f"{fmt(v, 2)}%" if v > 0 else "—")
        ativos_rows.append(row)
        # destaque verde nas células > 0
        for j, v in enumerate(a["liq_ponderada_por_bucket"], start=4):
            if v > 0:
                intensity = min(1.0, v / max(0.1, a["proporcao_pct"]))
                alpha = 0.1 + intensity * 0.4
                # rgb interpolation: (0,204,136) com alpha em escala de verde
                green_bg = colors.Color(0, 0.8, 0.53, alpha=alpha)
                ativos_extra.append(("BACKGROUND", (j, i), (j, i), green_bg))

    col_widths = [1.8*cm, 1.0*cm, 2.3*cm, 1.2*cm] + [0.95*cm] * len(buckets)
    ativos_tbl = Table(ativos_rows, colWidths=col_widths, repeatRows=1)
    ativos_tbl.setStyle(tbl_style(extra=ativos_extra))
    story.append(ativos_tbl)
    story.append(Spacer(1, 8))

    # ── Seção 5: Resgate projetado vs liquidez por bucket ──
    story.append(Paragraph("5. RESGATE PROJETADO vs LIQUIDEZ POR BUCKET", st_sec))
    resg_rows = [["BUCKET (DIAS)", "LIQUIDEZ ATIVOS (% PL)", "RESGATE PROJETADO (% PL)", "ÍNDICE LIQUIDEZ"]]
    resg_extra = []
    for i, (b, la, rp, idx) in enumerate(zip(
        buckets, snapshot["liquidez_ativos"], snapshot["resgate_projetado"], snapshot["indice_liquidez"]
    ), start=1):
        idx_str = ">100" if idx > 100 else fmt(idx, 2)
        resg_rows.append([str(b), f"{fmt(la)}%", f"{fmt(rp, 4)}%", idx_str])
        # cor do índice
        if idx >= 1.0:
            resg_extra.append(("TEXTCOLOR", (3, i), (3, i), C_OK))
        elif idx >= 0.7:
            resg_extra.append(("TEXTCOLOR", (3, i), (3, i), C_ALERTA))
        else:
            resg_extra.append(("TEXTCOLOR", (3, i), (3, i), C_VIOL))
    resg_tbl = Table(resg_rows, colWidths=[3.5*cm, 5*cm, 6*cm, 4*cm])
    resg_tbl.setStyle(tbl_style(extra=resg_extra))
    story.append(resg_tbl)

    # ── Rodapé com metodologia ──
    story.append(Spacer(1, 12))
    story.append(Paragraph(
        "Metodologia: dias para liquidar = max(prazo settlement, valor / (volume médio × participação máx 20% do ADV)). "
        "Resgate projetado = percentil das somas rolantes de resg_dia/PL no histórico CVM. "
        "Cenários: Neutro (vol 100%, P50), Stress (vol 50%, P75), Crise (vol 30%, P95).",
        st_foot
    ))

    doc.build(story)
    return buf.getvalue()


@app.route("/api/liquidity/pdf", methods=["GET"])
@require_admin
def api_liquidity_pdf():
    scenario = (request.args.get("scenario") or "neutro").lower()
    if scenario not in LIQUIDEZ_CENARIOS:
        return jsonify({"error": "scenario inválido"}), 400
    try:
        snap   = _build_liquidity_snapshot(scenario=scenario)
        market = _build_liquidity_market()
        compl  = _build_liquidity_compliance()
        pdf_bytes = _generate_liquidity_pdf(snap, market, compl)
    except Exception as e:
        return jsonify({"error": f"Erro ao gerar PDF: {str(e)}"}), 500
    from flask import make_response
    ts = _brt_now().strftime("%Y%m%d_%H%M")
    filename = f"liquidez_{scenario}_{ts}.pdf"
    resp = make_response(pdf_bytes)
    resp.headers["Content-Type"]        = "application/pdf"
    resp.headers["Content-Disposition"] = f'attachment; filename="{filename}"'
    return resp


def _record_liquidity_snapshot(today_str=None):
    """Calcula e grava 1 entrada em liquidity_history.json para a data informada."""
    if today_str is None:
        today_str = _brt_now().strftime("%Y-%m-%d")
    try:
        snap_n = _build_liquidity_snapshot(scenario="neutro")
        snap_s = _build_liquidity_snapshot(scenario="stress")
        snap_c = _build_liquidity_snapshot(scenario="crise")
        market = _build_liquidity_market()
        compl  = _build_liquidity_compliance()

        idx_5d = LIQUIDEZ_BUCKETS.index(5) if 5 in LIQUIDEZ_BUCKETS else 4
        entry = {
            "data":               today_str,
            "nav":                snap_n["nav"],
            "prazo_medio_zerar":  market["kpis"]["prazo_medio_zerar"],
            "pct_alta_liquidez":  market["kpis"]["pct_alta_liquidez"],
            "pct_baixa_liquidez": round(market["faixas"].get("baixa", 0) + market["faixas"].get("muito_baixa", 0), 2),
            "indice_liquidez_5d": {
                "neutro": snap_n["indice_liquidez"][idx_5d],
                "stress": snap_s["indice_liquidez"][idx_5d],
                "crise":  snap_c["indice_liquidez"][idx_5d],
            },
            "compliance_status":  compl["pior_status"],
        }
        hist = load_liquidity_history()
        hist = [h for h in hist if h.get("data") != today_str]  # dedupe
        hist.append(entry)
        hist.sort(key=lambda x: x.get("data") or "")
        save_liquidity_history(hist)
        return entry
    except Exception as e:
        print(f"[record_liquidity_snapshot] erro: {e}")
        return None


if __name__ == "__main__":
    app.run(debug=True, port=5000)
